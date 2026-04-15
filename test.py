"""
test.py — dump GA4 HTML for the Page Performance report and build Slide 5 only.

Usage:
    uv run test.py econet "Mar 1, 2026" "Mar 28, 2026"
"""
from __future__ import annotations

import argparse
import json
import os
import re
import sys
from pathlib import Path

from google import genai
from playwright.sync_api import sync_playwright
from pptx import Presentation

REPO_ROOT = Path(__file__).resolve().parent
if str(REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(REPO_ROOT))

from src.reports.generator import (  # noqa: E402
    _ensure_expected_ga4_property,
    _fill_text_run,
    _goto_ga4_section,
    _launch_persistent_context,
    _replace_image_in_slide,
    _switch_ga4_property_via_search,
)
from src.reports.generator_2026 import (  # noqa: E402
    TEMPLATES_2026,
    _classify_page,
    _open_snapshot_and_set_dates,
    _short_title,
    _write_para_with_highlights,
)
from src.reports.runtime import get_templates_dir, load_runtime_environment  # noqa: E402

PASS = "\033[92m✓\033[0m"
INFO = "\033[94m→\033[0m"
FAIL = "\033[91m✗\033[0m"

TEMPLATES_DIR = get_templates_dir()


def _switch_dimension_to_page_path(page) -> None:
    selectors = [
        "th.cdk-column-ROW_HEADER-unifiedScreenClass-primaryDimensionColumn button[data-guidedhelpid='table-dimension-picker']",
        "button[data-guidedhelpid='table-dimension-picker']",
    ]
    last_error: Exception | None = None

    for selector in selectors:
        try:
            button = page.locator(selector).first
            button.wait_for(state="visible", timeout=8000)
            button.click()
            page.wait_for_timeout(1500)
            break
        except Exception as exc:
            last_error = exc
    else:
        raise RuntimeError("Could not open the GA4 dimension picker") from last_error

    option_queries = [
        lambda: page.get_by_role("menuitem").filter(has_text="Page path").first,
        lambda: page.locator("button.mat-mdc-menu-item").filter(has_text="Page path").first,
        lambda: page.locator("[role='menuitem']").filter(has_text="Page path").first,
        lambda: page.locator("mat-option, .mat-mdc-option").filter(has_text="Page path").first,
        lambda: page.get_by_text("Page path and screen class").first,
    ]
    for query in option_queries:
        try:
            option = query()
            option.wait_for(state="visible", timeout=5000)
            option.click()
            page.wait_for_timeout(3000)
            page.locator("th.cdk-column-__row_index__").first.wait_for(state="visible", timeout=10000)
            return
        except Exception:
            continue

    raise RuntimeError("Could not switch the GA4 table dimension to 'Page path and screen class'")


def _scrape_page_path_rows(page) -> tuple[list[dict], int]:
    pages_data: list[dict] = []
    site_total_views = 0

    body = page.locator("body").inner_text()
    total_match = re.search(r"(?:^|\n)\s*Total\s*\n\s*([\d,]+)\s*\n\s*100% of total", body)
    if total_match:
        try:
            site_total_views = int(total_match.group(1).replace(",", ""))
        except ValueError:
            site_total_views = 0

    for line in body.splitlines():
        cols = [col.strip() for col in line.split("\t") if col.strip()]
        if not cols:
            continue

        if len(cols) < 5 or not cols[0].isdigit() or len(pages_data) >= 10:
            continue

        path_value = cols[1]
        views_match = re.match(r"^([\d,]+)", cols[2])
        active_users_match = re.match(r"^([\d,]+)", cols[3])
        views_pct_match = re.search(r"\(([^)]+%)\)", cols[2])
        active_users_pct_match = re.search(r"\(([^)]+%)\)", cols[3])
        if not views_match or not active_users_match:
            continue

        try:
            pages_data.append({
                "title": path_value,
                "path": path_value,
                "views": int(views_match.group(1).replace(",", "")),
                "views_pct": views_pct_match.group(1) if views_pct_match else "",
                "active_users": int(active_users_match.group(1).replace(",", "")),
                "active_users_pct": active_users_pct_match.group(1) if active_users_pct_match else "",
                "views_per_user": cols[4],
                "avg_engagement_time": cols[5] if len(cols) > 5 else "",
            })
        except ValueError:
            continue

    return pages_data, site_total_views


def _build_slide5_exact(slide, pages_data: list[dict], screenshots: dict) -> None:
    if pages_data:
        for row in pages_data:
            row["_type"] = _classify_page(row["title"])
            raw_label = _short_title(row["title"])
            row["_label"] = "Homepage" if row["_type"] == "home" else raw_label

        top = pages_data[0]
        second = pages_data[1] if len(pages_data) > 1 else None
        compliance_pages = [p for p in pages_data if p["_type"] == "compliance"]
        support_pages = [p for p in pages_data if p["_type"] == "support"]

        heading = (
            f"The {top['_label']}" + (f" and {second['_label']}" if second else "") + " pages drive the most meaningful engagement."
        )

        top_label_display = top["_label"] if top["_label"] != "Homepage" else "Homepage"
        page_suffix = "" if top["_label"] == "Homepage" else " page"
        para1 = (
            f"The {top_label_display}{page_suffix} drives the highest traffic, accounting for "
            f"{top.get('views_pct') or '0%'} of total views ({top['views']:,} views), with "
            f"{top['views_per_user']} views per active user and an average engagement duration of "
            f"{top['avg_engagement_time']}. This substantial volume indicates strong content discoverability and user interest."
        )

        if compliance_pages:
            top_compliance = compliance_pages[:2]
            compliance_bits = [
                f"{p['_label']} ({p.get('views_pct') or '0%'} of views)" for p in top_compliance
            ]
            para2 = (
                f"Compliance pages, including {' and '.join(compliance_bits)}, yield considerable traffic, "
                f"predominantly influenced by campaign landing directives over organic browsing."
            )
        else:
            para2 = (
                f"Secondary pages such as {second['_label'] if second else 'other pages'} contribute meaningful traffic, "
                f"but with more limited evidence of deep exploratory browsing."
            )

        if support_pages:
            support_bits = [
                p["_label"]
                for p in support_pages[:2]
            ]
            para3 = (
                f"Support-focused pages including {' and '.join(support_bits)} attract consistent traffic "
                f"but exhibit lower engagement depth, indicating they are mainly used for quick access to assistance rather than extended browsing."
            )
        else:
            para3 = (
                f"The leading pages drive visibility and entry traffic, while the remaining pages play a more supportive role in the user journey."
            )

        page_names = {p["_label"] for p in pages_data}
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.name == "object 3":
                content_paras = [p for p in shape.text_frame.paragraphs if p.text.strip()]
                if content_paras:
                    _fill_text_run(content_paras[0], heading)
            elif shape.name == "object 7":
                content_paras = [p for p in shape.text_frame.paragraphs if p.text.strip()]
                for i, text in enumerate([para1, para2, para3], start=1):
                    if i < len(content_paras):
                        _write_para_with_highlights(content_paras[i], text, bold_words=page_names)

    if "pages_table" in screenshots:
        picture_names = {
            shape.name
            for shape in slide.shapes
            if getattr(shape, "shape_type", None) == 13
        }
        for candidate in ("Picture 11", "Picture 10", "Picture 13"):
            if candidate in picture_names:
                _replace_image_in_slide(slide, screenshots["pages_table"], shape_name=candidate)
                break


def _fallback_page_label(path_value: str) -> str:
    clean = path_value.strip() or "/"
    if clean == "/":
        return "Homepage"
    clean = clean.split("?", 1)[0].split("#", 1)[0].strip("/")
    if not clean:
        return "Homepage"

    parts = [part for part in clean.split("/") if part]
    last = parts[-1].replace("-", " ").replace("_", " ").strip()
    if not last:
        return "Homepage"
    words = [word for word in last.split() if word]
    return " ".join(word.upper() if word.isupper() else word.capitalize() for word in words[:5])


def _label_page_paths_with_gemini(pages_data: list[dict]) -> None:
    if not pages_data:
        return

    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        for row in pages_data:
            row["title"] = _fallback_page_label(row.get("path", row["title"]))
        return

    client = genai.Client(api_key=api_key)
    path_list = "\n".join(f"- {row.get('path', row['title'])}" for row in pages_data)
    prompt = (
        "You are labeling website page paths for a PowerPoint report.\n"
        "Return strict JSON only: an array of objects with keys path and label.\n"
        "Rules:\n"
        "- label must be 2 to 5 words\n"
        "- use title case\n"
        "- make each label human-friendly\n"
        "- keep important brand/product names when obvious\n"
        "- '/' should become 'Homepage'\n"
        "- do not include commentary or markdown\n\n"
        f"Paths:\n{path_list}"
    )
    try:
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt,
        )
        raw = (response.text or "").strip()
        mappings = json.loads(raw)
        by_path = {
            item["path"]: item["label"].strip()
            for item in mappings
            if isinstance(item, dict) and item.get("path") and item.get("label")
        }
        for row in pages_data:
            path_value = row.get("path", row["title"])
            row["title"] = by_path.get(path_value, _fallback_page_label(path_value))
    except Exception:
        for row in pages_data:
            row["title"] = _fallback_page_label(row.get("path", row["title"]))


def _capture_page_performance(report_name: str, start_date: str, end_date: str) -> tuple[dict[str, Path], list[dict], int]:
    screenshots: dict[str, Path] = {}
    pages_data: list[dict] = []
    site_total_views = 0

    dump_html_path = REPO_ROOT / "dump_page_performance_page_path.html"
    dump_table_html_path = REPO_ROOT / "dump_page_performance_page_path_table.html"
    dump_text_path = REPO_ROOT / "dump_page_performance_page_path.txt"

    with sync_playwright() as playwright:
        context = _launch_persistent_context(playwright, headless=False)
        try:
            page = context.new_page()
            page.bring_to_front()

            print(f"{INFO} Switching GA4 property...")
            page = _switch_ga4_property_via_search(page, report_name)

            print(f"{INFO} Opening snapshot and applying date range...")
            page = _goto_ga4_section(page, report_name, "/home")
            try:
                page.keyboard.press("Escape")
                page.wait_for_timeout(1000)
            except Exception:
                pass
            _open_snapshot_and_set_dates(page, report_name, start_date, end_date)

            print(f"{INFO} Opening Pages and Screens report...")
            page.get_by_role("button", name="View pages and screens", exact=True).click()
            page.wait_for_timeout(4000)
            _ensure_expected_ga4_property(page, report_name)

            row_num_col = page.locator("th.cdk-column-__row_index__").first
            end_col = page.locator("th.cdk-column-DEFAULT-userEngagementDurationPerUser").first
            table = page.locator("table.adv-table").first

            row_num_col.wait_for(state="visible", timeout=10000)
            page.keyboard.press("End")
            page.wait_for_timeout(1000)
            page.mouse.wheel(0, 3000)
            page.wait_for_timeout(1000)

            print(f"{INFO} Switching table dimension to page path...")
            _switch_dimension_to_page_path(page)

            print(f"{INFO} Dumping HTML/text debug files...")
            dump_html_path.write_text(page.content(), encoding="utf-8")
            dump_table_html_path.write_text(table.inner_html(), encoding="utf-8")
            dump_text_path.write_text(page.locator("body").inner_text(), encoding="utf-8")

            start_box = row_num_col.bounding_box()
            end_box = end_col.bounding_box()
            table_box = table.bounding_box()
            if start_box and end_box and table_box:
                clip = {
                    "x": start_box["x"],
                    "y": table_box["y"],
                    "width": (end_box["x"] + end_box["width"]) - start_box["x"],
                    "height": table_box["height"],
                }
                screenshot_path = REPO_ROOT / "dump_page_performance_page_path_table.png"
                page.screenshot(path=str(screenshot_path), clip=clip, full_page=True)
                screenshots["pages_table"] = screenshot_path

            pages_data, site_total_views = _scrape_page_path_rows(page)
            _label_page_paths_with_gemini(pages_data)
        finally:
            try:
                context.close()
            except Exception:
                pass

    print(f"{PASS} Wrote page-path HTML dump → {dump_html_path.resolve()}")
    print(f"{PASS} Wrote page-path table HTML dump → {dump_table_html_path.resolve()}")
    print(f"{PASS} Wrote page-path text dump → {dump_text_path.resolve()}")
    if "pages_table" in screenshots:
        print(f"{PASS} Wrote page-path table screenshot → {screenshots['pages_table'].resolve()}")

    return screenshots, pages_data, site_total_views


def run(report_name: str, start_date: str, end_date: str) -> None:
    load_runtime_environment()
    print(f"\n{INFO} Report: {report_name}")
    print(f"{INFO} Date range: {start_date} → {end_date}\n")

    screenshots, pages_data, site_total_views = _capture_page_performance(
        report_name, start_date, end_date
    )

    print(f"{PASS} Scraped page rows: {len(pages_data)}")
    print(f"{PASS} Site total views: {site_total_views}")
    for page_row in pages_data[:5]:
        print(
            f"  {PASS} {page_row['title']} [{page_row.get('path', '')}] | views={page_row['views']} | "
            f"active_users={page_row['active_users']} | "
            f"views_per_user={page_row['views_per_user']} | "
            f"avg_engagement_time={page_row['avg_engagement_time']}"
        )

    template_path = TEMPLATES_DIR / TEMPLATES_2026[report_name]
    prs = Presentation(str(template_path))

    print(f"\n{INFO} Building slide 5 (Page Performance)...")
    slide_idx = 4
    _build_slide5_exact(prs.slides[slide_idx], pages_data, screenshots)
    print(f"{PASS} Slide 5 done.")

    slide_list = prs.slides._sldIdLst
    ids = list(slide_list)
    for i, el in enumerate(ids):
        if i != slide_idx:
            slide_list.remove(el)

    out_path = REPO_ROOT / "test_slide5_page_performance.pptx"
    prs.save(str(out_path))
    print(f"\n{PASS} Saved → {out_path.resolve()}")
    print("Open test_slide5_page_performance.pptx to inspect.")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("report_name", choices=sorted(TEMPLATES_2026.keys()))
    parser.add_argument("start_date", help='e.g. "Mar 1, 2026"')
    parser.add_argument("end_date", help='e.g. "Mar 28, 2026"')
    args = parser.parse_args()
    run(args.report_name, args.start_date, args.end_date)


if __name__ == "__main__":
    main()
