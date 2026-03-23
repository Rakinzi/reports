"""
Extract all text from each slide of a PPTX file, grouped by slide.

Usage:
    uv run python scripts/extract_pptx_text.py <path-to-pptx>

Example:
    uv run python scripts/extract_pptx_text.py "src/reports/report-templates/Econet-February 2026 Website Report.pptx"
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt


def extract_text(pptx_path: str) -> None:
    path = Path(pptx_path)
    if not path.exists():
        print(f"File not found: {pptx_path}")
        sys.exit(1)

    prs = Presentation(pptx_path)

    print(f"\n{'='*60}")
    print(f"  {path.name}")
    print(f"  {len(prs.slides)} slides")
    print(f"{'='*60}\n")

    for slide_num, slide in enumerate(prs.slides, start=1):
        print(f"--- Slide {slide_num} ---")

        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if line:
                    texts.append(line)

        if texts:
            for t in texts:
                print(f"  {t}")
        else:
            print("  (no text)")

        print()


if __name__ == "__main__":
    if len(sys.argv) < 2:
        # Default to listing all templates if no argument given
        templates_dir = Path("src/reports/report-templates")
        pptx_files = list(templates_dir.glob("*.pptx"))
        if not pptx_files:
            print("No PPTX files found. Pass a path as argument.")
            sys.exit(1)
        for f in pptx_files:
            extract_text(str(f))
    else:
        extract_text(sys.argv[1])
