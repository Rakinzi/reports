"""
Generates a report from a user-uploaded PPTX template.

Each template has a field_map (list of ShapeMapping dicts) that tells us which
shape on which slide holds which piece of data. We:
  1. Open GA4 with the saved Chrome session
  2. For each property section (or the single default property): switch GA4 property,
     set the date range, scrape home + snapshot metrics, capture screenshots/charts
  3. Build text_values dict per section using Gemini (batched)
  4. Copy the PPTX, walk shape mappings, and fill each shape with its section's data
  5. Return the output path

Multi-property (combined) reports: when property_sections are configured, each
section defines a slide range and a GA4 property ID. Shapes on those slides get
filled with that section's metrics. Backwards-compat: if no sections, the single
ga4_property_id on the template row is used as before.
"""

import json
import os
import re
import shutil
from pathlib import Path

from .db import get_template_by_slug, list_template_sections
from .logging_utils import configure_logging
from .runtime import get_output_dir, load_runtime_environment

logger = configure_logging()

# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def _section_for_slide(sections: list[dict], slide_index: int) -> dict | None:
    """Return the first section whose slide range contains slide_index, or None."""
    for sec in sections:
        if sec["start_slide"] <= slide_index <= sec["end_slide"]:
            return sec
    return None


def generate_user_template(
    report_name: str,
    date_range: str,
    report_date: str,
    start_date: str,
    end_date: str,
    _stage_callback=None,
) -> Path:
    def stage(msg: str) -> None:
        logger.info("[template_runner] %s", msg)
        if _stage_callback:
            _stage_callback(msg)

    template = get_template_by_slug(report_name)
    if template is None:
        raise ValueError(f"No template found for slug '{report_name}'")

    field_map: list[dict] = json.loads(template["field_map"] or "[]")
    pptx_path = Path(template["pptx_path"])
    if not pptx_path.exists():
        raise FileNotFoundError(f"Template PPTX not found: {pptx_path}")

    sections = list_template_sections(template["id"])

    # Build the list of properties to scrape.
    # Each entry is a dict with keys: ga4_property_id, gsc_url, section_key
    # section_key is None for the single-property fallback, or the section id.
    if sections:
        props_to_scrape = [
            {
                "ga4_property_id": s["ga4_property_id"],
                "gsc_url": s.get("gsc_url", ""),
                "section_key": s["id"],
                "start_slide": s["start_slide"],
                "end_slide": s["end_slide"],
            }
            for s in sections
            if s.get("ga4_property_id")
        ]
        if not props_to_scrape:
            raise ValueError(
                f"Template '{report_name}' has property sections configured but none have a GA4 property ID"
            )
    else:
        ga4_property_id: str = template.get("ga4_property_id", "")
        if not ga4_property_id:
            raise ValueError(f"Template '{report_name}' has no GA4 property ID configured")
        props_to_scrape = [
            {
                "ga4_property_id": ga4_property_id,
                "gsc_url": template.get("gsc_url", ""),
                "section_key": None,
                "start_slide": 0,
                "end_slide": 99999,
            }
        ]

    load_runtime_environment()

    from playwright.sync_api import sync_playwright
    from .generator import (
        _launch_persistent_context,
        _switch_ga4_property_by_id,
        _set_date_range,
        _scrape_home_metrics,
        _scrape_snapshot_metrics,
    )

    # metrics_by_key: section_key → metrics_context dict
    metrics_by_key: dict = {}
    # image_paths_by_key: section_key → {field_type: Path}
    image_paths_by_key: dict = {}

    n = len(props_to_scrape)
    with sync_playwright() as pw:
        ctx = _launch_persistent_context(pw)
        try:
            page = ctx.new_page()
            for idx, prop in enumerate(props_to_scrape):
                pid = prop["ga4_property_id"]
                key = prop["section_key"]
                label = f"section {idx + 1}/{n} ({pid})" if sections else pid

                stage(f"Switching GA4 to property {label}...")
                page = _switch_ga4_property_by_id(page, pid)

                stage(f"Scraping home metrics for {label}...")
                _set_date_range(page, start_date, end_date)
                home_metrics = _scrape_home_metrics(page)

                stage(f"Scraping snapshot metrics for {label}...")
                _set_date_range(page, start_date, end_date)
                snapshot_metrics = _scrape_snapshot_metrics(page)

                stage(f"Collecting live metrics for {label}...")
                metrics_context = _collect_live_metrics(
                    page=page,
                    ga4_property_id=pid,
                    start_date=start_date,
                    end_date=end_date,
                    home_metrics=home_metrics,
                    snapshot_metrics=snapshot_metrics,
                )
                metrics_by_key[key] = metrics_context

                # Only capture image fields whose shapes belong to this section's slide range
                section_field_map = [
                    m for m in field_map
                    if prop["start_slide"] <= m.get("slide_index", 0) <= prop["end_slide"]
                ]
                stage(f"Capturing screenshots/charts for {label}...")
                image_paths_by_key[key] = _capture_image_fields(
                    page=page,
                    field_map=section_field_map,
                    report_name=f"{report_name}_sec{key}" if key is not None else report_name,
                    start_date=start_date,
                    end_date=end_date,
                    ga4_property_id=pid,
                    metrics_context=metrics_context,
                )
        finally:
            ctx.close()

    stage("Generating text content with Gemini...")
    # Build text_values per section key
    text_values_by_key: dict = {}
    for prop in props_to_scrape:
        key = prop["section_key"]
        ctx = metrics_by_key[key]
        section_field_map = [
            m for m in field_map
            if prop["start_slide"] <= m.get("slide_index", 0) <= prop["end_slide"]
        ]
        text_values_by_key[key] = _build_text_values(
            section_field_map,
            ctx["home_metrics"],
            ctx["snapshot_metrics"],
            date_range,
            report_date,
        )

    stage("Filling template shapes...")
    output_path = get_output_dir() / f"{report_name}-{report_date.replace(' ', '_')}.pptx"
    shutil.copy(str(pptx_path), str(output_path))
    _fill_template_sections(
        output_path, field_map, props_to_scrape, sections,
        text_values_by_key, image_paths_by_key,
    )

    logger.info("User template report generated: %s", output_path)
    return output_path


# ---------------------------------------------------------------------------
# Live GA4 context
# ---------------------------------------------------------------------------

def _parse_metric_number(value: str | int | float | None) -> int:
    if value is None:
        return 0
    if isinstance(value, (int, float)):
        return int(value)

    normalized = str(value).strip().replace(",", "")
    if not normalized:
        return 0

    multiplier = 1
    suffix = normalized[-1].upper()
    if suffix == "K":
        multiplier = 1_000
        normalized = normalized[:-1]
    elif suffix == "M":
        multiplier = 1_000_000
        normalized = normalized[:-1]
    elif suffix == "B":
        multiplier = 1_000_000_000
        normalized = normalized[:-1]

    try:
        return int(float(normalized) * multiplier)
    except ValueError:
        return 0


def _restore_property_home(page, ga4_property_id: str):
    from .generator import _switch_ga4_property_by_id

    page = _switch_ga4_property_by_id(page, ga4_property_id)
    page.wait_for_timeout(1500)
    return page


def _collect_weekly_active_users(page, start_date: str, end_date: str) -> dict[str, int]:
    from .generator import _weekly_ranges, _set_date_range, _scrape_home_metrics

    weekly_active_users: dict[str, int] = {}
    for label, week_start, week_end in _weekly_ranges(start_date, end_date):
        try:
            _set_date_range(page, week_start, week_end)
            metrics = _scrape_home_metrics(page)
            weekly_active_users[label] = _parse_metric_number(metrics.get("Active users"))
        except Exception as exc:
            logger.warning("Could not collect weekly active users for %s: %s", label, exc)
            weekly_active_users[label] = 0

    _set_date_range(page, start_date, end_date)
    return weekly_active_users


def _collect_acquisition_metrics(page) -> dict[str, int]:
    acquisition: dict[str, int] = {}
    try:
        page.get_by_role("button", name="View user acquisition", exact=True).click()
        page.wait_for_timeout(4000)
        body = page.locator("body").inner_text()
        for line in [text.strip() for text in body.splitlines() if text.strip()]:
            match = re.match(r"^\d+\t(.+?)\t(\d+)\s*\(", line)
            if match:
                acquisition[match.group(1).strip()] = int(match.group(2))
    except Exception as exc:
        logger.warning("Could not collect acquisition metrics: %s", exc)
    return acquisition


def _collect_page_views(page) -> dict[str, int]:
    page_views: dict[str, int] = {}
    try:
        page.get_by_role("button", name="View pages and screens", exact=True).click()
        page.wait_for_timeout(4000)
        body = page.locator("body").inner_text()
        for line in [text.strip() for text in body.splitlines() if text.strip()]:
            match = re.match(r"^\d+\t(.+?)\t(\d+)\s*", line)
            if match and len(page_views) < 8:
                page_views[match.group(1).strip()] = int(match.group(2).strip())
    except Exception as exc:
        logger.warning("Could not collect page views: %s", exc)
    return page_views


def _collect_live_metrics(
    page,
    ga4_property_id: str,
    start_date: str,
    end_date: str,
    home_metrics: dict,
    snapshot_metrics: dict,
) -> dict[str, dict]:
    weekly_active_users: dict[str, int] = {}
    acquisition: dict[str, int] = {}
    page_views: dict[str, int] = {}

    try:
        page = _restore_property_home(page, ga4_property_id)
        weekly_active_users = _collect_weekly_active_users(page, start_date, end_date)
    except Exception as exc:
        logger.warning("Could not collect weekly active users: %s", exc)

    try:
        page = _restore_property_home(page, ga4_property_id)
        acquisition = _collect_acquisition_metrics(page)
    except Exception as exc:
        logger.warning("Could not collect acquisition metrics: %s", exc)

    try:
        page = _restore_property_home(page, ga4_property_id)
        page_views = _collect_page_views(page)
    except Exception as exc:
        logger.warning("Could not collect page views: %s", exc)

    return {
        "home_metrics": home_metrics,
        "snapshot_metrics": snapshot_metrics,
        "acquisition": acquisition,
        "page_views": page_views,
        "weekly_active_users": weekly_active_users,
    }


# ---------------------------------------------------------------------------
# Metric / text value building
# ---------------------------------------------------------------------------

def _parse_perf_month(report_date: str) -> str:
    """Convert '03 March 2026' → 'March 2026'."""
    parts = report_date.strip().split()
    if len(parts) >= 3:
        return f"{parts[1]} {parts[2]}"
    return report_date


def _compute_new_users_pct(home_metrics: dict) -> str:
    """Try to parse new-users percentage from metric labels like 'New users 78'."""
    import re
    for key, val in home_metrics.items():
        if "new user" in key.lower():
            m = re.search(r"(\d[\d,]*)", str(val))
            if m:
                return m.group(1) + "%"
    return "N/A"


def _build_text_values(
    field_map: list[dict],
    home_metrics: dict,
    snapshot_metrics: dict,
    date_range: str,
    report_date: str,
) -> dict[str, str]:
    """Return a dict of field_type → text value for all text-type mappings."""
    values: dict[str, str] = {}
    gemini_raws: list[tuple[str, str]] = []  # (field_type, raw_text)

    def hm(key: str) -> str:
        return home_metrics.get(key, "N/A")

    for mapping in field_map:
        if mapping.get("shape_type") != "text":
            continue
        ft = mapping["field_type"]
        if not ft:
            continue

        if ft == "perf_month":
            values[ft] = _parse_perf_month(report_date)
        elif ft == "date_range":
            values[ft] = date_range
        elif ft == "report_date":
            values[ft] = report_date
        elif ft == "active_users":
            values[ft] = hm("Active users")
        elif ft == "new_users":
            values[ft] = hm("New users")
        elif ft == "engagement_time":
            values[ft] = hm("Average engagement time per active user")
        elif ft == "new_users_pct":
            values[ft] = _compute_new_users_pct(home_metrics)
        elif ft == "ctr":
            values[ft] = snapshot_metrics.get("ctr", "N/A")
        elif ft.startswith("narrative_") or ft.startswith("subtitle_") or ft == "recommendations":
            # Collect for batch Gemini call
            raw = _build_gemini_raw(ft, home_metrics, snapshot_metrics)
            gemini_raws.append((ft, raw))
        else:
            values[ft] = "N/A"

    # Batch all Gemini fields in a single call
    if gemini_raws:
        from .generator_2026 import _gemini_paras_batch
        texts = [raw for _, raw in gemini_raws]
        paraphrased = _gemini_paras_batch(texts)
        for (ft, _), para in zip(gemini_raws, paraphrased):
            values[ft] = para

    return values


def _build_gemini_raw(field_type: str, home_metrics: dict, snapshot_metrics: dict) -> str:
    """Build the raw stats string that Gemini will paraphrase for a given Gemini field type."""
    hm = home_metrics
    sm = snapshot_metrics

    if field_type == "recommendations":
        return (
            f"Based on the following GA4 data, write 3 numbered recommendations:\n"
            f"Active users: {hm.get('Active users', 'N/A')}\n"
            f"New users: {hm.get('New users', 'N/A')}\n"
            f"Engagement time: {hm.get('Average engagement time per active user', 'N/A')}"
        )

    # Generic: provide all available metrics as context
    metrics_str = "\n".join(f"{k}: {v}" for k, v in {**hm, **sm}.items())
    section = field_type.replace("narrative_", "").replace("subtitle_", "").replace("_", " ")
    return f"Write a professional insight about {section} performance:\n{metrics_str}"


# ---------------------------------------------------------------------------
# Screenshot / chart image capturing
# ---------------------------------------------------------------------------

def _capture_image_fields(
    page,
    field_map: list[dict],
    report_name: str,
    start_date: str,
    end_date: str,
    ga4_property_id: str,
    metrics_context: dict[str, dict],
) -> dict[str, Path]:
    """Return a dict of field_type → image Path for all image-type mappings."""
    image_paths: dict[str, Path] = {}

    for mapping in field_map:
        if mapping.get("shape_type") != "image":
            continue
        ft = mapping.get("field_type", "")
        if not ft or ft in image_paths:
            continue

        try:
            path = _capture_single_image_field(
                page=page,
                field_type=ft,
                report_name=report_name,
                start_date=start_date,
                end_date=end_date,
                ga4_property_id=ga4_property_id,
                metrics_context=metrics_context,
            )
            if path:
                image_paths[ft] = path
        except Exception as exc:
            logger.warning("Could not capture image field '%s': %s", ft, exc)

    return image_paths


def _capture_single_image_field(
    page,
    field_type: str,
    report_name: str,
    start_date: str,
    end_date: str,
    ga4_property_id: str,
    metrics_context: dict[str, dict],
) -> Path | None:
    """Capture a single screenshot or generate a chart for the given field_type."""
    from .runtime import get_screenshots_dir
    from .generator import _set_date_range

    screenshots_dir = get_screenshots_dir() / f"template_{report_name}"
    screenshots_dir.mkdir(parents=True, exist_ok=True)
    out_path = screenshots_dir / f"{field_type}.png"

    if field_type == "screenshot_snapshot_card":
        page = _restore_property_home(page, ga4_property_id)
        _set_date_range(page, start_date, end_date)
        page.screenshot(path=str(out_path), clip={"x": 0, "y": 0, "width": 1920, "height": 600})
    elif field_type == "screenshot_countries_table":
        page = _restore_property_home(page, ga4_property_id)
        page.get_by_text("View reports snapshot").click()
        page.wait_for_timeout(3000)
        page.get_by_text("View countries").click()
        page.wait_for_timeout(4000)
        page.screenshot(path=str(out_path), full_page=True)
    elif field_type == "screenshot_pages_table":
        page = _restore_property_home(page, ga4_property_id)
        page.get_by_role("button", name="View pages and screens", exact=True).click()
        page.wait_for_timeout(4000)
        page.screenshot(path=str(out_path), full_page=True)
    elif field_type == "screenshot_search_console":
        page.screenshot(path=str(out_path))
    elif field_type.startswith("chart_"):
        out_path = _generate_chart(field_type, out_path, metrics_context)
    else:
        return None

    return out_path if out_path and out_path.exists() else None


def _generate_chart(field_type: str, out_path: Path, metrics_context: dict[str, dict]) -> Path | None:
    """Generate a chart image for a chart_ field type."""
    try:
        from .charts import (
            generate_country_bar_chart,
            generate_traffic_source_pie_chart,
            generate_line_chart,
            generate_page_views_bar_chart,
        )
    except ImportError:
        logger.warning("charts module not available — skipping chart generation")
        return None

    snapshot_metrics = metrics_context.get("snapshot_metrics", {})
    acquisition = metrics_context.get("acquisition", {})
    weekly_active_users = metrics_context.get("weekly_active_users", {})
    page_views = metrics_context.get("page_views", {})

    if field_type == "chart_country_bar":
        countries = snapshot_metrics.get("countries", {})
        data = dict(sorted(countries.items(), key=lambda item: item[1], reverse=True)[:5])
        if not data:
            return None
        generate_country_bar_chart(data, out_path)
    elif field_type == "chart_traffic_pie":
        preferred_channels = [
            "Organic Search",
            "Direct",
            "Referral",
            "Organic Social",
            "Unassigned",
            "Paid Search",
            "Email",
        ]
        data = {channel: acquisition[channel] for channel in preferred_channels if acquisition.get(channel)}
        if not data:
            return None
        generate_traffic_source_pie_chart(data, out_path)
    elif field_type == "chart_line":
        data = {label: value for label, value in weekly_active_users.items() if value > 0}
        if not data:
            return None
        generate_line_chart(data, out_path)
    elif field_type == "chart_page_views_bar":
        if not page_views:
            return None
        shortened_page_views = {}
        for name, views in page_views.items():
            short = name.split("-")[0].strip() if "-" in name else name
            shortened_page_views[short or name] = views
        data = dict(sorted(shortened_page_views.items(), key=lambda item: item[1], reverse=True)[:5])
        generate_page_views_bar_chart(data, out_path)
    else:
        return None

    return out_path


# ---------------------------------------------------------------------------
# PPTX shape filling
# ---------------------------------------------------------------------------

def _find_shape_by_name(slide, name: str):
    return next((s for s in slide.shapes if s.name == name), None)


def _fill_template_sections(
    output_path: Path,
    field_map: list[dict],
    props_to_scrape: list[dict],
    sections: list[dict],
    text_values_by_key: dict,
    image_paths_by_key: dict,
) -> None:
    """Fill all mapped shapes using per-section metrics."""
    from pptx import Presentation
    from .generator import _fill_text_run, _replace_image_in_slide

    prs = Presentation(str(output_path))

    for mapping in field_map:
        slide_index = mapping.get("slide_index", 0)
        shape_name = mapping.get("shape_name", "")
        field_type = mapping.get("field_type", "")
        shape_type = mapping.get("shape_type", "text")

        if not field_type:
            continue
        if slide_index >= len(prs.slides):
            logger.warning("Slide index %d out of range — skipping shape '%s'", slide_index, shape_name)
            continue

        # Resolve which section key owns this slide
        section_key = None
        for prop in props_to_scrape:
            if prop["start_slide"] <= slide_index <= prop["end_slide"]:
                section_key = prop["section_key"]
                break

        text_values = text_values_by_key.get(section_key, {})
        image_paths = image_paths_by_key.get(section_key, {})

        slide = prs.slides[slide_index]
        shape = _find_shape_by_name(slide, shape_name)
        if shape is None:
            logger.warning("Shape '%s' not found in slide %d — skipping", shape_name, slide_index)
            continue

        if shape_type == "text":
            value = text_values.get(field_type, "")
            if not value:
                continue
            if shape.has_text_frame and shape.text_frame.paragraphs:
                _fill_text_run(shape.text_frame.paragraphs[0], value)
            else:
                logger.warning("Shape '%s' has no text frame — skipping", shape_name)

        elif shape_type == "image":
            img_path = image_paths.get(field_type)
            if img_path and img_path.exists():
                try:
                    _replace_image_in_slide(slide, img_path, shape_name=shape_name)
                except Exception as exc:
                    logger.warning("Could not replace image for shape '%s': %s", shape_name, exc)
            else:
                logger.warning("No image available for field_type '%s' — skipping", field_type)

    prs.save(str(output_path))
    logger.info("Saved filled template to %s", output_path)
