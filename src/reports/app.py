import json
import os
import re
import shutil
import threading
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path

from fastapi import FastAPI, File, Form, HTTPException, Request, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, FileResponse, StreamingResponse

from .auth_session import auth_session_status, open_google_sign_in
from .chrome_profiles import list_profiles
from .db import (
    init_db, create_report, list_reports, get_report,
    update_report_completed, update_report_failed,
    update_report_slides_dir, update_report_edits, update_report_stage,
    delete_report,
    create_template, get_template, get_template_by_slug, list_templates,
    update_template_config, update_template_preview_dir, delete_template,
    upsert_template_shapes, list_template_shapes,
    upsert_template_sections, list_template_sections,
)
from .logging_utils import configure_logging, get_log_path, read_recent_logs, stream_logs
from .runtime import get_app_data_dir, get_runtime_status, load_runtime_environment, save_settings, get_user_templates_dir
from .schemas import AppSettingsUpdate, GenerateReportRequest, HARDCODED_REPORT_NAMES

_executor = ThreadPoolExecutor(max_workers=1)
logger = configure_logging()

# Maps report_id -> Event; set the event to request cancellation.
_cancel_flags: dict[int, threading.Event] = {}


def _run_generate(report_id: int, report_name: str, date_range: str, report_date: str, start_date: str, end_date: str):
    cancel_flag = _cancel_flags.get(report_id)

    def stage_callback(stage: str) -> None:
        if cancel_flag and cancel_flag.is_set():
            raise InterruptedError("Report generation cancelled by user.")
        update_report_stage(report_id, stage)

    try:
        from .generator_2026 import TEMPLATES_2026, generate_report_2026
        from .generator import generate_report

        if report_name in TEMPLATES_2026:
            generate_fn = generate_report_2026
        elif get_template_by_slug(report_name) is not None:
            from .template_runner import generate_user_template
            generate_fn = generate_user_template
        else:
            generate_fn = generate_report

        logger.info("Starting report generation for report_id=%s report_name=%s", report_id, report_name)
        update_report_stage(report_id, "Capturing GA4 data...")
        output_path = generate_fn(
            report_name=report_name,
            date_range=date_range,
            report_date=report_date,
            start_date=start_date,
            end_date=end_date,
            _stage_callback=stage_callback,
        )
        update_report_stage(report_id, "Finalising report...")
        update_report_completed(report_id, str(output_path))
        logger.info("Completed report generation for report_id=%s output_path=%s", report_id, output_path)

        # Render PDF preview
        try:
            from .slides import render_pdf
            from pathlib import Path as _Path
            update_report_stage(report_id, "Rendering slide previews...")
            pdf_path = render_pdf(report_id, _Path(output_path))
            update_report_slides_dir(report_id, str(pdf_path.parent))
            logger.info("Rendered PDF preview for report_id=%s pdf=%s", report_id, pdf_path)
        except Exception as render_err:
            logger.warning("PDF preview rendering failed for report_id=%s: %s", report_id, render_err)
    except InterruptedError as e:
        update_report_failed(report_id, str(e))
        logger.info("Report generation cancelled for report_id=%s", report_id)
    except Exception as e:
        update_report_failed(report_id, str(e))
        logger.exception("Report generation failed for report_id=%s", report_id)
    finally:
        _cancel_flags.pop(report_id, None)

app = FastAPI(title="Reports API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.on_event("startup")
def on_startup():
    load_runtime_environment()
    init_db()
    from .db import fail_orphaned_reports
    fail_orphaned_reports()
    logger.info("Reports API started. log_path=%s", get_log_path())


@app.get("/health")
def health():
    return JSONResponse({"status": "ok", **get_runtime_status()})


@app.get("/settings")
def get_settings():
    return JSONResponse(get_runtime_status())


@app.put("/settings")
def put_settings(body: AppSettingsUpdate):
    save_settings(body.model_dump())
    return JSONResponse(get_runtime_status())


@app.get("/settings/chrome-profiles")
def get_chrome_profiles(user_data_dir: str | None = None):
    try:
        return JSONResponse(list_profiles(user_data_dir))
    except FileNotFoundError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc


@app.get("/settings/test-gemini")
def test_gemini_key():
    load_runtime_environment()
    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        raise HTTPException(status_code=400, detail="Gemini API key is not configured.")
    try:
        from google import genai as _genai
        client = _genai.Client(api_key=api_key)
        client.models.generate_content(model="gemini-2.5-flash", contents="Reply with OK")
        return JSONResponse({"ok": True})
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e)) from e


@app.get("/settings/google-sign-in")
def get_google_sign_in_status():
    return JSONResponse(auth_session_status())


@app.post("/settings/google-sign-in")
def post_google_sign_in():
    status = get_runtime_status()
    if not status["browser_available"]:
        raise HTTPException(
            status_code=400,
            detail="No compatible browser installation was found. Install Google Chrome, Microsoft Edge, or Chromium.",
        )
    return JSONResponse(open_google_sign_in())


@app.get("/logs")
def get_logs(limit: int = 200):
    return JSONResponse({"path": str(get_log_path()), "lines": read_recent_logs(limit)})


@app.get("/logs/stream")
def get_logs_stream():
    """SSE endpoint — streams log lines in real time to any connected client."""
    return StreamingResponse(
        stream_logs(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )


@app.post("/reports/generate", status_code=202)
def post_generate_report(body: GenerateReportRequest):
    status = get_runtime_status()
    if not status["gemini_api_key_set"]:
        raise HTTPException(status_code=400, detail="Gemini API key is not configured")
    if not status["browser_available"]:
        raise HTTPException(
            status_code=400,
            detail="No compatible browser installation was found. Install Google Chrome, Microsoft Edge, or Chromium.",
        )

    report_id = create_report(body.report_name, body.date_range, body.report_date)
    _cancel_flags[report_id] = threading.Event()
    _executor.submit(
        _run_generate,
        report_id,
        body.report_name,
        body.date_range,
        body.report_date,
        body.start_date,
        body.end_date,
    )
    return JSONResponse({"id": report_id, "status": "pending"}, status_code=202)


@app.get("/reports")
def get_reports():
    return JSONResponse(list_reports())


@app.get("/reports/{report_id}")
def get_report_by_id(report_id: int):
    report = get_report(report_id)
    if not report:
        raise HTTPException(status_code=404, detail="Report not found")
    return JSONResponse(report)


@app.post("/reports/{report_id}/cancel")
def cancel_report(report_id: int):
    report = get_report(report_id)
    if not report:
        raise HTTPException(status_code=404, detail="Report not found")
    if report["status"] != "pending":
        raise HTTPException(status_code=400, detail="Report is not running")
    flag = _cancel_flags.get(report_id)
    if flag:
        flag.set()
    return JSONResponse({"id": report_id, "cancelled": True})


@app.delete("/reports/{report_id}")
def delete_report_by_id(report_id: int):
    report = get_report(report_id)
    if not report:
        raise HTTPException(status_code=404, detail="Report not found")
    if report["status"] == "pending":
        flag = _cancel_flags.get(report_id)
        if flag:
            flag.set()
    deleted = delete_report(report_id)
    if not deleted:
        raise HTTPException(status_code=404, detail="Report not found")
    return JSONResponse({"id": report_id, "deleted": True})


def _resolve_report_path(stored: str) -> Path:
    """Resolve a stored output_path to an absolute path.

    New records are always absolute. Old records may have been written with a relative
    path (e.g. 'artifacts/output/foo.pptx'). Try resolving those against the project
    root (parent of app_data_dir) so dev-mode old records still work.
    """
    p = Path(stored)
    if p.is_absolute():
        return p
    candidate = get_app_data_dir().parent / p
    if candidate.exists():
        return candidate
    return Path.cwd() / p


@app.get("/reports/{report_id}/download")
def download_report(report_id: int):
    report = get_report(report_id)
    if not report:
        raise HTTPException(status_code=404, detail="Report not found")
    if report["status"] != "completed" or not report["output_path"]:
        raise HTTPException(status_code=400, detail="Report is not ready for download")
    output_path = _resolve_report_path(report["output_path"])
    if not output_path.exists():
        raise HTTPException(status_code=404, detail="Report file not found on disk")
    return FileResponse(
        path=str(output_path),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=output_path.name,
    )


@app.get("/reports/{report_id}/preview.pdf")
def get_report_preview_pdf(report_id: int):
    from pathlib import Path as _Path
    report = get_report(report_id)
    if not report:
        raise HTTPException(status_code=404, detail="Report not found")
    if report["status"] != "completed" or not report["output_path"]:
        raise HTTPException(status_code=400, detail="Report is not ready")

    output_path = _resolve_report_path(report["output_path"])
    if not output_path.exists():
        raise HTTPException(status_code=404, detail="Report file not found on disk")

    # Check for a cached preview.pdf
    slides_dir = report.get("slides_dir")
    if slides_dir:
        cached = _Path(slides_dir) / "preview.pdf"
        if cached.exists():
            return FileResponse(str(cached), media_type="application/pdf")

    # Generate on demand
    try:
        from .slides import render_pdf
        pdf_path = render_pdf(report_id, output_path)
        return FileResponse(str(pdf_path), media_type="application/pdf")
    except Exception as e:
        raise HTTPException(status_code=503, detail=str(e)) from e


@app.get("/reports/{report_id}/slides")
def get_report_slides(report_id: int):
    from pathlib import Path as _Path
    report = get_report(report_id)
    if not report:
        raise HTTPException(status_code=404, detail="Report not found")
    if report["status"] != "completed" or not report["output_path"]:
        raise HTTPException(status_code=400, detail="Report is not ready")

    output_path = _resolve_report_path(report["output_path"])
    if not output_path.exists():
        raise HTTPException(status_code=404, detail="Report file not found on disk")

    from .slides import extract_slide_fields
    fields_by_slide = extract_slide_fields(output_path)

    slides_dir = report.get("slides_dir")
    result = []
    for slide_data in fields_by_slide:
        idx = slide_data["slide_index"]
        has_image = False
        if slides_dir:
            has_image = (_Path(slides_dir) / f"slide_{idx}.png").exists()
        result.append({
            "slide_index": idx,
            "image_url": f"/reports/{report_id}/slides/{idx}/image" if has_image else None,
            "fields": slide_data["fields"],
        })
    return JSONResponse(result)


@app.get("/reports/{report_id}/slides/{slide_index}/image")
def get_slide_image(report_id: int, slide_index: int):
    from pathlib import Path as _Path
    report = get_report(report_id)
    if not report or not report.get("slides_dir"):
        raise HTTPException(status_code=404, detail="Slide images not available")
    image_path = _Path(report["slides_dir"]) / f"slide_{slide_index}.png"
    if not image_path.exists():
        raise HTTPException(status_code=404, detail=f"Slide {slide_index} image not found")
    return FileResponse(str(image_path), media_type="image/png")


@app.post("/reports/{report_id}/slides/{slide_index}/rewrite")
def rewrite_slide_field(report_id: int, slide_index: int, body: dict):
    from .runtime import load_runtime_environment
    import os
    load_runtime_environment()

    current_text = body.get("current_text", "")
    instruction = body.get("instruction", "paraphrase for a professional report")
    if not current_text:
        raise HTTPException(status_code=400, detail="current_text is required")

    try:
        from google import genai
        client = genai.Client(api_key=os.environ["GEMINI_API_KEY"])
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=(
                f"Rewrite the following text for a professional PowerPoint report. "
                f"Instruction: {instruction}. "
                "Keep all numbers exactly as they are. Use clear, formal business English. "
                "No em dashes, bullets, or markdown. Output plain text only.\n\n"
                + current_text
            ),
        )
        return JSONResponse({"rewritten_text": response.text.strip()})
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Gemini rewrite failed: {e}") from e


@app.post("/reports/{report_id}/apply-edits")
def apply_report_edits(report_id: int, body: dict):
    import json as _json
    from pathlib import Path as _Path
    report = get_report(report_id)
    if not report:
        raise HTTPException(status_code=404, detail="Report not found")
    if not report.get("output_path"):
        raise HTTPException(status_code=400, detail="Report has no output file")

    edits = body.get("edits", {})
    if not edits:
        raise HTTPException(status_code=400, detail="No edits provided")

    original_path = _resolve_report_path(report["output_path"])
    if not original_path.exists():
        raise HTTPException(status_code=404, detail="Original PPTX not found")

    edited_path = original_path.with_stem(original_path.stem + "-edited")

    from .slides import apply_field_edits
    apply_field_edits(original_path, edits, edited_path)

    update_report_completed(report_id, str(edited_path))
    update_report_edits(report_id, _json.dumps(edits))

    try:
        from .slides import render_pdf
        pdf_path = render_pdf(report_id, edited_path)
        update_report_slides_dir(report_id, str(pdf_path.parent))
    except Exception as render_err:
        logger.warning("PDF re-render failed after edits for report_id=%s: %s", report_id, render_err)

    return JSONResponse({
        "output_path": str(edited_path),
        "download_url": f"/reports/{report_id}/download",
    })


_SLUG_RE = re.compile(r'^[a-z0-9_]+$')
_MAX_PPTX_BYTES = 50 * 1024 * 1024  # 50 MB


def _serialize_template(template: dict) -> dict:
    try:
        field_map = json.loads(template.get("field_map") or "[]")
    except Exception:
        field_map = []
    property_sections = list_template_sections(template["id"])
    return {
        **template,
        "field_map": field_map,
        "has_field_map": bool(field_map),
        "property_sections": property_sections,
    }



@app.post("/templates/upload", status_code=201)
async def upload_template(
    file: UploadFile = File(...),
    label: str = Form(...),
    slug: str = Form(...),
):
    # --- Validate inputs ---
    if not file.filename or not file.filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Only .pptx files are supported")
    if not _SLUG_RE.match(slug):
        raise HTTPException(status_code=400, detail="slug must match [a-z0-9_]+")
    if slug in HARDCODED_REPORT_NAMES:
        raise HTTPException(status_code=400, detail=f"slug '{slug}' conflicts with a built-in report name")
    if get_template_by_slug(slug) is not None:
        raise HTTPException(status_code=409, detail=f"A template with slug '{slug}' already exists")

    content = await file.read()
    if len(content) > _MAX_PPTX_BYTES:
        raise HTTPException(status_code=413, detail="File exceeds the 50 MB limit")

    # --- Save the PPTX ---
    templates_dir = get_user_templates_dir()
    pptx_path = templates_dir / f"{slug}.pptx"
    pptx_path.write_bytes(content)

    # --- Extract shapes and slide count ---
    try:
        from pptx import Presentation
        from .slides import extract_all_shapes
        prs = Presentation(str(pptx_path))
        slide_count = len(prs.slides)
        shapes = extract_all_shapes(pptx_path)
    except Exception as exc:
        pptx_path.unlink(missing_ok=True)
        raise HTTPException(status_code=422, detail=f"Could not parse PPTX: {exc}") from exc

    # --- Persist to DB ---
    template_id = create_template(label, slug, str(pptx_path), slide_count)
    upsert_template_shapes(template_id, shapes)

    # Previews are rendered client-side by the frontend (canvas → POST /slides/{i}/image).
    # No background Python render is started here.

    return JSONResponse({"id": template_id, "slug": slug, "slide_count": slide_count}, status_code=201)


@app.get("/templates")
def get_templates():
    return JSONResponse([_serialize_template(template) for template in list_templates()])


@app.get("/templates/report-options")
def get_report_options():
    """Return merged built-in + user template list for the report generation dropdown."""
    options = [
        {"value": name, "label": name.replace("_", " ").title(), "source": "builtin"}
        for name in sorted(HARDCODED_REPORT_NAMES)
    ]
    for t in list_templates():
        options.append({"value": t["slug"], "label": t["label"], "source": "user"})
    return JSONResponse(options)


@app.get("/templates/{template_id}")
def get_template_by_id(template_id: int):
    template = get_template(template_id)
    if not template:
        raise HTTPException(status_code=404, detail="Template not found")
    return JSONResponse(_serialize_template(template))


@app.get("/templates/{template_id}/shapes")
def get_template_shapes(template_id: int):
    template = get_template(template_id)
    if not template:
        raise HTTPException(status_code=404, detail="Template not found")
    shapes = list_template_shapes(template_id)
    preview_dir = template.get("preview_dir")
    slide_count: int = template.get("slide_count") or 0

    # Build a dict for all slides (0..slide_count-1) so the thumbnail panel
    # always shows every slide, not just those that have shapes.
    slides: dict[int, dict] = {}
    for idx in range(slide_count):
        has_image = bool(preview_dir and (Path(preview_dir) / f"slide_{idx}.png").exists())
        slides[idx] = {
            "slide_index": idx,
            "image_url": f"/templates/{template_id}/slides/{idx}/image" if has_image else None,
            "shapes": [],
        }

    # Populate shapes into their slide slots
    for shape in shapes:
        idx = shape["slide_index"]
        if idx not in slides:
            # shape on a slide beyond slide_count (shouldn't happen, but be safe)
            has_image = bool(preview_dir and (Path(preview_dir) / f"slide_{idx}.png").exists())
            slides[idx] = {
                "slide_index": idx,
                "image_url": f"/templates/{template_id}/slides/{idx}/image" if has_image else None,
                "shapes": [],
            }
        slides[idx]["shapes"].append(shape)

    return JSONResponse(sorted(slides.values(), key=lambda s: s["slide_index"]))


@app.get("/templates/{template_id}/slides/{slide_index}/image")
def get_template_slide_image(template_id: int, slide_index: int):
    template = get_template(template_id)
    if not template or not template.get("preview_dir"):
        raise HTTPException(status_code=404, detail="Slide images not available yet")
    image_path = Path(template["preview_dir"]) / f"slide_{slide_index}.png"
    if not image_path.exists():
        raise HTTPException(status_code=404, detail=f"Slide {slide_index} image not found")
    return FileResponse(str(image_path), media_type="image/png")


@app.post("/templates/{template_id}/slides/{slide_index}/image", status_code=200)
async def upload_template_slide_image(template_id: int, slide_index: int, request: Request):
    """Accept a PNG rendered by the frontend and cache it as a preview."""
    template = get_template(template_id)
    if not template:
        raise HTTPException(status_code=404, detail="Template not found")
    slug = template["slug"]
    preview_dir = get_user_templates_dir() / f"{slug}-previews"
    preview_dir.mkdir(parents=True, exist_ok=True)
    body = await request.body()
    if not body:
        raise HTTPException(status_code=400, detail="Empty body")
    (preview_dir / f"slide_{slide_index}.png").write_bytes(body)
    # Once all slides are present, set preview_dir in DB so polling loop sees them
    slide_count: int = template["slide_count"] or 0
    rendered = len(list(preview_dir.glob("slide_*.png")))
    if rendered >= slide_count:
        update_template_preview_dir(template_id, str(preview_dir))
    return JSONResponse({"ok": True, "rendered": rendered, "total": slide_count})


@app.get("/templates/{template_id}/pptx-file")
def get_template_pptx_file(template_id: int):
    """Serve the raw PPTX binary so the frontend can render thumbnails client-side."""
    template = get_template(template_id)
    if not template:
        raise HTTPException(status_code=404, detail="Template not found")
    pptx_path = Path(template["pptx_path"])
    if not pptx_path.exists():
        raise HTTPException(status_code=404, detail="PPTX file not found")
    return FileResponse(
        str(pptx_path),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


@app.put("/templates/{template_id}/config")
def put_template_config(template_id: int, body: dict):
    template = get_template(template_id)
    if not template:
        raise HTTPException(status_code=404, detail="Template not found")
    ga4_property_id = str(body.get("ga4_property_id", "")).strip()
    gsc_url = str(body.get("gsc_url", "")).strip()
    is_seven_slide = bool(body.get("is_seven_slide", False))
    field_map = body.get("field_map", [])
    property_sections = body.get("property_sections", [])
    update_template_config(template_id, ga4_property_id, gsc_url, is_seven_slide, json.dumps(field_map))
    upsert_template_sections(template_id, property_sections)
    updated = get_template(template_id)
    if not updated:
        raise HTTPException(status_code=404, detail="Template not found")
    return JSONResponse(_serialize_template(updated))


@app.post("/templates/{template_id}/rerender-previews", status_code=202)
def rerender_template_previews(template_id: int):
    """Clear the cached preview images and re-render them in the background."""
    template = get_template(template_id)
    if not template:
        raise HTTPException(status_code=404, detail="Template not found")
    pptx_path = Path(template["pptx_path"])
    if not pptx_path.exists():
        raise HTTPException(status_code=404, detail="Template PPTX file not found")
    # Clear preview_dir in DB — frontend will re-render slides client-side via pptxviewjs
    update_template_preview_dir(template_id, "")
    slug = template["slug"]
    preview_dir = get_user_templates_dir() / f"{slug}-previews"
    if preview_dir.exists():
        import shutil as _shutil
        _shutil.rmtree(str(preview_dir), ignore_errors=True)
    return JSONResponse({"queued": True})


@app.delete("/templates/{template_id}")
def delete_template_by_id(template_id: int):
    template = get_template(template_id)
    if not template:
        raise HTTPException(status_code=404, detail="Template not found")
    # Clean up files
    try:
        Path(template["pptx_path"]).unlink(missing_ok=True)
    except Exception:
        pass
    if template.get("preview_dir"):
        try:
            shutil.rmtree(template["preview_dir"], ignore_errors=True)
        except Exception:
            pass
    deleted = delete_template(template_id)
    if not deleted:
        raise HTTPException(status_code=404, detail="Template not found")
    return JSONResponse({"id": template_id, "deleted": True})


@app.post("/templates/ga4-search")
def search_ga4_properties(body: dict):
    """Search GA4 properties via Playwright (runs in the shared executor with 20s timeout)."""
    from concurrent.futures import TimeoutError as FutureTimeout
    query = str(body.get("query", "")).strip()
    if not query:
        raise HTTPException(status_code=400, detail="query is required")

    def _do_search():
        from .generator import _launch_persistent_context, _open_analytics_root
        from playwright.sync_api import sync_playwright
        load_runtime_environment()
        with sync_playwright() as playwright:
            ctx = _launch_persistent_context(playwright)
            try:
                page = ctx.new_page()
                _open_analytics_root(page)
                # Try to find the search input and type the query
                search_input = page.get_by_role("searchbox").first
                search_input.wait_for(state="visible", timeout=8000)
                search_input.click()
                page.wait_for_timeout(300)
                search_input.fill(query)
                page.wait_for_timeout(2000)
                # Collect visible option texts
                results = page.locator('[role="option"], [role="listitem"]').all_text_contents()
                return [r.strip() for r in results if r.strip()]
            finally:
                ctx.close()

    future = _executor.submit(_do_search)
    try:
        results = future.result(timeout=25)
        return JSONResponse({"results": results})
    except FutureTimeout:
        raise HTTPException(status_code=504, detail="GA4 search timed out")
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"GA4 search failed: {exc}") from exc


@app.get("/reports/econet")
def get_econet():
    from .extractors import extract_econet

    return JSONResponse(extract_econet())


@app.get("/reports/econet-ai")
def get_econet_ai():
    from .extractors import extract_econet_ai

    return JSONResponse(extract_econet_ai())


@app.get("/reports/infraco")
def get_infraco():
    from .extractors import extract_infraco

    return JSONResponse(extract_infraco())


@app.get("/reports/ecocash")
def get_ecocash():
    from .extractors import extract_ecocash

    return JSONResponse(extract_ecocash())


@app.get("/reports/ecosure")
def get_ecosure():
    from .extractors import extract_ecosure

    return JSONResponse(extract_ecosure())


@app.get("/reports/zimplats")
def get_zimplats():
    from .extractors import extract_zimplats

    return JSONResponse(extract_zimplats())


@app.get("/reports/union-hardware")
def get_union_hardware():
    from .extractors import extract_union_hardware

    return JSONResponse(extract_union_hardware())


@app.get("/reports/cancer-serve")
def get_cancer_serve():
    from .extractors import extract_cancer_serve

    return JSONResponse(extract_cancer_serve())
