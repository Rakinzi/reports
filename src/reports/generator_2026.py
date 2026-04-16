"""
generator_2026.py — Generation pipeline for 2026 PPTX template format.

Completely separate from the old pipeline in generator.py.
Called by app.py when the requested report name is in TEMPLATES_2026.
"""
from __future__ import annotations

import json
import os
import re
from copy import deepcopy
from pathlib import Path

import lxml.etree as etree
from google import genai
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.oxml.ns import qn

from .charts import generate_country_bar_chart
from .logging_utils import configure_logging
from .runtime import (
    get_output_dir,
    get_screenshots_dir,
    get_templates_dir,
    load_runtime_environment,
)

# Re-use browser/navigation helpers from the old module — they are pure utilities
# with no dependency on the old pipeline's data model.
from .generator import (
    _ensure_expected_ga4_property,
    _fill_text_run,
    _ga4_property_token,
    _ga4_url,
    _goto_ga4_section,
    _launch_persistent_context,
    _replace_image_in_slide,
    _scrape_home_metrics,
    _scrape_snapshot_metrics,
    _set_date_range,
    _switch_ga4_property_via_search,
    _generate_recommendations,
    _get_previous_month_data,
)

logger = configure_logging()

TEMPLATES_DIR = get_templates_dir()
OUTPUT_DIR = get_output_dir()
SCREENSHOTS_DIR = get_screenshots_dir()

# ---------------------------------------------------------------------------
# Template + property maps — 2026 format only
# ---------------------------------------------------------------------------

GA4_PROPERTIES_2026: dict[str, str] = {
    "cancer_serve": "454873082",
    "econet_ai":    "511212348",
    "infraco":      "516617515",
    "zimplats":     "385365994",
    "ecocash":      "386950925",
    "econet":       "386649040",
    "ecosure":      "384507667",
    "dicomm":       "382296904",
}

TEMPLATES_2026: dict[str, str] = {
    "econet":       "new/Econet March Website Report.pptx",
    "econet_ai":    "new/Econet AI March Website Report.pptx",
    "infraco":      "new/Econet Infraco March Website Report.pptx",
    "ecocash":      "new/EcoCash March Website Report.pptx",
    "ecosure":      "new/Ecosure January 2026 Website Report (1).pptx",
    "zimplats":     "new/Zimplats March Website Report.pptx",
    "cancer_serve": "new/Cancerserve March Website Report.pptx",
    "dicomm":       "new/Dicomm March Website Report.pptx",
}

# 7-slide variants skip Slide 6 (Search Performance)
SEVEN_SLIDE_REPORTS = {"zimplats", "dicomm"}

# Google Search Console site URLs — used to build the performance report URL
GSC_URLS: dict[str, str] = {
    "econet":       "https://www.econet.co.zw/",
    "econet_ai":    "https://econetai.co.zw/",
    "infraco":      "https://infraco.co.zw/",
    "ecocash":      "https://www.ecocash.co.zw/",
    "ecosure":      "https://www.ecosure.co.zw/",
    "cancer_serve": "https://www.cancerserve.org/",
    "dicomm":       "https://www.dicomm.co.zw/",
}


# ---------------------------------------------------------------------------
# Gemini text generators
# ---------------------------------------------------------------------------

def _gemini_para(raw: str) -> str:
    """Paraphrase raw text via Gemini, returning the same approximate length."""
    return _gemini_paras_batch([raw])[0]


def _gemini_paras_batch(raws: list[str]) -> list[str]:
    """Paraphrase multiple texts in a single Gemini call. Returns results in the same order.
    Falls back to raw text on any API error so report generation never crashes."""
    load_runtime_environment()
    if not raws:
        return []
    try:
        if len(raws) == 1:
            word_count = len(raws[0].split())
            client = genai.Client(api_key=os.environ["GEMINI_API_KEY"])
            resp = client.models.generate_content(
                model="gemini-2.5-flash",
                contents=(
                    f"Paraphrase the following for a professional PowerPoint report. "
                    f"Keep all numbers, percentages, and proper nouns exactly as they are. "
                    f"Use clear, formal business English. No em dashes, bullets, or markdown. "
                    f"Output must be approximately {word_count} words — do not shorten or expand. "
                    f"Output one plain paragraph only.\n\n" + raws[0]
                ),
            )
            return [resp.text.strip()]

        sections = "\n\n".join(
            f"[{i + 1}] (~{len(r.split())} words)\n{r}" for i, r in enumerate(raws)
        )
        prompt = (
            f"Paraphrase each of the following {len(raws)} numbered texts for a professional PowerPoint report. "
            f"Keep all numbers, percentages, and proper nouns exactly as they are. "
            f"Use clear, formal business English. No em dashes, bullets, or markdown. "
            f"Match each text's approximate word count. "
            f"Output ONLY the paraphrased texts, each preceded by its number in the format [1], [2], etc. "
            f"One plain paragraph per number. No other text.\n\n"
            + sections
        )
        client = genai.Client(api_key=os.environ["GEMINI_API_KEY"])
        resp = client.models.generate_content(model="gemini-2.5-flash", contents=prompt)
        raw_output = resp.text.strip()

        # Parse [1] ... [2] ... sections from the response
        import re as _re
        parts = _re.split(r"\[(\d+)\]", raw_output)
        result_map: dict[int, str] = {}
        for idx in range(1, len(parts), 2):
            num = int(parts[idx])
            text = parts[idx + 1].strip() if idx + 1 < len(parts) else ""
            result_map[num] = text

        return [result_map.get(i + 1, raws[i]) for i in range(len(raws))]
    except Exception as exc:
        logger.warning("[2026] Gemini paraphrase failed, using raw text: %s", exc)
        return list(raws)


def _write_para_with_highlights(para, text: str, bold_words: set[str] | None = None) -> None:
    """Replace paragraph content with text, bolding numbers/percentages and any extra bold_words."""
    import lxml.etree as etree

    if not para.runs:
        return

    # Save formatting from first run
    first_rpr = para.runs[0]._r.find(qn("a:rPr"))
    base_rpr = deepcopy(first_rpr) if first_rpr is not None else None

    # Remove all existing runs (keep pPr)
    for child in list(para._p):
        if child.tag != qn("a:pPr"):
            para._p.remove(child)

    # Build pattern: numbers/percentages + any extra words (e.g. country names)
    num_pat = r"\d[\d,]*(?:\.\d+)?(?:%|K|M|s|ms)?|\d+(?:\.\d+)?%"
    if bold_words:
        # Sort longest first to avoid partial matches (e.g. "South Africa" before "South")
        escaped = [re.escape(w) for w in sorted(bold_words, key=len, reverse=True)]
        token_pattern = re.compile(f"({'|'.join(escaped)}|{num_pat})")
    else:
        token_pattern = re.compile(f"({num_pat})")

    segments = []
    last = 0
    for m in token_pattern.finditer(text):
        if m.start() > last:
            segments.append((text[last:m.start()], False))
        segments.append((m.group(), True))
        last = m.end()
    if last < len(text):
        segments.append((text[last:], False))

    for seg_text, bold in segments:
        r = etree.SubElement(para._p, qn("a:r"))
        if base_rpr is not None:
            rpr = deepcopy(base_rpr)
            if bold:
                rpr.set("b", "1")
            else:
                rpr.attrib.pop("b", None)
            r.insert(0, rpr)
        t = etree.SubElement(r, qn("a:t"))
        t.text = seg_text


def _exec_summary_texts(report_name: str, home_metrics: dict, snapshot_metrics: dict) -> dict:
    """Slide 2 — KPI values + narratives for para 0 and para 2."""
    active_users = home_metrics.get("Active users", "N/A")
    new_users = home_metrics.get("New users", "N/A")
    brand = report_name.replace("_", " ").title()

    # GA4 already returns shorthand like "27K" — use as-is, only convert if raw number
    try:
        s = str(active_users).strip().replace(",", "")
        if s.upper().endswith("K"):
            active_users_short = s
        else:
            n = int(float(s))
            active_users_short = f"{n // 1000}K" if n >= 1000 else str(n)
    except (ValueError, TypeError):
        active_users_short = str(active_users)

    def _parse_num(val: str) -> int:
        s = str(val).strip().replace(",", "")
        if s.upper().endswith("K"):
            return int(float(s[:-1]) * 1000)
        return int(float(s))

    try:
        nu = _parse_num(new_users)
        au = _parse_num(active_users)
        new_pct = f"{round(nu / au * 100, 1)}%" if au > 0 else "N/A"
    except (ValueError, TypeError, ZeroDivisionError):
        new_pct = "N/A"

    new_user_summary = (
        f"of which {new_users} ({new_pct}) were new visitors"
        if new_pct != "N/A"
        else f"of which {new_users} were new visitors"
    )

    # Subtitle — same word count as template
    raw_subtitle = f"Performance Overview: {new_pct} of users are first-time visitors"

    # Para 0 — active users + new visitors count + new visitor %
    raw_para0 = (
        f"The {brand} website delivered solid overall performance during the period under review, "
        f"attracting {active_users} active users, {new_user_summary}. "
        f"This strong proportion of first-time users reflects effective audience acquisition "
        f"and sustained brand visibility across digital channels."
    )

    engagement = home_metrics.get("Average engagement time per active user", "N/A")

    # Para 1 — discovery / returning users insight
    raw_para1_no_gsc = (
        f"The high proportion of new users suggests that the platform is still in a discovery phase, "
        f"with most traffic coming from first-time visitors rather than returning users. "
        f"Building on this momentum through targeted retention strategies will be key to growing a loyal audience."
    )

    # Para 2 — engagement insight (no CTR dependency)
    raw_para2 = (
        f"The average engagement time of {engagement} indicates meaningful interaction with the content. "
        f"Users who visit the platform are spending time engaging rather than immediately exiting, "
        f"which reflects relevant and compelling content despite the size of the audience."
    )

    # Para 3 — closing insight
    raw_para3 = (
        f"Overall, the data reflects strong top-of-funnel performance with effective audience acquisition "
        f"and sustained brand visibility. The key opportunity going forward is to improve retention "
        f"and encourage repeat visits as the platform continues to grow."
    )

    subtitle, para0, para1_no_gsc, para2, para3 = _gemini_paras_batch(
        [raw_subtitle, raw_para0, raw_para1_no_gsc, raw_para2, raw_para3]
    )
    return {
        "active_users_short": active_users_short,
        "new_pct": new_pct,
        "subtitle": subtitle,
        "para0": para0,
        "para1_no_gsc": para1_no_gsc,
        "para2": para2,
        "para3": para3,
    }


def _site_overview_paras(report_name: str, home_metrics: dict, snapshot_metrics: dict) -> tuple[str, str, str, str]:
    """Slide 3 — subtitle + three narrative paragraphs, each same token count as template."""
    active_users = home_metrics.get("Active users", "N/A")
    new_users = home_metrics.get("New users", "N/A")
    engagement = home_metrics.get("Average engagement time per active user", "N/A")
    brand = report_name.replace("_", " ").title()

    def _parse_num(val: str) -> int:
        s = str(val).strip().replace(",", "")
        if s.upper().endswith("K"):
            return int(float(s[:-1]) * 1000)
        return int(float(s))

    try:
        au = _parse_num(active_users)
        nu = _parse_num(new_users)
        new_pct = f"{round(nu / au * 100, 1)}%"
    except (ValueError, TypeError, ZeroDivisionError):
        new_pct = "N/A"

    # Subtitle — "User Engagement Metrics: 34K Strong Active User Baseline"
    raw_subtitle = f"User Engagement Metrics: {active_users} Strong Active User Baseline"

    if new_pct != "N/A":
        new_user_clause = f"Of these, {new_users} users ({new_pct}) were new visitors"
    else:
        new_user_clause = f"Of these, {new_users} users were new visitors"

    raw_para0 = (
        f"The reporting period reflects solid and encouraging performance for the {brand} website, "
        f"with a total of {active_users} active users recorded during the period. "
        f"{new_user_clause}, highlighting continued strong "
        f"brand discovery and the effectiveness of current reach and awareness initiatives."
    )

    raw_para2 = (
        f"The consistently high proportion of new users suggests that marketing efforts, "
        f"search visibility, and broader brand exposure are successfully attracting first-time "
        f"audiences to the platform. This level of new user acquisition indicates that the website "
        f"remains highly discoverable and competitive within its category."
    )

    raw_para4 = (
        f"The average engagement time of {engagement} demonstrates moderate but meaningful interaction. "
        f"While slightly below longer-form corporate engagement benchmarks, this duration suggests "
        f"that users are spending enough time to review key information rather than exiting immediately."
    )

    return tuple(_gemini_paras_batch([raw_subtitle, raw_para0, raw_para2, raw_para4]))


def _geo_paras(countries_data: list[dict], raw_subtitle: str | None = None) -> tuple:
    """Slide 4 — returns 4 Gemini-paraphrased narrative paragraphs from real country data."""
    total_users = sum(r["users"] for r in countries_data) or 1
    sorted_rows = sorted(countries_data, key=lambda r: r["users"], reverse=True)
    top = sorted_rows[0]
    top_pct = round(top["users"] / total_users * 100, 2)

    secondary = sorted_rows[1:5]

    def _pct_val(s: str) -> float:
        try:
            return float(s.strip().rstrip("%"))
        except ValueError:
            return 0.0

    high_intent = [r for r in sorted_rows[1:] if _pct_val(r["engagement_rate"]) > 50]

    raw0 = (
        f"Geographic performance shows {top['country']} as the dominant market, "
        f"contributing {top['users']:,} users ({top_pct}% of total traffic) "
        f"with a {top['engagement_rate']} engagement rate and "
        f"{top['engaged_sessions_per_user']} engaged sessions per user, "
        f"delivering strong volume with steady interaction."
    )

    if secondary:
        sec_parts = [
            f"{r['country']} ({r['users']:,} users) records a {r['engagement_rate']} engagement rate"
            for r in secondary[:2]
        ]
        raw1 = (
            f"{sec_parts[0]}"
            + (f", while {sec_parts[1]}" if len(sec_parts) > 1 else "")
            + ", indicating stronger visit quality despite lower traffic volumes."
        )
    else:
        raw1 = f"Secondary markets are emerging and represent an opportunity for targeted regional campaigns."

    if high_intent:
        hi_parts = [f"{r['country']} ({r['engagement_rate']})" for r in high_intent[:4]]
        raw2 = (
            f"Notably, smaller markets such as {', '.join(hi_parts)} "
            f"demonstrate the strongest engagement efficiency, suggesting high-intent audiences."
        )
    else:
        raw2 = (
            f"Several international markets demonstrate strong engagement efficiency, "
            f"suggesting high-intent audiences with potential conversion value."
        )

    raw3 = (
        f"Overall, {top['country']} drives scale, while select international markets "
        f"show strong engagement depth and potential conversion value."
    )

    raws = ([raw_subtitle] if raw_subtitle else []) + [raw0, raw1, raw2, raw3]
    return tuple(_gemini_paras_batch(raws))


_PAGE_COMPLIANCE_KEYWORDS = (
    "terms and conditions", "ts and cs", "t&c", "privacy policy",
    "terms of service", "disclaimer", "legal", "cookie policy",
)
_PAGE_SUPPORT_KEYWORDS = (
    "contact us", "contact", "support", "help", "faq", "customer experience",
    "customer service", "feedback",
)
_PAGE_HOME_KEYWORDS = (
    "home", "homepage", "index", "main",
)


def _classify_page(title: str) -> str:
    """Classify a page by its title keywords. Returns: 'home', 'compliance', 'support', 'product'."""
    # Classify against the FULL title — brand suffix pattern varies by client:
    # "Devices - Econet Wireless Zimbabwe" (page first) OR
    # "Econet Wireless Zimbabwe - Contact Us" (brand first)
    # So search the whole string for keywords, not just one side.
    t = title.lower()

    if any(kw in t for kw in _PAGE_HOME_KEYWORDS):
        return "home"
    if any(kw in t for kw in _PAGE_COMPLIANCE_KEYWORDS):
        return "compliance"
    if any(kw in t for kw in _PAGE_SUPPORT_KEYWORDS):
        return "support"

    # If no descriptors found and no " - " separator → pure brand name = homepage
    if " - " not in t and "–" not in t and "|" not in t:
        return "home"

    return "product"


def _classify_page_from_row(row: dict) -> str:
    """Classify using the GA4 path first, falling back to title text only when needed."""
    path_value = str(row.get("path", "")).strip().lower()
    title_value = str(row.get("title", "")).strip().lower()
    combined = f"{path_value} {title_value}".strip()

    if path_value in {"", "/"}:
        return "home"
    if any(kw in combined for kw in _PAGE_COMPLIANCE_KEYWORDS):
        return "compliance"
    if any(kw in combined for kw in _PAGE_SUPPORT_KEYWORDS):
        return "support"

    # If the path is a concrete content URL, treat it as a normal content/product page.
    if path_value.startswith("/"):
        return "product"

    return _classify_page(title_value)


def _short_title(title: str) -> str:
    """Return a clean short label: strip brand suffix from either side, normalise T&C variants."""
    _TC_VARIANTS = ("terms and conditions", "ts and cs", "terms & conditions", "terms of service")

    t_lower = title.lower()

    # Normalise T&C first (before stripping brand name)
    for variant in _TC_VARIANTS:
        if variant in t_lower:
            idx = t_lower.find(variant)
            product = title[:idx].strip().rstrip("–-– ").strip()
            # Also strip brand suffix from product part if present
            if " - " in product:
                product = product.rsplit(" - ", 1)[0].strip()
            if product:
                return f"{product} Terms & Conditions"
            return "Terms & Conditions"

    # Strip brand: try both "Page - Brand" and "Brand - Page" patterns.
    # Heuristic: the shorter segment after splitting on " - " is usually the page name.
    if " - " in title:
        parts = title.split(" - ", 1)
        left, right = parts[0].strip(), parts[1].strip()
        # Pick the more descriptive (usually shorter) part as the page label
        label = left if len(left) <= len(right) else right
        return label

    return title


def _switch_dimension_to_page_path(page) -> None:
    selectors = [
        "th.cdk-column-ROW_HEADER-unifiedScreenClass-primaryDimensionColumn button[data-guidedhelpid='table-dimension-picker']",
        "button[data-guidedhelpid='table-dimension-picker']",
    ]
    last_error: Exception | None = None

    for selector in selectors:
        try:
            button = page.locator(selector).first
            button.wait_for(state="visible", timeout=8000)
            button.click()
            page.wait_for_timeout(1500)
            break
        except Exception as exc:
            last_error = exc
    else:
        raise RuntimeError("Could not open the GA4 dimension picker") from last_error

    option_queries = [
        lambda: page.get_by_role("menuitem").filter(has_text="Page path").first,
        lambda: page.locator("button.mat-mdc-menu-item").filter(has_text="Page path").first,
        lambda: page.locator("[role='menuitem']").filter(has_text="Page path").first,
        lambda: page.locator("mat-option, .mat-mdc-option").filter(has_text="Page path").first,
        lambda: page.get_by_text("Page path and screen class").first,
    ]
    for query in option_queries:
        try:
            option = query()
            option.wait_for(state="visible", timeout=5000)
            option.click()
            page.wait_for_timeout(3000)
            page.locator("th.cdk-column-__row_index__").first.wait_for(state="visible", timeout=10000)
            return
        except Exception:
            continue

    raise RuntimeError("Could not switch the GA4 table dimension to 'Page path and screen class'")


def _fallback_page_label(path_value: str) -> str:
    clean = path_value.strip() or "/"
    if clean == "/":
        return "Homepage"
    clean = clean.split("?", 1)[0].split("#", 1)[0].strip("/")
    if not clean:
        return "Homepage"

    parts = [part for part in clean.split("/") if part]
    last = parts[-1].replace("-", " ").replace("_", " ").strip()
    if not last:
        return "Homepage"
    words = [word for word in last.split() if word]
    return " ".join(word.upper() if word.isupper() else word.capitalize() for word in words[:5])


def _label_page_paths_with_gemini(pages_data: list[dict]) -> None:
    if not pages_data:
        return

    load_runtime_environment()
    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        for row in pages_data:
            row["title"] = _fallback_page_label(row.get("path", row["title"]))
        return

    client = genai.Client(api_key=api_key)
    path_list = "\n".join(f"- {row.get('path', row['title'])}" for row in pages_data)
    prompt = (
        "You are labeling website page paths for a PowerPoint report.\n"
        "Return strict JSON only: an array of objects with keys path and label.\n"
        "Rules:\n"
        "- label must be 2 to 5 words\n"
        "- use title case\n"
        "- make each label human-friendly\n"
        "- keep important brand/product names when obvious\n"
        "- '/' should become 'Homepage'\n"
        "- do not include commentary or markdown\n\n"
        f"Paths:\n{path_list}"
    )
    try:
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt,
        )
        raw = (response.text or "").strip()
        mappings = json.loads(raw)
        by_path = {
            item["path"]: item["label"].strip()
            for item in mappings
            if isinstance(item, dict) and item.get("path") and item.get("label")
        }
        for row in pages_data:
            path_value = row.get("path", row["title"])
            row["title"] = by_path.get(path_value, _fallback_page_label(path_value))
    except Exception:
        for row in pages_data:
            row["title"] = _fallback_page_label(row.get("path", row["title"]))


def _page_perf_paras(pages_data: list[dict], site_total_views: int = 0) -> tuple[str, str, str, str, str, str]:
    """Slide 5 — heading + 5 narrative paragraphs from real page data."""
    if not pages_data:
        return ("", "", "", "", "", "")

    for p in pages_data:
        p["_type"] = _classify_page_from_row(p)
        raw_label = _fallback_page_label(p.get("path", p["title"]))
        p["_label"] = "Homepage" if p["_type"] == "home" else raw_label

    top = pages_data[0]
    secondary_pages = [p for p in pages_data[1:] if p.get("_label") != "Homepage"]
    second = secondary_pages[0] if secondary_pages else (pages_data[1] if len(pages_data) > 1 else None)
    compliance_pages = [p for p in pages_data if p["_type"] == "compliance"]
    support_pages = [p for p in pages_data if p["_type"] == "support"]

    raw_heading = (
        f"The {top['_label']}"
        + (f" and {second['_label']}" if second else "")
        + " pages drive the most meaningful engagement."
    )

    top_label_display = top['_label'] if top['_label'] != "Homepage" else "Homepage"
    page_suffix = "" if top['_label'] == "Homepage" else " page"
    raw_para1 = (
        f"The {top_label_display}{page_suffix} drives the highest traffic, accounting for "
        f"{top.get('views_pct') or '0%'} of total views ({top['views']:,} views), with "
        f"{top['views_per_user']} views per active user and an average engagement duration of "
        f"{top['avg_engagement_time']}. This substantial volume indicates strong content discoverability and user interest."
    )

    if compliance_pages:
        c_parts = [f"{p['_label']} ({p.get('views_pct') or '0%'} of views)" for p in compliance_pages[:2]]
        raw_para2 = (
            f"Compliance pages, including {' and '.join(c_parts)}, yield considerable traffic, "
            f"predominantly influenced by campaign landing directives over organic browsing."
        )
    else:
        secondary_labels = [p["_label"] for p in secondary_pages[:2]]
        secondary_phrase = " and ".join(secondary_labels) if secondary_labels else "other pages"
        raw_para2 = (
            f"Secondary pages such as {secondary_phrase} contribute meaningful traffic, "
            f"but with more limited evidence of deep exploratory browsing."
        )

    if support_pages:
        s_parts = [p["_label"] for p in support_pages[:2]]
        raw_para3 = (
            f"Support-focused pages including {' and '.join(s_parts)} attract consistent traffic "
            f"but exhibit lower engagement depth, indicating they are mainly used for quick access "
            f"to assistance rather than extended browsing."
        )
    else:
        raw_para3 = (
            f"The leading pages drive visibility and entry traffic, while the remaining pages play "
            f"a more supportive role in the user journey."
        )

    # Para 4 — highest engagement depth pages (by views_per_user)
    def _parse_float(val: str) -> float:
        try:
            return float(str(val).replace(",", ""))
        except (ValueError, TypeError):
            return 0.0

    deep_pages = sorted(
        [p for p in pages_data if _parse_float(p.get("views_per_user", "0")) > 1.0],
        key=lambda p: _parse_float(p.get("views_per_user", "0")),
        reverse=True,
    )[:3]
    if deep_pages:
        deep_parts = [
            f"{p['_label']} ({p['views_per_user']} views per user)"
            for p in deep_pages
        ]
        raw_para4 = (
            f"Pages with the deepest browsing engagement include {', '.join(deep_parts)}, "
            f"indicating stronger exploration intent and higher content value for those visitors."
        )
    else:
        raw_para4 = (
            f"Product and service pages show deeper per-user interaction, indicating visitors "
            f"exploring core offerings are more engaged than those landing on support or contact pages."
        )

    # Para 5 — overall summary
    total_label = f"{site_total_views:,} total views" if site_total_views else "total views recorded"
    raw_para5 = (
        f"Overall, with {total_label} across all pages, traffic is largely driven by core solution "
        f"pages with strong top-of-funnel acquisition, while opportunities remain to improve engagement "
        f"depth and guide users more effectively through key conversion journeys."
    )

    return raw_heading, raw_para1, raw_para2, raw_para3, raw_para4, raw_para5


# ---------------------------------------------------------------------------
# GSC capture helper
# ---------------------------------------------------------------------------

def _capture_gsc(context, report_name: str, start_date: str, end_date: str, out_dir: Path) -> tuple[dict, dict]:
    """Open GSC, select the correct property, set custom date range, scrape metrics + screenshot."""
    from datetime import datetime as _dt
    from urllib.parse import quote as _quote

    start_dt = _dt.strptime(start_date, "%b %d, %Y")
    end_dt   = _dt.strptime(end_date,   "%b %d, %Y")
    gsc_site = GSC_URLS[report_name]

    # Format dates for the GSC date picker inputs: "3/1/2026" (no leading zeros)
    start_picker = f"{start_dt.month}/{start_dt.day}/{start_dt.year}"
    end_picker   = f"{end_dt.month}/{end_dt.day}/{end_dt.year}"

    search_metrics: dict = {}
    screenshots: dict = {}

    gsc_page = context.new_page()

    # 1. Go to GSC performance page for this property directly via URL
    gsc_page.goto(
        "https://search.google.com/search-console/performance/search-analytics"
        f"?resource_id={_quote(gsc_site, safe='')}",
        wait_until="domcontentloaded", timeout=30000,
    )
    gsc_page.wait_for_selector("text=Total clicks", state="attached", timeout=20000)
    gsc_page.wait_for_timeout(1000)

    # 2. Click the property selector pill to confirm/switch property
    try:
        gsc_page.locator('div[jscontroller="Jdbz6e"]').click()
        gsc_page.wait_for_timeout(800)

        # Use the active combobox input — exact=True to avoid matching "Inspect any URL in..."
        prop_input = gsc_page.get_by_role("combobox", name=gsc_site, exact=True)
        prop_input.click(click_count=3)
        prop_input.fill(gsc_site)
        gsc_page.wait_for_timeout(800)

        # Click the matching result that appears in the dropdown list
        gsc_page.locator(f'[data-initialvalue="{gsc_site}"]').first.click()
        gsc_page.wait_for_timeout(2000)
        gsc_page.wait_for_selector("text=Total clicks", state="attached", timeout=10000)
    except Exception as e:
        logger.warning("[2026] GSC property selector failed (continuing): %s", e)

    # 3. Set custom date range.
    #    Some GSC instances show "Custom" directly in the toolbar; others hide it under "More".
    #    Both open the same "Date range" modal with YYYY-MM-DD inputs.
    try:
        start_val = start_dt.strftime("%Y-%m-%d")
        end_val   = end_dt.strftime("%Y-%m-%d")

        # Try clicking "Custom" directly first; fall back to "More time ranges" → "Custom"
        try:
            gsc_page.locator('button[role="radio"]').filter(has_text="Custom").click(timeout=3000)
        except Exception:
            gsc_page.get_by_role("button", name="More time ranges").click(timeout=5000)
            gsc_page.wait_for_timeout(500)
            # Two "Custom" labels exist (Filter + Compare tabs) — click the Filter one
            gsc_page.get_by_label("Filter", exact=True).get_by_text("Custom", exact=True).click(timeout=5000)

        gsc_page.wait_for_timeout(600)

        # Fill the date inputs inside the modal.
        # The inputs have class "qdOxv-fmcmS-wGMbrd" and use aria-labelledby.
        # Wait for them to appear, then fill by order (start=0, end=1).
        gsc_page.wait_for_selector("input.qdOxv-fmcmS-wGMbrd", timeout=8000)
        date_inputs = gsc_page.locator("input.qdOxv-fmcmS-wGMbrd")

        date_inputs.nth(0).click(click_count=3)
        date_inputs.nth(0).fill(start_val)
        gsc_page.wait_for_timeout(200)
        date_inputs.nth(1).click(click_count=3)
        date_inputs.nth(1).fill(end_val)
        gsc_page.wait_for_timeout(200)

        # Click Apply — page reloads with new date range
        gsc_page.get_by_role("button", name="Apply").click()
        gsc_page.wait_for_timeout(5000)
    except Exception as e:
        logger.warning("[2026] GSC date picker failed, reloading with URL date params: %s", e)
        # Fallback: open a fresh page with date params baked into URL
        try:
            gsc_page.close()
        except Exception:
            pass
        gsc_page = context.new_page()
        gsc_page.goto(
            "https://search.google.com/search-console/performance/search-analytics"
            f"?resource_id={_quote(gsc_site, safe='')}"
            f"&start_date={start_dt.strftime('%Y%m%d')}"
            f"&end_date={end_dt.strftime('%Y%m%d')}",
            wait_until="domcontentloaded", timeout=30000,
        )
        gsc_page.wait_for_selector("text=Total clicks", state="attached", timeout=15000)
        gsc_page.wait_for_timeout(2000)

    # 5. Scrape metrics
    body_text = gsc_page.locator("body").inner_text()
    lines = [l.strip() for l in body_text.splitlines()]

    def _next_val_after(label: str) -> str:
        for i, line in enumerate(lines):
            if line.lower() == label.lower():
                for j in range(i + 1, min(i + 4, len(lines))):
                    if lines[j]:
                        return lines[j]
        return "N/A"

    search_metrics = {
        "impressions":  _next_val_after("Total impressions"),
        "clicks":       _next_val_after("Total clicks"),
        "ctr":          _next_val_after("Average CTR"),
        "avg_position": _next_val_after("Average position"),
    }

    # 6. Screenshot: metric cards + chart section
    try:
        gsc_page.mouse.move(0, 0)
        gsc_page.wait_for_timeout(500)
        path = out_dir / "search_console.png"

        gsc_page.evaluate("window.scrollTo(0, 0)")
        gsc_page.wait_for_timeout(500)

        clip = gsc_page.evaluate("""() => {
            // Find any leaf or near-leaf element whose trimmed text is 'Total clicks'
            const all = Array.from(document.querySelectorAll('*'));
            const label = all.find(el =>
                el.textContent.trim() === 'Total clicks' &&
                el.getBoundingClientRect().width > 0
            );
            if (!label) return null;
            const vw = window.innerWidth;
            let el = label;
            for (let i = 0; i < 20; i++) {
                el = el.parentElement;
                if (!el) break;
                const r = el.getBoundingClientRect();
                if (r.width >= vw * 0.6 && r.height >= 400) {
                    return { x: r.x, y: r.y, width: r.width, height: r.height };
                }
            }
            return null;
        }""")

        if clip and clip["height"] > 50:
            gsc_page.screenshot(path=str(path), clip=clip)
        else:
            # Fallback: screenshot just the top portion of the viewport
            gsc_page.screenshot(path=str(path), clip={"x": 0, "y": 150, "width": 1920, "height": 550})

        screenshots["search_screenshot"] = path
    except Exception:
        pass

    gsc_page.close()
    return search_metrics, screenshots


# ---------------------------------------------------------------------------
# GA4 capture — 2026 pipeline only
# ---------------------------------------------------------------------------

def capture_2026(
    report_name: str,
    start_date: str,
    end_date: str,
    _stage_callback=None,
) -> tuple[dict[str, Path], dict, dict, dict]:
    """Navigate GA4, capture all screenshots and metrics needed for the 2026 pipeline.

    Returns: (screenshots, home_metrics, snapshot_metrics, page_views)
    """
    def _stage(msg: str):
        if _stage_callback:
            _stage_callback(msg)

    screenshots: dict[str, Path] = {}
    out_dir = SCREENSHOTS_DIR / report_name
    out_dir.mkdir(parents=True, exist_ok=True)

    home_metrics: dict = {}
    snapshot_metrics: dict = {}
    page_views: dict = {}
    pages_data: list[dict] = []  # rich per-page rows: {title, views, active_users, views_per_user, avg_engagement_time}
    site_total_views: int = 0   # total views across ALL pages (from GA4 Total row)
    countries_data: list[dict] = []  # rich per-country rows: {country, users, engagement_rate, engaged_sessions_per_user}
    search_metrics: dict = {}   # GSC metrics: impressions, clicks, ctr, avg_position

    with sync_playwright() as p:
        context = _launch_persistent_context(p, headless=False)
        try:
            page = context.new_page()
            page.bring_to_front()

            # --- Switch to correct GA4 property ---
            _stage("Switching GA4 property...")
            page = _switch_ga4_property_via_search(page, report_name)

            # --- Home page: set date range + scrape metrics ---
            _stage("Capturing GA4 home metrics...")
            page = _goto_ga4_section(page, report_name, "/home")
            try:
                page.keyboard.press("Escape")
                page.wait_for_timeout(1000)
            except Exception:
                pass
            _set_date_range(page, start_date, end_date)
            home_metrics = _scrape_home_metrics(page)

            # --- Home line chart screenshot (Slide 3) ---
            try:
                chart_el = page.locator("ga-card.card_0 ga-tab-chart")
                chart_el.wait_for(state="visible", timeout=10000)
                page.mouse.move(0, 0)
                page.mouse.click(0, 0)
                page.wait_for_timeout(800)
                path = out_dir / "home_chart.png"
                chart_el.screenshot(path=str(path))
                screenshots["home_chart"] = path
            except Exception:
                pass

            # --- Navigate to Reports Snapshot via the confirmed button ---
            _stage("Capturing GA4 snapshot metrics...")
            _ensure_expected_ga4_property(page, report_name)
            page.locator("span.view-link-text", has_text="View reports snapshot").click()
            page.wait_for_timeout(4000)
            _ensure_expected_ga4_property(page, report_name)

            # Set date range on snapshot page
            _set_date_range(page, start_date, end_date)
            page.wait_for_timeout(3000)

            snapshot_metrics = _scrape_snapshot_metrics(page)

            # --- Snapshot KPI card screenshot (Slide 1) ---
            try:
                card_el = page.locator("ga-card[data-guidedhelpid='summary']").first
                card_el.wait_for(state="visible", timeout=10000)
                page.mouse.move(0, 0)
                page.mouse.click(0, 0)
                page.wait_for_timeout(800)
                path = out_dir / "snapshot_card.png"
                card_el.screenshot(path=str(path))
                screenshots["snapshot_card"] = path
            except Exception:
                pass

            # --- Countries table: screenshot + scrape rich data (Slide 4) ---
            _stage("Capturing countries data...")
            try:
                page.get_by_text("View countries").click()
                page.wait_for_timeout(4000)
                _ensure_expected_ga4_property(page, report_name)
                row_num_col = page.locator("th.cdk-column-__row_index__").first
                end_col = page.locator("th.cdk-column-DEFAULT-engagedSessionsPerUser").first
                row_num_col.wait_for(state="visible", timeout=10000)
                page.keyboard.press("End")
                page.wait_for_timeout(1000)
                page.mouse.wheel(0, 3000)
                page.wait_for_timeout(1000)
                start_box = row_num_col.bounding_box()
                end_box = end_col.bounding_box()
                table_box = page.locator("table.adv-table").bounding_box()
                clip = {
                    "x": start_box["x"],
                    "y": table_box["y"],
                    "width": (end_box["x"] + end_box["width"]) - start_box["x"],
                    "height": table_box["height"],
                }
                path = out_dir / "countries_table.png"
                page.screenshot(path=str(path), clip=clip, full_page=True)
                screenshots["countries_table"] = path

                # Scrape rich country data from table rows
                # GA4 table inner_text rows: "N\tCountry\tActive Users\tNew Users\tEngaged Sessions\tEngagement Rate\tEng.Sessions/User\tAvg Eng Time"
                import re as _re
                body = page.locator("body").inner_text()
                # Rows format (leading tab + index):
                # "\t1\tZimbabwe\t21,735 (80.66%)\t19,675 (79.45%)\t15,135 (80.46%)\t47.98%\t0.70\t..."
                for line in body.splitlines():
                    m = _re.match(
                        r"^\t\d+\t(.+?)\t([\d,]+)\s*\([^)]+\)\t([\d,]+)\s*\([^)]+\)\t([\d,]+)\s*\([^)]+\)\t([\d.]+%)\t([\d.]+)",
                        line,
                    )
                    if m:
                        countries_data.append({
                            "country": m.group(1).strip(),
                            "users": int(m.group(2).replace(",", "")),
                            "new_users": int(m.group(3).replace(",", "")),
                            "engaged_sessions": int(m.group(4).replace(",", "")),
                            "engagement_rate": m.group(5),
                            "engaged_sessions_per_user": m.group(6),
                        })

                page.go_back()
                page.wait_for_timeout(3000)
                _ensure_expected_ga4_property(page, report_name)
            except Exception:
                pass

            # --- Pages and screens: screenshot + scrape page views (Slide 5) ---
            _stage("Capturing pages & screens data...")
            try:
                page.get_by_role("button", name="View pages and screens", exact=True).click()
                page.wait_for_timeout(4000)
                _ensure_expected_ga4_property(page, report_name)
                row_num_col = page.locator("th.cdk-column-__row_index__").first
                end_col = page.locator("th.cdk-column-DEFAULT-userEngagementDurationPerUser").first
                table = page.locator("table.adv-table").first
                row_num_col.wait_for(state="visible", timeout=10000)
                page.keyboard.press("End")
                page.wait_for_timeout(1000)
                page.mouse.wheel(0, 3000)
                page.wait_for_timeout(1000)
                _switch_dimension_to_page_path(page)
                start_box = row_num_col.bounding_box()
                end_box = end_col.bounding_box()
                table_box = table.bounding_box()
                clip = {
                    "x": start_box["x"],
                    "y": table_box["y"],
                    "width": (end_box["x"] + end_box["width"]) - start_box["x"],
                    "height": table_box["height"],
                }
                path = out_dir / "pages_table.png"
                page.screenshot(path=str(path), clip=clip, full_page=True)
                screenshots["pages_table"] = path

                page_views, pages_data, site_total_views = _scrape_pages_table(page)
                _label_page_paths_with_gemini(pages_data)
            except Exception:
                pass

            # --- Google Search Console: scrape metrics + screenshot (Slide 6) ---
            if report_name not in SEVEN_SLIDE_REPORTS and report_name in GSC_URLS:
                _stage("Capturing Google Search Console data...")
                try:
                    search_metrics, gsc_screenshots = _capture_gsc(
                        context, report_name, start_date, end_date, out_dir
                    )
                    screenshots.update(gsc_screenshots)
                except Exception as e:
                    logger.warning("[2026] GSC capture failed for %s: %s", report_name, e)

        finally:
            try:
                context.close()
            except Exception:
                pass

    return screenshots, home_metrics, snapshot_metrics, page_views, pages_data, site_total_views, countries_data, search_metrics


# ---------------------------------------------------------------------------
# PPTX builder helpers
# ---------------------------------------------------------------------------

def _performance_month(date_range: str) -> str:
    """Extract 'Month,YYYY' from date_range e.g. '1 February 2026 - 28 February 2026' -> 'February,2026'"""
    match = re.search(r"[A-Za-z]+ \d{4}", date_range)
    return match.group(0).replace(" ", ",") if match else ""


# Per-template picture name for the KPI card slot on Slide 1 (right-side image).
# "Picture 11" is always the footer logo — the KPI card is the other right-side picture.
_SLIDE1_KPI_CARD_PICTURE: dict[str, str] = {
    "econet":       "Picture 10",
    "econet_ai":    "Picture 14",
    "infraco":      "Picture 13",
    "ecocash":      "Picture 13",
    "zimplats":     "Picture 15",
    "cancer_serve": "Picture 12",
    "dicomm":       "Picture 10",
}


def _build_slide1(slide, performance_month: str, screenshots: dict, report_name: str = "") -> None:
    """Replace date text and KPI card screenshot on Slide 1."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            if re.match(r"^[A-Za-z]+,?\s*\d{4}$", para.text.strip()):
                _fill_text_run(para, performance_month)

    if "snapshot_card" in screenshots:
        pic_name = _SLIDE1_KPI_CARD_PICTURE.get(report_name)
        if pic_name:
            _replace_image_in_slide(slide, screenshots["snapshot_card"], shape_name=pic_name)
        else:
            # Fallback: replace the largest right-side picture (left > 5 inches)
            EMU = 914400
            right_pics = [
                s for s in slide.shapes
                if s.shape_type == 13 and s.left > 5 * EMU and s.name != "Picture 11"
            ]
            if right_pics:
                _replace_image_in_slide(slide, screenshots["snapshot_card"], shape_name=right_pics[0].name)


def _build_slide2(slide, home_metrics: dict, snapshot_metrics: dict, report_name: str, search_metrics: dict | None = None) -> None:
    """Replace KPI stat boxes and narrative on Slide 2 (Executive Summary)."""
    exec_texts = _exec_summary_texts(report_name, home_metrics, snapshot_metrics)
    ctr = (search_metrics or {}).get("ctr", "N/A")
    impressions = (search_metrics or {}).get("impressions", "N/A")
    avg_position = (search_metrics or {}).get("avg_position", "N/A")

    raw_para1 = (
        f"From a search visibility perspective, the website recorded {impressions} impressions "
        f"with a {ctr} click-through rate and an average search position of {avg_position}, "
        f"reflecting strong organic discoverability and continued relevance in search results."
    )
    para1_text = _gemini_para(raw_para1) if ctr != "N/A" else ""

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        # KPI stat boxes — target by shape name + paragraph index (para 0 = value, para 1 = label)
        # object 9 = Active Users, object 10 = New Visitors %, object 11 = CTR
        if shape.name == "object 9":
            paras = shape.text_frame.paragraphs
            if len(paras) > 0 and paras[0].text.strip():
                _fill_text_run(paras[0], exec_texts["active_users_short"])
            continue
        if shape.name == "object 10":
            paras = shape.text_frame.paragraphs
            if len(paras) > 0 and paras[0].text.strip():
                _fill_text_run(paras[0], exec_texts["new_pct"])
            continue
        if shape.name == "object 11":
            paras = shape.text_frame.paragraphs
            if ctr != "N/A" and len(paras) > 0 and paras[0].text.strip():
                _fill_text_run(paras[0], ctr)
            continue

        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            # Subtitle line
            if text.lower().startswith("performance overview:"):
                _fill_text_run(para, exec_texts["subtitle"])
            # Para 0 — active users + new visitors acquisition summary
            elif any(kw in text.lower() for kw in (
                "under review", "attracting", "recorded", "active users",
                "new visitors", "highlighting", "early stage visibility",
                "first-time users", "strong proportion",
            )):
                _write_para_with_highlights(para, exec_texts["para0"])
            # Para 1 — discovery phase / returning users (no GSC)
            elif any(kw in text.lower() for kw in (
                "discovery phase", "returning users", "high proportion of new",
                "first time visitors", "first-time visitors",
            )):
                _write_para_with_highlights(para, exec_texts["para1_no_gsc"])
            # Para — CTR / search visibility (from GSC, overwrites discovery para if GSC available)
            elif para1_text and any(kw in text.lower() for kw in (
                "click-through", "search visibility", "organic", "ctr", "impressions",
                "search perspective", "search console",
            )):
                _write_para_with_highlights(para, para1_text)
            # Para — engagement time
            elif any(kw in text.lower() for kw in (
                "engagement time", "minute", "seconds", "meaningful interaction",
                "average engagement", "spending time engaging",
            )):
                _write_para_with_highlights(para, exec_texts["para2"])
            # Para — closing / overall
            elif any(kw in text.lower() for kw in (
                "overall", "data reflects", "top of funnel", "key opportunity",
                "combination of high", "discoverability", "brand loyalty",
            )):
                _write_para_with_highlights(para, exec_texts["para3"])


def _build_slide3(slide, home_metrics: dict, snapshot_metrics: dict, report_name: str, screenshots: dict) -> None:
    """Replace stat values, subtitle, narratives, and chart on Slide 3 (Site Overview)."""
    active_users = home_metrics.get("Active users", "N/A")
    new_users = home_metrics.get("New users", "N/A")
    engagement = home_metrics.get("Average engagement time per active user", "N/A")

    def _parse_num(val: str) -> int:
        s = str(val).strip().replace(",", "")
        if s.upper().endswith("K"):
            return int(float(s[:-1]) * 1000)
        return int(float(s))

    try:
        au = _parse_num(active_users)
        nu = _parse_num(new_users)
        new_pct = f"{round(nu / au * 100, 1)}%"
    except (ValueError, TypeError, ZeroDivisionError):
        new_pct = "N/A"

    subtitle, para0, para2, para4 = _site_overview_paras(report_name, home_metrics, snapshot_metrics)

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        # KPI stat boxes — target by shape name + paragraph index to avoid stale-text matching.
        # Each box has: para 0 = value, para 1 = label (e.g. "Total Active Users")
        if shape.name == "object 6":
            # Total Active Users: replace value (para 0), leave label (para 1) as-is
            paras = shape.text_frame.paragraphs
            if len(paras) > 0 and paras[0].text.strip():
                _fill_text_run(paras[0], str(active_users))
            continue

        if shape.name == "object 10":
            # New Users: para 0 = count, para 1 = "New Users (XX%)" label — update both
            paras = shape.text_frame.paragraphs
            if len(paras) > 0 and paras[0].text.strip():
                _fill_text_run(paras[0], str(new_users))
            if len(paras) > 1 and paras[1].text.strip():
                label = f"New Users ({new_pct})" if new_pct != "N/A" else "New Users"
                _fill_text_run(paras[1], label)
            continue

        if shape.name == "object 14":
            # Avg Engagement Time: para 0 = time value, para 1 = "Avg Engagement Time" label
            paras = shape.text_frame.paragraphs
            if len(paras) > 0 and paras[0].text.strip():
                _fill_text_run(paras[0], str(engagement))
            continue

        for para in shape.text_frame.paragraphs:
            text = para.text.strip()

            # Subtitle — Gemini paraphrase
            if "user engagement metrics" in text.lower():
                _fill_text_run(para, subtitle)

            # Narrative paragraphs — Gemini paraphrase with highlights
            elif any(kw in text.lower() for kw in (
                "reflects solid", "active users recorded", "encouraging performance",
                "active users", "new visitors", "recorded during",
            )):
                _write_para_with_highlights(para, para0)
            elif any(kw in text.lower() for kw in (
                "consistently high proportion", "search visibility", "first-time audiences",
                "marketing efforts", "growth and discovery",
            )):
                _write_para_with_highlights(para, para2)
            elif any(kw in text.lower() for kw in (
                "average engagement time", "engagement benchmarks", "exiting immediately",
                "meaningful interaction", "spending time",
            )):
                _write_para_with_highlights(para, para4)

    # Replace the snapshot card screenshot — each template uses either Picture 18 or Picture 19
    if "snapshot_card" in screenshots:
        picture_names = {
            shape.name
            for shape in slide.shapes
            if getattr(shape, "shape_type", None) == 13
        }
        for candidate in ("Picture 18", "Picture 19"):
            if candidate in picture_names:
                _replace_image_in_slide(slide, screenshots["snapshot_card"], shape_name=candidate)
                break


def _build_slide4(slide, countries_data: list[dict], screenshots: dict) -> None:
    """Replace country table screenshot and narrative on Slide 4 (Geographic Performance)."""
    if countries_data:
        sorted_rows = sorted(countries_data, key=lambda r: r["users"], reverse=True)
        top = sorted_rows[0]

        # Subtitle + 4 narrative paragraphs — all batched in one Gemini call via _geo_paras
        raw_subtitle = f"{top['country']} Dominance with International Opportunity"
        subtitle, para0, para1, para2, para3 = _geo_paras(countries_data, raw_subtitle=raw_subtitle)
        narratives = [para0, para1, para2, para3]

        # Country names to bold in narratives
        country_names = {r["country"] for r in countries_data}

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            # Subtitle: object 2, para index 1
            if shape.name == "object 2":
                paras = shape.text_frame.paragraphs
                if len(paras) > 1:
                    _fill_text_run(paras[1], subtitle)

            # Narratives: object 6, non-empty paragraphs (even-indexed, odd ones are spacers)
            elif shape.name == "object 6":
                content_paras = [p for p in shape.text_frame.paragraphs if p.text.strip()]
                for i, para in enumerate(content_paras):
                    if i < len(narratives):
                        _write_para_with_highlights(para, narratives[i], bold_words=country_names)

    # Use the captured countries table screenshot from GA4.
    if "countries_table" in screenshots:
        picture_names = {
            shape.name
            for shape in slide.shapes
            if getattr(shape, "shape_type", None) == 13
        }
        for candidate in ("Picture 10", "Picture 9"):
            if candidate in picture_names:
                _replace_image_in_slide(slide, screenshots["countries_table"], shape_name=candidate)
                break


def _build_slide5(slide, pages_data: list[dict], screenshots: dict, site_total_views: int = 0) -> None:
    """Replace pages table screenshot and narratives on Slide 5 (Page Performance)."""
    if pages_data:
        heading, para1, para2, para3, para4, para5 = _page_perf_paras(pages_data, site_total_views)
        # Collect clean page labels for bolding (classification already ran inside _page_perf_paras)
        page_names = {p.get("_label", p["title"].split(" - ")[0].strip()) for p in pages_data}

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            # object 3 — short heading line
            if shape.name == "object 3":
                content_paras = [p for p in shape.text_frame.paragraphs if p.text.strip()]
                if content_paras:
                    _fill_text_run(content_paras[0], heading)

            # object 7 — "Overall Insight:" label + up to 5 narrative paragraphs
            elif shape.name == "object 7":
                # All paragraphs including spacers; content slots sit at odd indices (1,3,5,7,9)
                all_paras = shape.text_frame.paragraphs
                narratives = [para1, para2, para3, para4, para5]
                # Odd-indexed paragraphs are the content slots (index 0 is the label)
                content_slots = [p for i, p in enumerate(all_paras) if i % 2 == 1]
                for i, para in enumerate(content_slots):
                    if i < len(narratives):
                        _write_para_with_highlights(para, narratives[i], bold_words=page_names)
                    else:
                        # Clear any leftover text from a previous run in unused slots
                        for run in para.runs:
                            run.text = ""

    if "pages_table" in screenshots:
        picture_names = {
            shape.name
            for shape in slide.shapes
            if getattr(shape, "shape_type", None) == 13
        }
        for candidate in ("Picture 11", "Picture 10", "Picture 13"):
            if candidate in picture_names:
                _replace_image_in_slide(slide, screenshots["pages_table"], shape_name=candidate)
                break


def _search_perf_paras(search_metrics: dict) -> tuple[str, str, str, str]:
    """Slide 6 — subtitle + 4 narrative paragraphs from real GSC data."""
    impressions = search_metrics.get("impressions", "N/A")
    clicks      = search_metrics.get("clicks", "N/A")
    ctr         = search_metrics.get("ctr", "N/A")
    position    = search_metrics.get("avg_position", "N/A")

    # Subtitle — dynamic based on CTR
    def _pct_val(s: str) -> float:
        try:
            return float(str(s).strip().rstrip("%"))
        except ValueError:
            return 0.0

    ctr_val = _pct_val(ctr)
    if ctr_val >= 5:
        subtitle_raw = f"Strong Search Visibility with {ctr} Click-Through Rate"
    elif ctr_val >= 3:
        subtitle_raw = f"Solid Foundation with Room for Visibility Improvement"
    else:
        subtitle_raw = f"Search Presence Established — CTR Optimisation Opportunity"

    raw_para0 = (
        f"With {impressions} impressions, the brand maintains substantial presence in search results, "
        f"indicating strong keyword coverage and consistent discoverability. "
        f"From this visibility, the site generated {clicks} clicks, "
        f"reflecting meaningful traffic acquisition directly from search engines."
    )

    raw_para1 = (
        f"The {ctr} click-through rate suggests "
        f"{'strong' if ctr_val >= 5 else 'moderate'} conversion of impressions into clicks. "
        f"{'This reflects compelling meta titles and descriptions that resonate with searcher intent.' if ctr_val >= 5 else 'There is room to optimise meta titles and descriptions to further improve click appeal and capture a larger share of search demand.'}"
    )

    raw_para2 = (
        f"An average position of {position} is "
        f"{'particularly encouraging, placing the website on the first page of search results for many queries. This reinforces strong SEO foundations and competitive ranking strength.' if _pct_val(position) <= 10 else 'an indication that the website is gaining search traction. Focused content and technical SEO efforts can improve ranking further toward the first page.'}"
    )

    raw_para3 = (
        f"Trend-wise, performance appears relatively stable throughout the period, "
        f"with normal fluctuations but no major declines, "
        f"indicating consistent search demand and sustained visibility."
    )

    return tuple(_gemini_paras_batch([subtitle_raw, raw_para0, raw_para1, raw_para2, raw_para3]))


def _build_slide6(slide, search_metrics: dict, screenshots: dict) -> None:
    """Replace search performance narrative on Slide 6 (8-slide variants only)."""
    if not search_metrics or all(v == "N/A" for v in search_metrics.values()):
        return

    subtitle, para0, para1, para2, para3 = _search_perf_paras(search_metrics)

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        # object 4 — subtitle
        if shape.name == "object 4":
            content_paras = [p for p in shape.text_frame.paragraphs if p.text.strip()]
            if content_paras:
                _fill_text_run(content_paras[0], subtitle)

        # object 5 — 4 narrative paragraphs (paras 0, 2, 3, 5 — odds are spacers)
        elif shape.name == "object 5":
            content_paras = [p for p in shape.text_frame.paragraphs if p.text.strip()]
            for i, para in enumerate(content_paras):
                narr = [para0, para1, para2, para3][i] if i < 4 else None
                if narr:
                    _write_para_with_highlights(para, narr)

    if "search_screenshot" in screenshots:
        picture_names = {
            shape.name
            for shape in slide.shapes
            if getattr(shape, "shape_type", None) == 13
        }
        for candidate in ("Picture 10", "Picture 8", "Picture 9", "Picture 20"):
            if candidate in picture_names:
                _replace_image_in_slide(slide, screenshots["search_screenshot"], shape_name=candidate)
                break


def _prev_month_date_range(start_date: str) -> tuple[str, str]:
    """Return (prev_start, prev_end) strings given a start_date like 'Mar 1, 2026'."""
    from datetime import datetime, timedelta
    dt = datetime.strptime(start_date, "%b %d, %Y")
    first_of_current = dt.replace(day=1)
    last_of_prev = first_of_current - timedelta(days=1)
    first_of_prev = last_of_prev.replace(day=1)
    return (
        f"{first_of_prev.strftime('%b')} {first_of_prev.day}, {first_of_prev.year}",
        f"{last_of_prev.strftime('%b')} {last_of_prev.day}, {last_of_prev.year}",
    )


def _scrape_prev_metrics_with_context(context, report_name: str, start_date: str):
    """Scrape previous-month GA4 metrics for slide 7. Returns (prev_ga4_metrics, page)."""
    prev_ga4_metrics: dict = {}
    page = context.new_page()
    try:
        prev_start, prev_end = _prev_month_date_range(start_date)
        (
            prev_home_metrics,
            prev_snapshot_metrics,
            prev_page_views,
            prev_pages_data,
            prev_site_total_views,
            prev_countries_data,
            page,
        ) = _capture_ga4_metrics_no_screenshots(context, report_name, prev_start, prev_end, existing_page=page)
        prev_search_metrics: dict = {}
        if report_name not in SEVEN_SLIDE_REPORTS and report_name in GSC_URLS:
            try:
                prev_search_metrics, _ = _capture_gsc(
                    context,
                    report_name,
                    prev_start,
                    prev_end,
                    SCREENSHOTS_DIR / "prev_period_gsc",
                )
            except Exception as e:
                logger.warning("[2026] Previous-period GSC scrape failed: %s", e)
        prev_ga4_metrics = {
            "home_metrics": prev_home_metrics,
            "snapshot_metrics": prev_snapshot_metrics,
            "page_views": prev_page_views,
            "pages_data": prev_pages_data,
            "site_total_views": prev_site_total_views,
            "countries_data": prev_countries_data,
            "search_metrics": prev_search_metrics,
            "start_date": prev_start,
            "end_date": prev_end,
        }
    except Exception as e:
        logger.warning("[2026] Previous month metrics scrape failed: %s", e)
    return prev_ga4_metrics, page


def _scrape_previous_month_metrics(
    report_name: str,
    start_date: str,
    end_date: str,
) -> dict:
    """Scrape GA4 home metrics for the previous month. Opens its own browser session."""
    prev_ga4_metrics: dict = {}
    with sync_playwright() as p:
        context = _launch_persistent_context(p, headless=False)
        try:
            prev_ga4_metrics, _page = _scrape_prev_metrics_with_context(context, report_name, start_date)
        finally:
            try:
                context.close()
            except Exception:
                pass
    return prev_ga4_metrics


def _open_snapshot_and_set_dates(page, report_name: str, start_date: str, end_date: str):
    """Open Reports snapshot from Home and apply the requested date range."""
    from datetime import datetime as _dt

    _ensure_expected_ga4_property(page, report_name)
    page.locator("span.view-link-text", has_text="View reports snapshot").click()
    page.wait_for_timeout(3000)
    _ensure_expected_ga4_property(page, report_name)
    snapshot_date_btn = page.get_by_role("combobox", name="Open date range picker").first
    snapshot_date_btn.wait_for(state="visible", timeout=15000)
    snapshot_date_btn.click()
    page.wait_for_timeout(1000)
    page.get_by_role("menuitem").filter(has_text="Custom").click()
    page.wait_for_timeout(1000)
    start_input = page.get_by_label("Start date")
    start_input.wait_for(state="visible", timeout=10000)
    start_input.click(click_count=3)
    page.keyboard.press("Meta+a" if os.name == "posix" else "Control+a")
    page.wait_for_timeout(200)
    start_input.fill(start_date)
    page.keyboard.press("Tab")
    page.wait_for_timeout(500)
    end_input = page.get_by_label("End date")
    end_input.click(click_count=3)
    page.keyboard.press("Meta+a" if os.name == "posix" else "Control+a")
    page.wait_for_timeout(200)
    end_input.fill(end_date)
    page.keyboard.press("Tab")
    page.wait_for_timeout(500)
    page.get_by_role("button", name="Apply").click()
    expected_start = _dt.strptime(start_date, "%b %d, %Y").strftime("%Y%m%d")
    expected_end = _dt.strptime(end_date, "%b %d, %Y").strftime("%Y%m%d")
    page.wait_for_function(
        """
        ({ expectedStart, expectedEnd }) => {
            const href = window.location.href;
            return href.includes(`date00%3D${expectedStart}`) &&
                   href.includes(`date01%3D${expectedEnd}`);
        }
        """,
        arg={"expectedStart": expected_start, "expectedEnd": expected_end},
        timeout=20000,
    )
    page.wait_for_timeout(3000)
    logger.info(
        "[2026] Snapshot date confirmed. start_date=%s end_date=%s current_url=%s",
        start_date,
        end_date,
        page.url,
    )


def _scrape_countries_table(page) -> list[dict]:
    countries_data: list[dict] = []
    import re as _re

    body = page.locator("body").inner_text()
    for line in body.splitlines():
        m = _re.match(
            r"^\t\d+\t(.+?)\t([\d,]+)\s*\([^)]+\)\t([\d,]+)\s*\([^)]+\)\t([\d,]+)\s*\([^)]+\)\t([\d.]+%)\t([\d.]+)",
            line,
        )
        if m:
            countries_data.append({
                "country": m.group(1).strip(),
                "users": int(m.group(2).replace(",", "")),
                "new_users": int(m.group(3).replace(",", "")),
                "engaged_sessions": int(m.group(4).replace(",", "")),
                "engagement_rate": m.group(5),
                "engaged_sessions_per_user": m.group(6),
            })
    return countries_data


def _scrape_pages_table(page) -> tuple[dict, list[dict], int]:
    page_views: dict = {}
    pages_data: list[dict] = []
    site_total_views = 0

    def _clean_cell(text: str) -> str:
        return re.sub(r"\s+", " ", text).strip()

    def _extract_count_and_pct(text: str) -> tuple[int | None, str]:
        normalized = _clean_cell(text)
        count_match = re.match(r"^([\d,]+)", normalized)
        pct_match = re.search(r"\(([^)]+%)\)", normalized)
        count = None
        if count_match:
            try:
                count = int(count_match.group(1).replace(",", ""))
            except ValueError:
                count = None
        return count, (pct_match.group(1) if pct_match else "")

    try:
        rows = page.locator("table.adv-table tbody tr")
        row_count = rows.count()
        for i in range(row_count):
            row = rows.nth(i)
            cols = [_clean_cell(text) for text in row.locator("td").all_inner_texts()]
            cols = [col for col in cols if col]
            if len(cols) < 6:
                continue

            first = cols[0].lower()
            if first == "total":
                total_views, _ = _extract_count_and_pct(cols[1])
                if total_views is not None:
                    site_total_views = total_views
                continue

            if not cols[0].isdigit() or len(pages_data) >= 10:
                continue

            path_value = cols[1]
            views, views_pct = _extract_count_and_pct(cols[2])
            active_users, active_users_pct = _extract_count_and_pct(cols[3])
            if views is None or active_users is None:
                continue

            row_data = {
                "title": path_value,
                "path": path_value,
                "views": views,
                "views_pct": views_pct,
                "active_users": active_users,
                "active_users_pct": active_users_pct,
                "views_per_user": cols[4],
                "avg_engagement_time": cols[5],
            }
            pages_data.append(row_data)
            if len(page_views) < 4:
                page_views[path_value] = views

        # GA4 sometimes renders the "Total" summary outside the normal tbody rows.
        # Fall back to the full page text when the direct row scan didn't capture it.
        if site_total_views == 0:
            body = page.locator("body").inner_text()
            total_match = re.search(r"(?:^|\n)\s*Total\s*\n\s*([\d,]+)\s*\n\s*100% of total", body)
            if total_match:
                try:
                    site_total_views = int(total_match.group(1).replace(",", ""))
                except ValueError:
                    site_total_views = 0
    except Exception:
        # Fall back to the old body-text path if GA4 changes the row structure again.
        body = page.locator("body").inner_text()
        total_match = re.search(r"(?:^|\n)\s*Total\s*\n\s*([\d,]+)\s*\n\s*100% of total", body)
        if total_match:
            try:
                site_total_views = int(total_match.group(1).replace(",", ""))
            except ValueError:
                site_total_views = 0

        for line in body.splitlines():
            cols = [col.strip() for col in line.split("\t") if col.strip()]
            if not cols or len(cols) < 5 or not cols[0].isdigit() or len(pages_data) >= 10:
                continue

            path_value = cols[1]
            views, views_pct = _extract_count_and_pct(cols[2])
            active_users, active_users_pct = _extract_count_and_pct(cols[3])
            if views is None or active_users is None:
                continue

            pages_data.append({
                "title": path_value,
                "path": path_value,
                "views": views,
                "views_pct": views_pct,
                "active_users": active_users,
                "active_users_pct": active_users_pct,
                "views_per_user": cols[4],
                "avg_engagement_time": cols[5] if len(cols) > 5 else "",
            })
            if len(page_views) < 4:
                page_views[path_value] = views

    return page_views, pages_data, site_total_views


def _capture_ga4_metrics_no_screenshots(context, report_name: str, start_date: str, end_date: str, existing_page=None) -> tuple[dict, dict, dict, list[dict], int, list[dict], object]:
    """Repeat the GA4 metric collection flow without taking screenshots."""
    home_metrics: dict = {}
    snapshot_metrics: dict = {}
    page_views: dict = {}
    pages_data: list[dict] = []
    site_total_views: int = 0
    countries_data: list[dict] = []

    page = existing_page or context.new_page()
    page.bring_to_front()

    if existing_page is None:
        page = _switch_ga4_property_via_search(page, report_name)

    page = _goto_ga4_section(page, report_name, "/home")
    try:
        page.keyboard.press("Escape")
        page.wait_for_timeout(1000)
    except Exception:
        pass

    _open_snapshot_and_set_dates(page, report_name, start_date, end_date)
    snapshot_metrics = _scrape_snapshot_metrics(page)

    try:
        page.get_by_text("View countries").click()
        page.wait_for_timeout(4000)
        _ensure_expected_ga4_property(page, report_name)
        page.locator("th.cdk-column-__row_index__").first.wait_for(state="visible", timeout=10000)
        countries_data = _scrape_countries_table(page)
        page.go_back()
        page.wait_for_timeout(3000)
        _ensure_expected_ga4_property(page, report_name)
    except Exception as e:
        logger.warning("[2026] Previous-period countries scrape failed: %s", e)

    try:
        page.get_by_role("button", name="View pages and screens", exact=True).click()
        page.wait_for_timeout(4000)
        _ensure_expected_ga4_property(page, report_name)
        page.locator("th.cdk-column-__row_index__").first.wait_for(state="visible", timeout=10000)
        page_views, pages_data, site_total_views = _scrape_pages_table(page)
    except Exception as e:
        logger.warning("[2026] Previous-period pages scrape failed: %s", e)

    try:
        page = _goto_ga4_section(page, report_name, "/home")
        page.wait_for_timeout(2000)
        home_metrics = _scrape_home_metrics(page)
    except Exception as e:
        logger.warning("[2026] Previous-period home metrics scrape failed: %s", e)

    return home_metrics, snapshot_metrics, page_views, pages_data, site_total_views, countries_data, page


def _scrape_ga4_page_paths(context, report_name: str, start_date: str, end_date: str, existing_page=None) -> list[dict]:
    """Navigate GA4 Pages and screens report, switch dimension to page path, return top 5 with title + path.

    Pass existing_page to reuse an already-navigated GA4 page (avoids a second property search).
    """
    import re as _re

    top5: list[dict] = []

    try:
        if existing_page is not None:
            # Already on GA4 with correct property — go to home, click snapshot, set dates, click pages and screens
            page = existing_page
            page.bring_to_front()
            page = _goto_ga4_section(page, report_name, "/home")
            page.wait_for_timeout(2000)
            _open_snapshot_and_set_dates(page, report_name, start_date, end_date)
            page.get_by_role("button", name="View pages and screens", exact=True).click()
            page.wait_for_timeout(4000)
        else:
            page = context.new_page()
            page.bring_to_front()
            page = _switch_ga4_property_via_search(page, report_name)
            page = _goto_ga4_section(page, report_name, "/home")
            page.wait_for_timeout(2000)
            _open_snapshot_and_set_dates(page, report_name, start_date, end_date)
            page.get_by_role("button", name="View pages and screens", exact=True).click()

        page.wait_for_timeout(4000)
        _ensure_expected_ga4_property(page, report_name)
        page.locator("th.cdk-column-__row_index__").first.wait_for(state="visible", timeout=10000)

        # --- Scrape page titles first (default view) ---
        body_titles = page.locator("body").inner_text()
        titles: list[str] = []
        for line in body_titles.splitlines():
            m = _re.match(
                r"^\t\d+\t(.+?)\t([\d,]+)\s*\([^)]+\)\t",
                line,
            )
            if m and len(titles) < 5:
                titles.append(m.group(1).strip())

        # --- Dump table header HTML to file for selector discovery ---
        try:
            import pathlib as _pl
            _header_html = page.locator("thead").first.inner_html()
            (SCREENSHOTS_DIR / "ga4_thead_debug.html").write_text(_header_html, encoding="utf-8")
            logger.info("[2026] Dumped GA4 thead HTML → %s", SCREENSHOTS_DIR / "ga4_thead_debug.html")
        except Exception:
            pass

        # --- Switch dimension to "Page path and screen class" to get real URL paths ---
        # GA4 uses data-guidedhelpid="table-dimension-picker" on the button in the primary column header
        switched = False
        try:
            # The primary dimension column has two "table-dimension-picker" buttons —
            # the first one is the dimension selector (second is the secondary dimension "+")
            dim_btn = page.locator(
                "th.cdk-column-ROW_HEADER-unifiedScreenClass-primaryDimensionColumn "
                "button[data-guidedhelpid='table-dimension-picker']"
            ).first
            dim_btn.wait_for(state="visible", timeout=8000)
            dim_btn.click()
            page.wait_for_timeout(1500)
            # Menu opens — try multiple selectors since GA4 mat-menu items may not surface as role=menuitem
            _clicked_dim = False
            for _sel in [
                lambda: page.get_by_role("menuitem").filter(has_text="Page path").first,
                lambda: page.locator("button.mat-mdc-menu-item").filter(has_text="Page path").first,
                lambda: page.locator("[role='menuitem']").filter(has_text="Page path").first,
                lambda: page.locator("mat-option, .mat-mdc-option").filter(has_text="Page path").first,
                lambda: page.get_by_text("Page path and screen class").first,
            ]:
                try:
                    el = _sel()
                    el.wait_for(state="visible", timeout=5000)
                    el.click()
                    _clicked_dim = True
                    break
                except Exception:
                    continue
            if not _clicked_dim:
                raise RuntimeError("Could not find 'Page path' menu item after clicking dimension picker")
            page.wait_for_timeout(3000)
            page.locator("th.cdk-column-__row_index__").first.wait_for(state="visible", timeout=10000)
            switched = True
        except Exception as e:
            logger.warning("[2026] Could not switch GA4 dimension to page path: %s", e)

        body_paths = page.locator("body").inner_text()
        paths: list[tuple[str, int]] = []
        for line in body_paths.splitlines():
            m = _re.match(
                r"^\t\d+\t(.+?)\t([\d,]+)\s*\([^)]+\)\t",
                line,
            )
            if m and len(paths) < 5:
                paths.append((m.group(1).strip(), int(m.group(2).replace(",", ""))))

        # Zip titles + paths by row index
        for i, (path, views) in enumerate(paths):
            title = titles[i] if i < len(titles) else path
            top5.append({"title": title, "path": path, "views": views})

    except Exception as e:
        logger.warning("[2026] GA4 page paths scrape failed: %s", e)

    return top5


def _google_search_url(title: str, base_url: str) -> str | None:
    """Use DuckDuckGo search (site: query) to find the real URL for a given page title.

    Verifies the result by checking that the DDG result title contains the meaningful keyword.
    Falls back through up to 3 results before giving up.
    """
    try:
        import time as _time
        from ddgs import DDGS
        import urllib.parse as _up

        domain = _up.urlparse(base_url).netloc
        # Strip brand name from either end:
        # "Devices - Econet Wireless Zimbabwe" → "Devices"
        # "Econet Wireless Zimbabwe - Contact Us" → "Contact Us"
        parts = [p.strip() for p in title.split(" - ")]
        domain_root = domain.split(".")[1] if domain.count(".") >= 1 else domain
        meaningful = [p for p in parts if domain_root.lower() not in p.lower()]
        short_title = meaningful[0] if meaningful else parts[0]

        query = f"site:{domain} {short_title}"
        _time.sleep(1)
        with DDGS() as ddgs:
            results = list(ddgs.text(query, max_results=5))

        for r in results:
            result_title = r.get("title", "")
            result_url = r.get("href", "")
            if not result_url:
                continue
            # Must be on the same domain
            if domain not in result_url:
                continue
            # Accept if the result title contains our keyword (case-insensitive)
            if short_title.lower() in result_title.lower():
                return result_url

        # Second pass — domain match only, ignore title
        for r in results:
            result_url = r.get("href", "")
            if result_url and domain in result_url and result_url.rstrip("/") != base_url.rstrip("/"):
                return result_url

    except Exception as e:
        logger.warning("[2026] DDG search for title %r failed: %s", title, e)
    return None


def _audit_page_ctas(page, base_url: str) -> list[dict]:
    """Click primary CTA buttons on the current page and verify they lead somewhere valid.

    Returns a list of {label, href, resolved_url, status, broken} dicts.
    Only tests buttons/links that look like CTAs — ignores nav, footer, social icons.
    """
    import urllib.parse as _up

    results = []
    origin = _up.urlparse(base_url).scheme + "://" + _up.urlparse(base_url).netloc

    # Collect candidate CTA anchors: prominent buttons, "a" tags with button-like text
    cta_keywords = ["buy", "shop", "get", "start", "activate", "download", "contact",
                    "learn", "explore", "order", "sign", "register", "apply", "more"]
    try:
        # Grab all visible <a> and <button> elements
        elements = page.locator("a[href], button").all()
        seen_hrefs: set = set()

        for el in elements[:40]:  # cap at 40 to avoid nav spam
            try:
                if not el.is_visible():
                    continue
                label = (el.inner_text() or "").strip().lower()
                if not label or len(label) > 60:
                    continue
                if not any(k in label for k in cta_keywords):
                    continue

                href = el.get_attribute("href") or ""
                # Skip anchors, mailto, tel, javascript
                if not href or href.startswith(("#", "mailto:", "tel:", "javascript:")):
                    continue
                # Resolve relative URLs
                resolved = href if href.startswith("http") else _up.urljoin(origin + "/", href)
                if resolved in seen_hrefs:
                    continue
                seen_hrefs.add(resolved)

                # Open in a new tab to test without leaving the page
                new_tab = page.context.new_page()
                status = 0
                final_url = resolved
                broken = True
                try:
                    resp = new_tab.goto(resolved, wait_until="domcontentloaded", timeout=10000)
                    status = resp.status if resp else 0
                    final_url = new_tab.url
                    broken = status >= 400 or status == 0
                except Exception:
                    broken = True
                finally:
                    new_tab.close()

                results.append({
                    "label": label,
                    "href": href,
                    "resolved_url": final_url,
                    "status": status,
                    "broken": broken,
                })
            except Exception:
                continue
    except Exception as e:
        logger.warning("[2026] CTA audit failed: %s", e)

    return results


def _scrape_website_pages(
    report_name: str,
    start_date: str,
    end_date: str,
    _stage_callback=None,
) -> tuple[dict, dict]:
    """Navigate GA4 Pages report to get top 5 page titles, resolve real URLs via DDG search, then scrape content.

    Returns (website_pages, prev_ga4_metrics) — both scraped in a single browser session.
    """
    def _stage(msg: str):
        if _stage_callback:
            _stage_callback(msg)
        logger.info("[2026] Stage: %s", msg)
    base_url = GSC_URLS.get(report_name, "")
    if not base_url:
        return {"top": []}, {}

    result: dict = {"top": []}
    prev_ga4_metrics: dict = {}

    with sync_playwright() as p:
        context = _launch_persistent_context(p, headless=False)
        ga4_page = None
        try:
            _stage("Scraping top pages from GA4...")
            top5 = _scrape_ga4_page_paths(
                context,
                report_name,
                start_date,
                end_date,
                existing_page=ga4_page,
            )

            import urllib.parse as _up
            domain_root = _up.urlparse(base_url).netloc.split(".")[1]

            screenshots_dir = out_dir = SCREENSHOTS_DIR / "site_pages"
            screenshots_dir.mkdir(parents=True, exist_ok=True)

            site_page = context.new_page()
            site_page.set_viewport_size({"width": 1440, "height": 900})
            qr_codes: list[dict] = []

            for row in top5:
                title = row.get("title", "")
                path = row.get("path", "")

                # Homepage — all title parts are brand name, use base URL directly
                all_parts_are_brand = all(
                    domain_root.lower() in p.lower()
                    for p in [p.strip() for p in title.split(" - ")]
                )
                if all_parts_are_brand:
                    url = base_url
                elif path.startswith("/"):
                    # Dimension switch succeeded — path is a real URL path, use it directly
                    url = base_url.rstrip("/") + path
                else:
                    # Dimension switch failed — path is a title, use DDG to resolve
                    url = _google_search_url(title, base_url)
                    if not url:
                        url = base_url.rstrip("/") + f"/{path}"

                # Skip if this URL was already visited (DDG returned duplicate)
                already_visited = any(r["url"] == url for r in result["top"])
                if already_visited:
                    logger.warning("[2026] Duplicate URL %s for title %r — skipping", url, title)
                    result["top"].append({"title": title, "url": url, "content": "", "views": row.get("views", 0), "status": 0, "screenshot": None})
                    continue

                status = 0
                content = ""
                screenshot_path = None
                short_title = title.split(" - ")[0].strip()
                _stage(f"Visiting {short_title}...")
                try:
                    response = site_page.goto(url, wait_until="domcontentloaded", timeout=20000)
                    status = response.status if response else 0

                    # Block custom fonts so screenshot never hangs waiting for them
                    try:
                        site_page.add_style_tag(content="@font-face { src: none !important; } * { font-family: Arial, sans-serif !important; }")
                    except Exception:
                        pass

                    # Dismiss cookie consent banners before anything else
                    for cookie_sel in [
                        "button:has-text('Accept')", "button:has-text('Accept All')",
                        "button:has-text('I Accept')", "button:has-text('OK')",
                        "button:has-text('Got it')", "button:has-text('Allow')",
                        "[id*='cookie'] button", "[class*='cookie'] button",
                        "[id*='consent'] button", "[class*='consent'] button",
                    ]:
                        try:
                            btn = site_page.locator(cookie_sel).first
                            if btn.is_visible(timeout=1000):
                                btn.click(timeout=1000)
                                site_page.wait_for_timeout(500)
                                break
                        except Exception:
                            continue

                    # Wait for images to finish loading before scrolling/screenshotting
                    site_page.wait_for_timeout(2000)
                    try:
                        site_page.wait_for_function(
                            "() => [...document.images].every(img => img.complete)",
                            timeout=8000,
                        )
                    except Exception:
                        pass  # some images may still be lazy — proceed anyway

                    # Scroll to trigger lazy-loaded content
                    site_page.evaluate("window.scrollTo(0, document.body.scrollHeight / 2)")
                    site_page.wait_for_timeout(1000)
                    site_page.evaluate("window.scrollTo(0, document.body.scrollHeight)")
                    site_page.wait_for_timeout(1000)
                    site_page.evaluate("window.scrollTo(0, 0)")
                    site_page.wait_for_timeout(1000)

                    # Wait again for any images triggered by scrolling
                    try:
                        site_page.wait_for_function(
                            "() => [...document.images].every(img => img.complete)",
                            timeout=5000,
                        )
                    except Exception:
                        pass

                    content = site_page.inner_text("body")
                    content = " ".join(content.split())[:3000]
                    # Full-page screenshot — inject CSS to block remote font loading which can hang
                    safe_name = "".join(c if c.isalnum() else "_" for c in title)[:60]
                    shot_path = screenshots_dir / f"{safe_name}.png"
                    try:
                        site_page.add_style_tag(content="@font-face { font-family: any; src: none !important; }")
                    except Exception:
                        pass
                    try:
                        site_page.screenshot(path=str(shot_path), full_page=True, timeout=30000, animations="disabled")
                        screenshot_path = str(shot_path)
                    except Exception as _se:
                        logger.warning("[2026] Screenshot timed out for %s, falling back to viewport: %s", url, _se)
                        try:
                            site_page.screenshot(path=str(shot_path), full_page=False, timeout=15000, animations="disabled")
                            screenshot_path = str(shot_path)
                        except Exception:
                            pass

                    # --- QR code detection (only runs if screenshot exists) ---
                    try:
                        import cv2 as _cv2
                        _img = _cv2.imread(str(shot_path))
                        if _img is not None:
                            _qr_data, _, _ = _cv2.QRCodeDetector().detectAndDecode(_img)
                            if _qr_data:
                                qr_codes.append({"url": url, "data": _qr_data})
                                logger.info("[2026] QR code found on %s: %s", url, _qr_data)
                    except Exception:
                        pass

                    # --- CTA button/link audit ---
                    cta_results = _audit_page_ctas(site_page, base_url)
                    if cta_results:
                        result.setdefault("cta_audit", []).append({
                            "page_url": url,
                            "page_title": title,
                            "ctas": cta_results,
                        })
                except Exception as e:
                    logger.warning("[2026] Failed to load %s: %s", url, e)

                result["top"].append({
                    "title": title,
                    "url": url,
                    "content": content,
                    "views": row.get("views", 0),
                    "status": status,
                    "screenshot": screenshot_path,
                })
        except Exception as e:
            logger.warning("[2026] Website page scrape failed: %s", e)
        finally:
            try:
                context.close()
            except Exception:
                pass

    result["qr_codes"] = qr_codes
    return result, prev_ga4_metrics


def _generate_recommendations_2026(
    report_name: str,
    home_metrics: dict,
    snapshot_metrics: dict,
    search_metrics: dict,
    pages_data: list[dict],
    countries_data: list[dict],
    date_range: str,
    prev_ga4_metrics: dict | None = None,
    website_pages: dict | None = None,
) -> list[dict]:
    """Use Gemini to generate 3 structured recommendations, each with a title + 3 bullet points."""
    load_runtime_environment()
    brand = report_name.replace("_", " ").title()

    top_pages = ", ".join(
        f"{p.get('_label', p['title'].split(' - ')[0])} ({p['views']:,} views, {p['views_per_user']} views/user)"
        for p in pages_data[:5]
    ) if pages_data else "N/A"

    top_countries = ", ".join(
        f"{r['country']} ({r['users']:,} users, {r['engagement_rate']} engagement)"
        for r in sorted(countries_data, key=lambda x: x["users"], reverse=True)[:3]
    ) if countries_data else "N/A"

    channels = snapshot_metrics.get("channels", {})
    channels_str = ", ".join(f"{k}: {v}" for k, v in list(channels.items())[:5]) if channels else "N/A"

    prev_home_metrics = (prev_ga4_metrics or {}).get("home_metrics", {})
    prev_snapshot_metrics = (prev_ga4_metrics or {}).get("snapshot_metrics", {})
    prev_pages_data = (prev_ga4_metrics or {}).get("pages_data", [])
    prev_countries_data = (prev_ga4_metrics or {}).get("countries_data", [])
    prev_search_metrics = (prev_ga4_metrics or {}).get("search_metrics", {})
    prev_channels = prev_snapshot_metrics.get("channels", {})
    prev_channels_str = ", ".join(f"{k}: {v}" for k, v in list(prev_channels.items())[:5]) if prev_channels else "N/A"

    mom_section = ""
    if prev_ga4_metrics:
        def _mom(label, key):
            curr = home_metrics.get(key, "N/A")
            prev = prev_home_metrics.get(key, "N/A")
            return f"{label}: {prev} (prev) → {curr} (current)"

        prev_top_pages = ", ".join(
            f"{p['title'].split(' - ')[0]} ({p['views']:,} views)"
            for p in prev_pages_data[:5]
        ) if prev_pages_data else "N/A"

        prev_top_countries = ", ".join(
            f"{r['country']} ({r['users']:,} users, {r['engagement_rate']} engagement)"
            for r in sorted(prev_countries_data, key=lambda x: x["users"], reverse=True)[:3]
        ) if prev_countries_data else "N/A"

        mom_section = (
            "\nMonth-on-month comparison:\n"
            + _mom("Active users", "Active users") + "\n"
            + _mom("New users", "New users") + "\n"
            + _mom("Avg engagement time", "Average engagement time per active user") + "\n"
            + f"Acquisition channels: {prev_channels_str} (prev) → {channels_str} (current)\n"
            + f"Top pages: {prev_top_pages} (prev) → {top_pages} (current)\n"
            + f"Top countries: {prev_top_countries} (prev) → {top_countries} (current)\n"
            + (
                f"Search metrics: impressions {prev_search_metrics.get('impressions', 'N/A')}, "
                f"clicks {prev_search_metrics.get('clicks', 'N/A')}, "
                f"CTR {prev_search_metrics.get('ctr', 'N/A')}, "
                f"avg position {prev_search_metrics.get('avg_position', 'N/A')} (prev) → "
                f"impressions {search_metrics.get('impressions', 'N/A')}, "
                f"clicks {search_metrics.get('clicks', 'N/A')}, "
                f"CTR {search_metrics.get('ctr', 'N/A')}, "
                f"avg position {search_metrics.get('avg_position', 'N/A')} (current)\n"
                if prev_search_metrics else ""
            )
        )

    # CTA audit findings
    cta_section = ""
    if website_pages and website_pages.get("cta_audit"):
        lines = []
        for page_audit in website_pages["cta_audit"]:
            broken = [c for c in page_audit["ctas"] if c["broken"]]
            working = [c for c in page_audit["ctas"] if not c["broken"]]
            if broken:
                lines.append(f"  {page_audit['page_title']} ({page_audit['page_url']}):")
                for c in broken:
                    lines.append(f"    BROKEN CTA: '{c['label']}' → {c['resolved_url']} (status {c['status']})")
            if working:
                lines.append(f"  {page_audit['page_title']} — working CTAs: {', '.join(c['label'] for c in working[:5])}")
        if lines:
            cta_section = "\nCTA button audit (clicked and verified on live pages):\n" + "\n".join(lines) + "\n"

    # QR code findings
    qr_section = ""
    if website_pages and website_pages.get("qr_codes"):
        qr_lines = "\n".join(
            f"- {q['url']}: QR encodes {q['data']!r}"
            for q in website_pages["qr_codes"]
        )
        qr_section = f"\nFunctional QR codes detected on the live site:\n{qr_lines}\n"

    # Website page content section — only include pages with real, unique content
    pages_section = ""
    if website_pages:
        seen_urls: set = set()
        valid_pages = []
        for p in website_pages.get("top", []):
            url = p.get("url", "")
            content = p.get("content", "")
            if content and url and url not in seen_urls:
                seen_urls.add(url)
                valid_pages.append(p)
        if valid_pages:
            top_content = "\n".join(
                f"- {p['title']} ({p['url']}): {p['content'][:600]}"
                for p in valid_pages
            )
            pages_section = "\nTop performing pages (actual site content scraped from live site):\n" + top_content + "\n"

    prompt = (
        f"You are a digital analytics expert writing a monthly website performance report for {brand}.\n\n"
        f"Period: {date_range}\n"
        f"Active users: {home_metrics.get('Active users', 'N/A')}\n"
        f"New users: {home_metrics.get('New users', 'N/A')}\n"
        f"Avg engagement time: {home_metrics.get('Average engagement time per active user', 'N/A')}\n"
        f"Top acquisition channels: {channels_str}\n"
        f"Top pages by views: {top_pages}\n"
        f"Top countries: {top_countries}\n"
        f"{mom_section}"
        f"{cta_section}"
        f"{qr_section}"
        f"{pages_section}\n"
        "Screenshots of the top performing pages are attached. Use them to identify UX, content, "
        "and conversion issues visible on the actual pages. "
        "Ignore any cookie consent banners visible in screenshots — these are browser artifacts. "
        "Only flag missing images or broken elements if the CTA audit above explicitly confirms them as broken.\n\n"
        "Write exactly 3 actionable recommendations. Each must follow this EXACT format with no deviations:\n\n"
        "TITLE: <action-oriented title, 4-7 words>\n"
        "BODY: <one sentence framing the recommendation>\n"
        "- <specific action bullet 1>\n"
        "- <specific action bullet 2>\n"
        "- <specific action bullet 3>\n"
        "- <specific action bullet 4>\n"
        "- <specific action bullet 5>\n"
        "---\n"
        "TITLE: <action-oriented title, 4-7 words>\n"
        "BODY: <one sentence framing the recommendation>\n"
        "- <specific action bullet 1>\n"
        "- <specific action bullet 2>\n"
        "- <specific action bullet 3>\n"
        "- <specific action bullet 4>\n"
        "---\n"
        "TITLE: <action-oriented title, 4-7 words>\n"
        "BODY: <one sentence framing the recommendation>\n"
        "- <specific action bullet 1>\n"
        "- <specific action bullet 2>\n"
        "- <specific action bullet 3>\n"
        "---\n"
        "Use formal business English. No markdown bold, no em dashes. "
        "Be specific — reference actual page names, metrics, and visual observations from the screenshots."
    )

    # Build multimodal contents — text prompt + screenshots of top pages
    from google.genai import types as _gtypes
    contents: list = [prompt]
    if website_pages:
        for pg in website_pages.get("top", []):
            shot = pg.get("screenshot")
            if shot and Path(shot).exists():
                try:
                    image_bytes = Path(shot).read_bytes()
                    contents.append(_gtypes.Part.from_bytes(data=image_bytes, mime_type="image/png"))
                    contents.append(f"[Screenshot of: {pg['title']} — {pg['url']}]")
                except Exception:
                    pass

    client = genai.Client(api_key=os.environ["GEMINI_API_KEY"])
    resp = client.models.generate_content(model="gemini-2.5-flash", contents=contents)
    text = resp.text.strip()

    # Expected bullets per rec matches template: rec1=5, rec2=4, rec3=3
    bullet_counts = [5, 4, 3]

    recs = []
    for i, block in enumerate(text.split("---")):
        block = block.strip()
        if not block:
            continue
        lines = [l.strip() for l in block.splitlines() if l.strip()]
        title = ""
        body = ""
        bullets = []
        for line in lines:
            if line.upper().startswith("TITLE:"):
                title = line[6:].strip()
            elif line.upper().startswith("BODY:"):
                body = line[5:].strip()
            elif line.startswith("-"):
                bullets.append(line[1:].strip())
        if title:
            max_bullets = bullet_counts[i] if i < len(bullet_counts) else 3
            # Pad with body sentence if not enough bullets
            all_bullets = ([body] + bullets) if body else bullets
            recs.append({"title": title, "body": body, "bullets": all_bullets[:max_bullets]})
        if len(recs) == 3:
            break

    return recs


def _build_recommendations_slide(slide) -> None:
    """Write placeholder text on the recommendations slide for manual completion."""
    placeholders = [
        ("Recommendation 1", [
            "Add your first recommendation here.",
            "Supporting point for recommendation 1.",
            "Supporting point for recommendation 1.",
        ]),
        ("2. Recommendation 2", [
            "Add your second recommendation here.",
            "Supporting point for recommendation 2.",
            "Supporting point for recommendation 2.",
        ]),
        ("3. Recommendation 3", [
            "Add your third recommendation here.",
            "Supporting point for recommendation 3.",
            "Supporting point for recommendation 3.",
        ]),
    ]

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        if shape.name == "object 3":
            content_paras = [p for p in shape.text_frame.paragraphs if p.text.strip()]
            if content_paras:
                _fill_text_run(content_paras[0], "Three Key Initiatives to Enhance Performance")

        elif shape.name == "object 7":
            paras = shape.text_frame.paragraphs
            layout = [
                (0,  0, "title"),
                (1,  0, "bullet_0"), (2,  0, "bullet_1"), (3,  0, "bullet_2"),
                (4,  0, "bullet_3"), (5,  0, "bullet_4"),
                (6,  1, "title"),
                (7,  1, "bullet_0"), (8,  1, "bullet_1"), (9,  1, "bullet_2"),
                (10, 1, "bullet_3"),
                (11, 2, "title"),
                (12, 2, "bullet_0"), (13, 2, "bullet_1"), (14, 2, "bullet_2"),
            ]
            for para_idx, rec_idx, role in layout:
                if para_idx >= len(paras) or rec_idx >= len(placeholders):
                    continue
                title, bullets = placeholders[rec_idx]
                para = paras[para_idx]
                if role == "title":
                    _fill_text_run(para, title)
                else:
                    bullet_n = int(role.split("_")[1])
                    if bullet_n < len(bullets):
                        _fill_text_run(para, bullets[bullet_n])
                    else:
                        _fill_text_run(para, "")


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------

def generate_report_2026(
    report_name: str,
    date_range: str,
    report_date: str,
    start_date: str,
    end_date: str,
    _stage_callback=None,
) -> Path:
    """Full 2026 pipeline: capture GA4 → fill PPTX → save."""
    def _stage(msg: str):
        if _stage_callback:
            _stage_callback(msg)
        logger.info("[2026] Stage: %s", msg)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    # Step 1: Capture screenshots + metrics
    screenshots, home_metrics, snapshot_metrics, page_views, pages_data, site_total_views, countries_data, search_metrics = capture_2026(
        report_name, start_date, end_date, _stage_callback=_stage_callback
    )


    # Step 2: Load template
    template_path = TEMPLATES_DIR / TEMPLATES_2026[report_name]
    prs = Presentation(str(template_path))
    slide_count = len(prs.slides)
    is_7_slide = report_name in SEVEN_SLIDE_REPORTS

    perf_month = _performance_month(date_range)

    # Step 3: Build each slide
    logger.info("[2026] Building slides for %s (%d slides)", report_name, slide_count)

    _stage("Building slide 1 up to complete...")
    _build_slide1(prs.slides[0], perf_month, screenshots, report_name=report_name)
    _stage("Building slide 2 up to complete...")
    _build_slide2(prs.slides[1], home_metrics, snapshot_metrics, report_name, search_metrics)
    _stage("Building slide 3 up to complete...")
    _build_slide3(prs.slides[2], home_metrics, snapshot_metrics, report_name, screenshots)
    _stage("Building slide 4 up to complete...")
    _build_slide4(prs.slides[3], countries_data, screenshots)
    _stage("Building slide 5 up to complete...")
    _build_slide5(prs.slides[4], pages_data, screenshots, site_total_views)

    if not is_7_slide and slide_count >= 8:
        _stage("Building slide 6 up to complete...")
        _build_slide6(prs.slides[5], search_metrics, screenshots)

    rec_slide_idx = slide_count - 2
    _stage(f"Building slide {rec_slide_idx + 1} up to complete...")
    _build_recommendations_slide(prs.slides[rec_slide_idx])

    # Step 4: Save
    safe_name = report_name.replace("_", "-")
    output_path = OUTPUT_DIR / f"{safe_name}-{report_date.replace(' ', '-')}.pptx"
    prs.save(str(output_path))
    logger.info("[2026] Saved report to %s", output_path)
    return output_path
