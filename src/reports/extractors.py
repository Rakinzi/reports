"""
Per-client PPTX text extraction methods.
Each function opens its dedicated template and returns a dict of slide number -> list of text lines.
"""

from pathlib import Path

from pptx import Presentation

from .runtime import get_templates_dir

TEMPLATES_DIR = get_templates_dir()


def _extract_slides(pptx_path: Path) -> dict[int, list[str]]:
    prs = Presentation(str(pptx_path))
    slides = {}
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if line:
                    texts.append(line)
        slides[i] = texts
    return slides


def extract_econet() -> dict[int, list[str]]:
    return _extract_slides(TEMPLATES_DIR / "Econet-February 2026 Website Report.pptx")


def extract_econet_ai() -> dict[int, list[str]]:
    return _extract_slides(TEMPLATES_DIR / "Econet AI February 2026 Website Report.pptx")


def extract_ecocash() -> dict[int, list[str]]:
    return _extract_slides(TEMPLATES_DIR / "Ecocash February 2026 Website Report.pptx")


def extract_ecosure() -> dict[int, list[str]]:
    return _extract_slides(TEMPLATES_DIR / "Ecosure January 2026 Website Report (1).pptx")


def extract_zimplats() -> dict[int, list[str]]:
    return _extract_slides(TEMPLATES_DIR / "Zimplats February 2026 Website Report.pptx")


def extract_union_hardware() -> dict[int, list[str]]:
    return _extract_slides(TEMPLATES_DIR / "Union Hardware February 2026 Report.pptx")


def extract_cancer_serve() -> dict[int, list[str]]:
    return _extract_slides(TEMPLATES_DIR / "Cancer Serve February 2025 Website Report.pptx")
