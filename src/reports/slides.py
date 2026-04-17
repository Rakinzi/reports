"""
PPTX → PNG slide rendering and structured field extraction.

Rendering uses python-pptx + Pillow (no external dependencies required).
Each slide is saved as slide_0.png, slide_1.png, etc. in a report-specific directory.
A preview.pdf is synthesised by stitching the PNGs into a multi-page PDF via Pillow.
"""

import io
import re
from pathlib import Path

import lxml.etree as etree
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Pt

from .logging_utils import configure_logging
from .runtime import get_slides_dir

logger = configure_logging()

# Render at 150 DPI — gives 1500×844 canvas for a 10×5.63 inch slide (sharper thumbnails)
_DPI = 150
_EMU_PER_INCH = 914400

# Fallback font search order
_FONT_CANDIDATES = [
    "C:/Windows/Fonts/arial.ttf",
    "C:/Windows/Fonts/calibri.ttf",
    "C:/Windows/Fonts/segoeui.ttf",
    "/System/Library/Fonts/Helvetica.ttc",
    "/System/Library/Fonts/SFNS.ttf",
    "/Library/Fonts/Arial.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
]


def _emu_to_px(emu: int) -> int:
    return round(emu / _EMU_PER_INCH * _DPI)


def _resolve_scheme_colors(prs: Presentation) -> dict[str, tuple[int, int, int]]:
    """Parse the theme XML and return a mapping of scheme name → RGB tuple."""
    colors: dict[str, tuple[int, int, int]] = {}
    master = prs.slides[0].slide_layout.slide_master if prs.slides else None
    if master is None:
        return colors
    for rel in master.part.rels.values():
        if "theme" not in rel.reltype:
            continue
        try:
            xml = etree.fromstring(rel.target_part.blob)
        except Exception:
            continue
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        clr_scheme = xml.find(f".//{{{ns}}}clrScheme")
        if clr_scheme is None:
            continue
        for child in clr_scheme:
            name = child.tag.split("}")[-1]  # e.g. "dk1", "lt1", "accent1"
            rgb_el = child.find(f"{{{ns}}}srgbClr")
            sys_el = child.find(f"{{{ns}}}sysClr")
            if rgb_el is not None:
                val = rgb_el.get("val", "")
                if len(val) == 6:
                    colors[name] = (int(val[0:2], 16), int(val[2:4], 16), int(val[4:6], 16))
            elif sys_el is not None:
                last = sys_el.get("lastClr", "")
                if len(last) == 6:
                    colors[name] = (int(last[0:2], 16), int(last[2:4], 16), int(last[4:6], 16))
        break
    return colors


def _get_font(size_pt: float) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    size_px = max(8, round(size_pt * _DPI / 72))
    for path in _FONT_CANDIDATES:
        try:
            return ImageFont.truetype(path, size_px)
        except (OSError, IOError):
            continue
    return ImageFont.load_default()


def _resolve_run_color(
    run_el: etree._Element,
    scheme_colors: dict[str, tuple[int, int, int]],
    fallback: tuple[int, int, int],
) -> tuple[int, int, int]:
    """Return the RGB color for a text run element."""
    rpr = run_el.find(qn("a:rPr"))
    if rpr is None:
        return fallback
    solid_fill = rpr.find(qn("a:solidFill"))
    if solid_fill is None:
        return fallback
    rgb_el = solid_fill.find(qn("a:srgbClr"))
    if rgb_el is not None:
        val = rgb_el.get("val", "")
        if len(val) == 6:
            return (int(val[0:2], 16), int(val[2:4], 16), int(val[4:6], 16))
    scheme_el = solid_fill.find(qn("a:schemeClr"))
    if scheme_el is not None:
        name = scheme_el.get("val", "")
        return scheme_colors.get(name, fallback)
    return fallback


def _resolve_solid_fill_from_xml(
    sp_el: etree._Element,
    scheme_colors: dict[str, tuple[int, int, int]],
) -> tuple[int, int, int, int] | None:
    """Read solidFill directly from shape XML.
    Returns (R, G, B, A) where A is 0-255, or None if no solid fill."""
    sp_pr = sp_el.find(qn("p:spPr"))
    if sp_pr is None:
        return None
    solid = sp_pr.find(qn("a:solidFill"))
    if solid is None:
        return None

    rgb: tuple[int, int, int] | None = None
    alpha = 255

    rgb_el = solid.find(qn("a:srgbClr"))
    if rgb_el is not None:
        val = rgb_el.get("val", "")
        if len(val) == 6:
            rgb = (int(val[0:2], 16), int(val[2:4], 16), int(val[4:6], 16))
        # Check for alpha child (val is 0–100000, where 100000 = fully opaque)
        alpha_el = rgb_el.find(qn("a:alpha"))
        if alpha_el is not None:
            alpha = round(int(alpha_el.get("val", "100000")) / 100000 * 255)

    scheme_el = solid.find(qn("a:schemeClr"))
    if scheme_el is not None and rgb is None:
        name = scheme_el.get("val", "")
        rgb = scheme_colors.get(name)
        alpha_el = scheme_el.find(qn("a:alpha"))
        if alpha_el is not None:
            alpha = round(int(alpha_el.get("val", "100000")) / 100000 * 255)

    if rgb is None:
        return None
    return (*rgb, alpha)


def _render_shapes(
    shapes,
    img: Image.Image,
    draw: ImageDraw.ImageDraw,
    scheme_colors: dict,
) -> None:
    """Render a collection of shapes onto img (recursive for groups)."""
    for shape in shapes:
        left = _emu_to_px(shape.left or 0)
        top = _emu_to_px(shape.top or 0)
        width = _emu_to_px(shape.width or 0)
        height = _emu_to_px(shape.height or 0)

        # --- Groups: child coords are absolute, just recurse ---
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                _render_shapes(shape.shapes, img, draw, scheme_colors)
            except Exception:
                pass
            continue

        # --- Pictures ---
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                pic_img = Image.open(io.BytesIO(shape.image.blob)).convert("RGBA")
                pic_img = pic_img.resize((max(1, width), max(1, height)), Image.LANCZOS)
                img.paste(pic_img, (left, top), pic_img)
            except Exception:
                pass
            continue

        # --- Filled shapes (AUTO_SHAPE, FREEFORM, PLACEHOLDER) ---
        fill_color = _resolve_solid_fill_from_xml(shape._element, scheme_colors)
        if fill_color is not None:
            r, g, b, a = fill_color
            if a >= 250:
                # Fully opaque — draw directly
                draw.rectangle([left, top, left + width, top + height], fill=(r, g, b))
            elif a > 5:
                # Semi-transparent — composite via overlay layer
                overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
                ov_draw = ImageDraw.Draw(overlay)
                ov_draw.rectangle([left, top, left + width, top + height], fill=(r, g, b, a))
                img.paste(Image.alpha_composite(img.convert("RGBA"), overlay).convert("RGB"))

        # --- Text frames ---
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        y_cursor = top + _emu_to_px(tf.margin_top or 0)

        for para in tf.paragraphs:
            if not para.runs:
                y_cursor += round(_DPI * 14 / 72)
                continue

            try:
                size_pt = para.runs[0].font.size
                size_pt = Pt(size_pt).pt if size_pt else 12.0
            except Exception:
                size_pt = 12.0

            font = _get_font(size_pt)
            line_height = round(size_pt * _DPI / 72 * 1.2)

            full_text = "".join(r.text for r in para.runs)
            if not full_text.strip():
                y_cursor += line_height
                continue

            first_run_el = para.runs[0]._r
            # When no explicit color is set, default to dark (dk1) so text is
            # readable on light backgrounds. Slides with white-on-dark designs
            # always set the run color explicitly in the XML.
            text_color = _resolve_run_color(
                first_run_el,
                scheme_colors,
                scheme_colors.get("dk1", (0, 0, 0)),
            )

            words = full_text.split()
            line = ""
            x_start = left + _emu_to_px(tf.margin_left or 0)
            max_x = left + width - _emu_to_px(tf.margin_right or 0)

            for word in words:
                test = f"{line} {word}".strip() if line else word
                bbox = font.getbbox(test)
                if bbox[2] - bbox[0] <= max_x - x_start:
                    line = test
                else:
                    if line and y_cursor + line_height <= top + height + line_height:
                        draw.text((x_start, y_cursor), line, font=font, fill=text_color)
                    y_cursor += line_height
                    line = word

            if line and y_cursor + line_height <= top + height + line_height:
                draw.text((x_start, y_cursor), line, font=font, fill=text_color)
            y_cursor += line_height


def _bg_image_from_part(part, rel_id: str, slide_w: int, slide_h: int) -> Image.Image | None:
    """Extract and resize the background image referenced by a relationship ID."""
    try:
        rel = part.rels[rel_id]
        blob = rel.target_part.blob
        bg = Image.open(io.BytesIO(blob)).convert("RGB")
        bg = bg.resize((slide_w, slide_h), Image.LANCZOS)
        return bg
    except Exception:
        return None


def _render_background_element(bg_el, part, slide_w: int, slide_h: int, scheme_colors: dict) -> Image.Image | None:
    """
    Parse a <p:bg> element and return a fully-sized background image, or None.
    Handles solid fills, image (blip) fills, and gradient fills.
    """
    bgPr = bg_el.find(qn("p:bgPr"))
    if bgPr is None:
        return None

    # --- Image fill (blipFill) ---
    blipFill = bgPr.find(qn("a:blipFill"))
    if blipFill is not None:
        blip = blipFill.find(qn("a:blip"))
        if blip is not None:
            rel_id = blip.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
            if rel_id:
                return _bg_image_from_part(part, rel_id, slide_w, slide_h)

    # --- Solid fill ---
    solidFill = bgPr.find(qn("a:solidFill"))
    if solidFill is not None:
        rgb = None
        srgb = solidFill.find(qn("a:srgbClr"))
        if srgb is not None:
            val = srgb.get("val", "")
            if len(val) == 6:
                rgb = (int(val[0:2], 16), int(val[2:4], 16), int(val[4:6], 16))
        scheme = solidFill.find(qn("a:schemeClr"))
        if scheme is not None and rgb is None:
            rgb = scheme_colors.get(scheme.get("val", ""))
        if rgb:
            img = Image.new("RGB", (slide_w, slide_h), rgb)
            return img

    # --- Gradient fill (render as simple two-stop linear gradient) ---
    gradFill = bgPr.find(qn("a:gradFill"))
    if gradFill is not None:
        gsLst = gradFill.find(qn("a:gsLst"))
        stops: list[tuple[int, int, int]] = []
        if gsLst is not None:
            for gs in gsLst.findall(qn("a:gs")):
                srgb = gs.find(qn("a:srgbClr"))
                scheme = gs.find(qn("a:schemeClr"))
                rgb = None
                if srgb is not None:
                    val = srgb.get("val", "")
                    if len(val) == 6:
                        rgb = (int(val[0:2], 16), int(val[2:4], 16), int(val[4:6], 16))
                if scheme is not None and rgb is None:
                    rgb = scheme_colors.get(scheme.get("val", ""))
                if rgb:
                    stops.append(rgb)
        if len(stops) >= 2:
            img = Image.new("RGB", (slide_w, slide_h))
            arr = img.load()
            c1, c2 = stops[0], stops[-1]
            for y in range(slide_h):
                t = y / slide_h
                r = round(c1[0] + (c2[0] - c1[0]) * t)
                g = round(c1[1] + (c2[1] - c1[1]) * t)
                b = round(c1[2] + (c2[2] - c1[2]) * t)
                for x in range(slide_w):
                    arr[x, y] = (r, g, b)
            return img

    return None


def _get_slide_background(slide, slide_w: int, slide_h: int, scheme_colors: dict) -> Image.Image:
    """
    Resolve the slide background by checking the slide, then its layout, then the master.
    Falls back to lt1 (usually white) if no background is found anywhere.
    """
    # Check the slide itself
    bg_el = slide._element.find(qn("p:bg"))
    if bg_el is not None:
        result = _render_background_element(bg_el, slide.part, slide_w, slide_h, scheme_colors)
        if result is not None:
            return result

    # Check the slide layout
    try:
        layout = slide.slide_layout
        bg_el = layout._element.find(qn("p:bg"))
        if bg_el is not None:
            result = _render_background_element(bg_el, layout.part, slide_w, slide_h, scheme_colors)
            if result is not None:
                return result
    except Exception:
        pass

    # Check the slide master
    try:
        master = slide.slide_layout.slide_master
        bg_el = master._element.find(qn("p:bg"))
        if bg_el is not None:
            result = _render_background_element(bg_el, master.part, slide_w, slide_h, scheme_colors)
            if result is not None:
                return result
    except Exception:
        pass

    # Final fallback: lt1 (theme light color, typically white)
    bg_color = scheme_colors.get("lt1", (255, 255, 255))
    return Image.new("RGB", (slide_w, slide_h), bg_color)


def _render_slide(slide, slide_w: int, slide_h: int, scheme_colors: dict) -> Image.Image:
    """Render a single slide to a PIL Image."""
    img = _get_slide_background(slide, slide_w, slide_h, scheme_colors)
    draw = ImageDraw.Draw(img)
    _render_shapes(slide.shapes, img, draw, scheme_colors)
    return img





def render_slides(report_id: int, pptx_path: Path) -> Path:
    """
    Render a PPTX to PNG images using python-pptx + Pillow.
    Returns the directory containing slide_0.png, slide_1.png, etc.
    """
    slides_dir = get_slides_dir() / str(report_id)
    slides_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(pptx_path))
    slide_w = _emu_to_px(prs.slide_width)
    slide_h = _emu_to_px(prs.slide_height)
    scheme_colors = _resolve_scheme_colors(prs)

    for i, slide in enumerate(prs.slides):
        img = _render_slide(slide, slide_w, slide_h, scheme_colors)
        out = slides_dir / f"slide_{i}.png"
        img.save(str(out), "PNG")

    logger.info("Rendered %d slides for report_id=%s to %s", len(prs.slides), report_id, slides_dir)
    return slides_dir


def render_pdf(report_id: int, pptx_path: Path) -> Path:
    """
    Render a PPTX to a multi-page PDF by stitching PNG slide images via Pillow.
    Returns the path to the generated preview.pdf.
    """
    slides_dir = render_slides(report_id, pptx_path)
    pngs = sorted(slides_dir.glob("slide_*.png"), key=lambda p: int(p.stem.split("_")[1]))
    if not pngs:
        raise RuntimeError("No slides were rendered")

    images = []
    for p in pngs:
        try:
            images.append(Image.open(p).convert("RGB"))
        except Exception as e:
            logger.warning("Skipping unreadable slide image %s: %s", p.name, e)
    if not images:
        raise RuntimeError("No slide images could be read")
    canonical = slides_dir / "preview.pdf"
    images[0].save(
        str(canonical),
        "PDF",
        save_all=True,
        append_images=images[1:],
        resolution=_DPI,
    )
    logger.info("Rendered PDF preview for report_id=%s to %s", report_id, canonical)
    return canonical


def extract_slide_fields(pptx_path: Path) -> list[dict]:
    """
    Extract editable text fields from each slide.
    Returns a list of slide dicts:
      [{ "slide_index": 0, "fields": [{ "field_id": "s0_shape0_para0", "label": "...", "value": "..." }] }]
    """
    prs = Presentation(str(pptx_path))
    result = []

    for slide_idx, slide in enumerate(prs.slides):
        fields = []
        shapes_with_text = [s for s in slide.shapes if s.has_text_frame]
        for shape_idx, shape in enumerate(shapes_with_text):
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                text = para.text.strip()
                if not text:
                    continue

                field_id = f"s{slide_idx}_shape{shape_idx}_para{para_idx}"
                label = _label_for_field(slide_idx, shape.name, text)
                if label:
                    fields.append({
                        "field_id": field_id,
                        "label": label,
                        "value": text,
                        "slide_index": slide_idx,
                        "shape_name": shape.name,
                        "para_index": para_idx,
                    })

        result.append({"slide_index": slide_idx, "fields": fields})
    return result


def _label_for_field(slide_idx: int, shape_name: str, text: str) -> str | None:
    """Return a human-readable label for a text field, or None if it should not be editable."""
    if len(text) < 3:
        return None

    # Slide 1 — date
    if slide_idx == 0 and re.match(r"^[A-Za-z]+,?\s*\d{4}$", text):
        return "Report Month"

    # Slide 2 — executive summary
    if slide_idx == 1:
        if re.match(r"^\d+K?$", text):
            return "Active Users (short)"
        if re.match(r"^\d+%$", text) and len(text) <= 5:
            return "New Visitors %"
        if re.match(r"^\d+\.\d+%$", text):
            return "CTR"
        if len(text) > 60:
            return "Executive Summary Paragraph"

    # Slide 3 — site overview
    if slide_idx == 2:
        if "total active users" in text.lower():
            return "Active Users Label"
        if "new users" in text.lower():
            return "New Users Label"
        if "engagement" in text.lower():
            return "Engagement Time Label"
        if len(text) > 60:
            return "Site Overview Paragraph"
        if re.match(r"^[\d,]+$", text):
            return "Stat Value"

    # Slide 4 — geographic
    if slide_idx == 3 and len(text) > 40:
        return "Geographic Performance Paragraph"

    # Slide 5 — page performance
    if slide_idx == 4 and len(text) > 40:
        return "Page Performance Insight"

    # Slide 6 — search performance
    if slide_idx == 5 and len(text) > 40:
        return "Search Performance Paragraph"

    # Recommendations slide (index 6 for 8-slide, 5 for 7-slide)
    if slide_idx in (5, 6) and re.match(r"^\d\.", text):
        return f"Recommendation {text[0]}"

    return None


def extract_all_shapes(pptx_path: Path) -> list[dict]:
    """
    Extract ALL shapes from every slide (text + picture) without filtering.
    Used for the template mapping UI so the user can assign field types to each shape.
    Returns a flat list sorted by slide_index, then top/left position.
    """
    prs = Presentation(str(pptx_path))
    result = []

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                continue  # skip groups — leaf shapes are what matter for mapping

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_type = "image"
                placeholder_text = ""
            elif shape.has_text_frame:
                shape_type = "text"
                placeholder_text = shape.text_frame.text[:120].strip()
            else:
                continue  # connector, table, etc. — not mappable

            result.append({
                "slide_index": slide_idx,
                "shape_name": shape.name,
                "shape_type": shape_type,
                "placeholder_text": placeholder_text,
                "left_emu": shape.left,
                "top_emu": shape.top,
                "width_emu": shape.width,
                "height_emu": shape.height,
            })

    return result


def _render_via_spire(pptx_path: Path, target_dir: Path) -> None:
    """Render slides using Spire.Presentation (pure Python, high quality).
    Uses a short temp path to avoid Spire issues with spaces in directory names.
    """
    import tempfile, shutil
    from spire.presentation import Presentation as SpirePresentation  # type: ignore[import]

    with tempfile.TemporaryDirectory() as tmpdir:
        prs = SpirePresentation()
        prs.LoadFromFile(str(pptx_path))
        count = len(prs.Slides)
        try:
            for i in range(count):
                img = prs.Slides[i].SaveAsImage()
                img.Save(str(Path(tmpdir) / f"slide_{i}.png"))
                img.Dispose()
        finally:
            prs.Dispose()
        # Move rendered PNGs to the actual target dir
        for i in range(count):
            shutil.move(str(Path(tmpdir) / f"slide_{i}.png"), str(target_dir / f"slide_{i}.png"))
    logger.info("Rendered %d slides (Spire) to %s", count, target_dir)


def _render_slides_pillow(pptx_path: Path, target_dir: Path) -> None:
    prs = Presentation(str(pptx_path))
    slide_w = _emu_to_px(prs.slide_width)
    slide_h = _emu_to_px(prs.slide_height)
    scheme_colors = _resolve_scheme_colors(prs)
    for i, slide in enumerate(prs.slides):
        img = _render_slide(slide, slide_w, slide_h, scheme_colors)
        img.save(str(target_dir / f"slide_{i}.png"), "PNG")
    logger.info("Rendered %d slides (Pillow) to %s", len(prs.slides), target_dir)


def render_slides_to_dir(pptx_path: Path, target_dir: Path) -> None:
    """
    Render all slides of a PPTX to PNG files in target_dir.
    Uses Spire.Presentation for high-quality output, falls back to Pillow.
    """
    target_dir.mkdir(parents=True, exist_ok=True)
    try:
        _render_via_spire(pptx_path, target_dir)
    except Exception as exc:
        logger.warning("Spire rendering failed, falling back to Pillow: %s", exc)
        _render_slides_pillow(pptx_path, target_dir)


def apply_field_edits(pptx_path: Path, edits: dict[str, str], output_path: Path) -> None:
    """
    Re-open the PPTX, apply field edits (field_id -> new_text), and save to output_path.
    edits keys match the field_id format: "s{slide_idx}_shape{shape_idx}_para{para_idx}"
    """
    prs = Presentation(str(pptx_path))

    for field_id, new_text in edits.items():
        m = re.match(r"^s(\d+)_shape(\d+)_para(\d+)$", field_id)
        if not m:
            logger.warning("Unrecognised field_id format: %s", field_id)
            continue
        slide_idx, shape_idx, para_idx = int(m.group(1)), int(m.group(2)), int(m.group(3))

        if slide_idx >= len(prs.slides):
            continue
        slide = prs.slides[slide_idx]
        shapes_with_text = [s for s in slide.shapes if s.has_text_frame]
        if shape_idx >= len(shapes_with_text):
            continue
        shape = shapes_with_text[shape_idx]
        if para_idx >= len(shape.text_frame.paragraphs):
            continue
        para = shape.text_frame.paragraphs[para_idx]
        _fill_para(para, new_text)

    prs.save(str(output_path))


def _fill_para(para, new_text: str) -> None:
    """Replace paragraph text preserving first run's rPr formatting."""
    from copy import deepcopy
    import lxml.etree as etree
    from pptx.oxml.ns import qn

    if not para.runs:
        return
    first_rpr = para.runs[0]._r.find(qn("a:rPr"))
    saved_rpr = deepcopy(first_rpr) if first_rpr is not None else None
    for child in list(para._p):
        if child.tag != qn("a:pPr"):
            para._p.remove(child)
    new_r = etree.SubElement(para._p, qn("a:r"))
    if saved_rpr is not None:
        new_r.insert(0, saved_rpr)
    new_t = etree.SubElement(new_r, qn("a:t"))
    new_t.text = new_text
