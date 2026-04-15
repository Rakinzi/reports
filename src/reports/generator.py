"""
Core report generation pipeline:
1. Playwright captures screenshots + scrapes metrics from GA4
2. Gemini paraphrases the Slide 3 performance summary paragraph
3. python-pptx opens the template, replaces text (stats, dates, slide 3 paragraph) and images
4. Gemini compares current vs previous month data and writes recommendations
5. Saves the final PPTX to artifacts/output/
"""

import matplotlib
matplotlib.use('Agg')

import os
import re
import shutil
from pathlib import Path

from google import genai
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches

from .charts import (
    generate_country_bar_chart,
    generate_line_chart,
    generate_page_views_bar_chart,
    generate_traffic_source_pie_chart,
    generate_user_type_pie_chart,
)
from .logging_utils import configure_logging
from .runtime import (
    get_managed_chrome_profile_directory,
    get_managed_chrome_user_data_dir,
    get_output_dir,
    get_screenshots_dir,
    get_templates_dir,
    load_runtime_environment,
)
from .browser_support import build_launch_prefs

TEMPLATES_DIR = get_templates_dir()
OUTPUT_DIR = get_output_dir()
SCREENSHOTS_DIR = get_screenshots_dir()
logger = configure_logging()

GA4_PROPERTIES = {
    "cancer_serve": "454873082",
    "econet_ai":    "511212348",
    "infraco":      "516617515",
    "zimplats":     "385365994",
    "ecocash":      "386950925",
    "econet":       "386649040",
    "ecosure":      "384507667",
    "dicomm":       "382296904",
}

TEMPLATES = {
    "econet":       "Econet-February 2026 Website Report.pptx",
    "econet_ai":    "Econet AI February 2026 Website Report.pptx",
    "infraco":      "Econet_InfraCo_February_2026 (1).pptx",
    "ecocash":      "Ecocash February 2026 Website Report.pptx",
    "ecosure":      "Ecosure January 2026 Website Report (1).pptx",
    "zimplats":     "Zimplats February 2026 Website Report.pptx",
    "cancer_serve": "Cancer Serve February 2025 Website Report.pptx",
}


# GA4 sections to screenshot per report — (label, url fragment)
GA4_SECTIONS = [
    ("home",         "/home"),
    ("demographics", "/reports/user-demographics-overview"),
    ("acquisition",  "/reports/acquisition-overview"),
    ("engagement",   "/reports/engagement-overview"),
]


def _ga4_url(property_key: str, section_fragment: str) -> str:
    pid = GA4_PROPERTIES[property_key]
    return f"https://analytics.google.com/analytics/web/#/p{pid}{section_fragment}"


def _ga4_property_token(property_key: str) -> str:
    return f"p{GA4_PROPERTIES[property_key]}"


def _ga4_navigation_url(page, property_key: str, section_fragment: str) -> str:
    property_token = _ga4_property_token(property_key)
    match = re.search(rf"(https://analytics\.google\.com/analytics/web/#/[^/]*{re.escape(property_token)})", page.url)
    if match:
        return f"{match.group(1)}{section_fragment}"
    return _ga4_url(property_key, section_fragment)


def _ensure_expected_ga4_property(page, property_key: str, timeout: int = 15000) -> None:
    expected_token = _ga4_property_token(property_key)
    page.wait_for_function(
        "expected => window.location.href.includes(expected)",
        arg=expected_token,
        timeout=timeout,
    )


def _ga4_section_aliases(section_fragment: str) -> list[str]:
    if section_fragment == "/home":
        return ["/home", "/reports/intelligenthome"]
    return [section_fragment]


def _ensure_expected_ga4_location(page, property_key: str, section_fragment: str, timeout: int = 15000) -> None:
    expected_token = _ga4_property_token(property_key)
    expected_sections = _ga4_section_aliases(section_fragment)
    page.wait_for_function(
        """
        ({ expectedToken, expectedSections }) => {
            const href = window.location.href;
            return href.includes(expectedToken) && expectedSections.some(section => href.includes(section));
        }
        """,
        arg={"expectedToken": expected_token, "expectedSections": expected_sections},
        timeout=timeout,
    )


def _open_analytics_root(page) -> None:
    try:
        page.goto("https://analytics.google.com/analytics/web/", wait_until="domcontentloaded", timeout=30000)
    except Exception:
        page.goto("https://analytics.google.com/analytics/web/", wait_until="load", timeout=45000)
    # Wait for the search bar to actually appear instead of a fixed sleep
    for sel in [
        lambda: page.get_by_role("searchbox").first,
        lambda: page.get_by_role("textbox", name=re.compile("search", re.I)).first,
        lambda: page.locator('input[aria-label*="Search"], input[placeholder*="Search"]').first,
    ]:
        try:
            sel().wait_for(state="visible", timeout=12000)
            return
        except Exception:
            continue
    # fallback — give GA4 more time if nothing found yet
    page.wait_for_timeout(5000)


def _click_first_visible(locator) -> bool:
    count = locator.count()
    for index in range(count):
        candidate = locator.nth(index)
        try:
            candidate.wait_for(state="visible", timeout=1500)
            candidate.click()
            return True
        except Exception:
            continue
    return False


def _click_nth_visible(locator, target_visible_index: int) -> bool:
    visible_index = 0
    count = locator.count()
    for index in range(count):
        candidate = locator.nth(index)
        try:
            candidate.wait_for(state="visible", timeout=1500)
        except Exception:
            continue
        if visible_index == target_visible_index:
            candidate.click()
            return True
        visible_index += 1
    return False


def _resolve_post_click_page(page, previous_page_count: int):
    page.wait_for_timeout(2500)
    pages = page.context.pages
    if len(pages) > previous_page_count:
        new_page = pages[-1]
        try:
            new_page.wait_for_load_state("domcontentloaded", timeout=15000)
        except Exception:
            pass
        new_page.bring_to_front()
        logger.info("GA4 search result opened a new tab. current_url=%s", new_page.url)
        return new_page
    return page


def _switch_ga4_property_via_search(page, property_key: str):
    property_id = GA4_PROPERTIES[property_key]
    logger.info("Switching GA4 property via search. property_key=%s property_id=%s", property_key, property_id)

    _open_analytics_root(page)

    search_input = None
    search_selectors = [
        lambda: page.get_by_role("searchbox").first,
        lambda: page.get_by_role("textbox", name=re.compile("search", re.I)).first,
        lambda: page.locator('input[aria-label*="Search"], input[placeholder*="Search"]').first,
    ]
    for resolve in search_selectors:
        try:
            candidate = resolve()
            candidate.wait_for(state="visible", timeout=4000)
            search_input = candidate
            break
        except Exception:
            continue

    if search_input is None:
        raise RuntimeError("Could not find the GA4 search bar to switch properties.")

    url_before_fill = page.url
    search_input.click()
    page.wait_for_timeout(500)
    search_input.press_sequentially(property_id, delay=80)

    # Give GA4 a moment to react — it may auto-navigate via the Insights panel
    # (typing the property ID triggers an "Ask Analytics Advisor" flow that lands
    # directly on the correct property without requiring a dropdown click).
    page.wait_for_timeout(3000)

    # Only skip click logic if GA4 *navigated* to the correct property after the fill
    # (URL must have changed to include the property token).
    if f"p{property_id}" in page.url and page.url != url_before_fill:
        logger.info(
            "GA4 auto-navigated to correct property after search fill. property_id=%s url=%s",
            property_id, page.url,
        )
        try:
            search_input.press("Escape")
        except Exception:
            pass
        page.wait_for_timeout(500)
        return page

    result_candidates = page.locator('[role="option"], [role="menuitem"], a, button, li')
    # Match results whose own label contains "(GA4 Property <id>)" — avoids matching
    # results that only mention the id in a description sub-element.
    ga4_label_result = result_candidates.filter(
        has_text=re.compile(rf"GA4 Property\s+{re.escape(property_id)}", re.I)
    )
    # Fallback: any result whose text contains the bare property id
    exact_result = result_candidates.filter(
        has_text=re.compile(rf"\b{re.escape(property_id)}\b")
    )
    previous_page_count = len(page.context.pages)

    if _click_first_visible(ga4_label_result):
        logger.info("Selected GA4 property result by label for property_id=%s", property_id)
    elif _click_first_visible(exact_result):
        logger.info("Selected GA4 search result matching property_id=%s", property_id)
    elif _click_first_visible(result_candidates):
        logger.info("Selected first visible GA4 search result for property_id=%s", property_id)
    else:
        try:
            search_input.press("Enter")
            logger.info("Submitted GA4 search via Enter for property_id=%s", property_id)
        except Exception as exc:
            raise RuntimeError(
                f"Could not select a GA4 search result for property '{property_key}' ({property_id})."
            ) from exc

    page = _resolve_post_click_page(page, previous_page_count)
    _ensure_expected_ga4_property(page, property_key, timeout=20000)
    page.wait_for_timeout(2500)
    return page


def _switch_ga4_property_by_id(page, property_id: str):
    """Switch to a GA4 property using a raw numeric property ID (user-template path).
    Injects the id temporarily into GA4_PROPERTIES under a sentinel key."""
    _SENTINEL = "__user_template__"
    GA4_PROPERTIES[_SENTINEL] = property_id
    try:
        page = _switch_ga4_property_via_search(page, _SENTINEL)
    finally:
        GA4_PROPERTIES.pop(_SENTINEL, None)
    return page


def _goto_ga4_section(page, property_key: str, section_fragment: str, timeout: int = 45000):
    expected_token = _ga4_property_token(property_key)
    expected_sections = _ga4_section_aliases(section_fragment)

    last_error: Exception | None = None
    for attempt in range(1, 4):
        try:
            # Only switch property via search if not already on the correct property.
            if expected_token not in page.url:
                page = _switch_ga4_property_via_search(page, property_key)
            url = _ga4_navigation_url(page, property_key, section_fragment)
            logger.info(
                "Navigating to GA4 property=%s section=%s attempt=%s url=%s",
                GA4_PROPERTIES[property_key],
                section_fragment,
                attempt,
                url,
            )
            try:
                page.goto(url, wait_until="domcontentloaded", timeout=30000)
            except Exception:
                page.goto(url, wait_until="load", timeout=timeout)

            _ensure_expected_ga4_location(page, property_key, section_fragment, timeout=15000)
            page.wait_for_timeout(3000)
            logger.info(
                "GA4 navigation confirmed for property=%s section=%s current_url=%s",
                GA4_PROPERTIES[property_key],
                section_fragment,
                page.url,
            )
            return page
        except Exception as exc:
            last_error = exc
            logger.warning(
                "GA4 navigation mismatch for property=%s section=%s attempt=%s current_url=%s error=%s",
                GA4_PROPERTIES[property_key],
                section_fragment,
                attempt,
                page.url,
                exc,
            )
            try:
                page.goto("https://analytics.google.com/analytics/web/", wait_until="domcontentloaded", timeout=20000)
                page.wait_for_timeout(2000)
            except Exception:
                pass

    current_url = page.url
    raise RuntimeError(
        f"Could not open the selected GA4 property for '{property_key}'. "
        f"Expected URL containing '{expected_token}' and one of {expected_sections}, "
        f"but Analytics stayed on '{current_url}'."
    ) from last_error


def _launch_persistent_context(playwright, headless: bool = False):
    load_runtime_environment()
    chrome_user_data_dir = os.getenv(
        "CHROME_USER_DATA_DIR",
        str(get_managed_chrome_user_data_dir().resolve()),
    )
    chrome_profile_directory = os.getenv(
        "CHROME_PROFILE_DIRECTORY", get_managed_chrome_profile_directory()
    )

    Path(chrome_user_data_dir).mkdir(parents=True, exist_ok=True)
    stealth_args = [
        f"--profile-directory={chrome_profile_directory}",
        "--disable-blink-features=AutomationControlled",
    ]
    last_error: Exception | None = None
    for launch_pref in build_launch_prefs():
        try:
            return playwright.chromium.launch_persistent_context(
                user_data_dir=chrome_user_data_dir,
                headless=headless,
                args=stealth_args,
                ignore_default_args=["--enable-automation"],
                viewport={"width": 1920, "height": 1080},
                **launch_pref,
            )
        except Exception as exc:
            last_error = exc
    raise RuntimeError(
        "Could not launch a compatible Chromium browser. Install Google Chrome, Microsoft Edge, or Chromium."
    ) from last_error


def _set_date_range(page, start: str, end: str) -> None:
    # Dismiss any GA4 modals/tooltips (welcome dialogs, what's new, guided tours)
    for dismiss_sel in [
        "button[aria-label='Close']",
        "button[aria-label='Dismiss']",
        "button:has-text('Got it')",
        "button:has-text('Dismiss')",
        "button:has-text('Close')",
        "button:has-text('No thanks')",
        "[data-guidedhelpid='agi-dismiss-button']",
        "gds-guide-button[dismiss]",
    ]:
        try:
            btn = page.locator(dismiss_sel).first
            if btn.is_visible(timeout=1500):
                btn.click()
                page.wait_for_timeout(500)
        except Exception:
            continue

    date_picker = page.get_by_role("combobox", name="Open date range picker").first
    date_picker.wait_for(state="visible", timeout=15000)
    date_picker.click()
    page.wait_for_timeout(1000)
    page.get_by_role("menuitem").filter(has_text="Custom").click()
    page.wait_for_timeout(1000)
    start_input = page.get_by_label("Start date")
    start_input.wait_for(state="visible", timeout=10000)
    start_input.click()
    start_input.select_text()
    start_input.fill(start)
    page.keyboard.press("Tab")
    page.wait_for_timeout(500)
    end_input = page.get_by_label("End date")
    end_input.click()
    end_input.select_text()
    end_input.fill(end)
    page.keyboard.press("Tab")
    page.wait_for_timeout(500)
    page.get_by_role("button", name="Apply").click()
    page.wait_for_timeout(3000)


def _scrape_home_metrics(page) -> dict:
    page.wait_for_selector(".metric-container", timeout=15000)
    page.wait_for_timeout(1000)
    metrics = {}
    cards = page.locator(".metric-container [aria-label]").all()
    for card in cards:
        try:
            aria = card.get_attribute("aria-label")
            if not aria:
                continue
            match = re.match(r"^(.+?)\s+([\d,\.]+(?:m \d+s|\s?s|K|M)?),?\s*", aria)
            if match:
                metrics[match.group(1).strip()] = match.group(2).strip().rstrip(",")
        except Exception:
            continue
    return metrics


def _scrape_snapshot_metrics(page) -> dict:
    page.wait_for_timeout(3000)
    channels = {}
    for bar in page.locator("g.bargroup[aria-label]").all():
        try:
            aria = bar.get_attribute("aria-label")
            if not aria:
                continue
            match = re.match(r"^(.+?)\s+(\d[\d,]*)\s+", aria)
            if match:
                channels[match.group(1).strip()] = match.group(2).strip()
        except Exception:
            continue

    top_country = top_country_users = None
    countries: dict[str, int] = {}
    try:
        body_text = page.locator("body").inner_text()
        lines = [l.strip() for l in body_text.splitlines() if l.strip()]
        for i, line in enumerate(lines):
            if line.startswith("COUNTRY\t"):
                # Parse all country rows until we hit a non-data line
                j = i + 1
                while j < len(lines):
                    parts = lines[j].split("\t")
                    if len(parts) == 2 and parts[1].strip().isdigit():
                        country, users = parts[0].strip(), int(parts[1].strip())
                        countries[country] = users
                        if top_country is None:
                            top_country = country
                            top_country_users = str(users)
                    else:
                        break
                    j += 1
                break
    except Exception:
        pass

    return {
        "channels": channels,
        "top_country": top_country,
        "top_country_users": top_country_users,
        "countries": countries,
    }


def _generate_slide3_paragraph(report_name: str, home_metrics: dict, snapshot_metrics: dict) -> str:
    """Use Gemini to paraphrase the Slide 3 performance summary paragraph."""
    load_runtime_environment()
    active_users    = home_metrics.get("Active users", "N/A")
    new_users       = home_metrics.get("New users", "N/A")
    returning_users = home_metrics.get("Returning users", "N/A")
    engagement      = home_metrics.get("Average engagement time per active user", "N/A")
    channels        = snapshot_metrics.get("channels", {})
    organic         = channels.get("Organic Search", "N/A")
    direct          = channels.get("Direct", "N/A")
    top_country     = snapshot_metrics.get("top_country") or "N/A"
    top_country_n   = snapshot_metrics.get("top_country_users") or "N/A"
    brand           = report_name.replace("_", " ").title()

    raw = (
        f"In the past month, the {brand} Website attracted a total of {active_users} active users, "
        f"with {new_users} being new users and {returning_users} returning users who visited the site "
        f"multiple times. The majority of users accessed the site through organic search ({organic} users), "
        f"while smaller numbers came via direct search ({direct}). "
        f"{top_country} contributed the highest traffic, with {top_country_n} users. "
        f"On average, users engaged with the site for {engagement} per session."
    )

    client = genai.Client(api_key=os.environ["GEMINI_API_KEY"])
    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=(
            "Paraphrase the following website analytics summary for a professional PowerPoint report. "
            "Keep all numbers exactly as they are. Use clear, formal business English. "
            "Do not use em dashes, bullet points, or markdown. Output a single plain paragraph only.\n\n"
            + raw
        ),
    )
    return response.text.strip()


def _generate_slide4_texts(home_metrics: dict) -> dict:
    """Use Gemini to paraphrase each slide 4 text block using home metrics."""
    load_runtime_environment()
    active_users    = home_metrics.get("Active users", "N/A")
    new_users       = home_metrics.get("New users", "N/A")
    engagement      = home_metrics.get("Average engagement time per active user", "N/A")

    raw_intro = (
        f"Emphasizing the important numerical data, the large line chart on the right displays "
        f"the key metrics for the Econet AI website from the previous month."
    )
    raw_users       = f"Users: {active_users} total users visited the site."
    raw_new_users   = f"New Users: {new_users} new users, indicating most visitors were first-time visitors."
    raw_engagement  = f"Average Engagement Time: Users spent an average of {engagement} on the site."

    client = genai.Client(api_key=os.environ["GEMINI_API_KEY"])

    def paraphrase(text: str) -> str:
        resp = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=(
                "Paraphrase the following sentence for a professional PowerPoint report. "
                "Keep all numbers exactly as they are. Use clear, formal business English. "
                "No em dashes, bullets, or markdown. Output one plain sentence only.\n\n" + text
            ),
        )
        return resp.text.strip()

    return {
        "intro":      paraphrase(raw_intro),
        "users":      paraphrase(raw_users),
        "new_users":  paraphrase(raw_new_users),
        "engagement": paraphrase(raw_engagement),
    }


def _generate_slide7_text(slide, acquisition: dict) -> str:
    """Read the existing slide 7 paragraph, inject live stats, and let Gemini paraphrase it."""
    load_runtime_environment()

    # Read existing text from the slide
    original = ""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = "".join(r.text for r in para.runs).strip()
            if "traffic channels" in text.lower() or "pie chart" in text.lower():
                original = text
                break

    organic = acquisition.get("Organic Search", 0)
    direct = acquisition.get("Direct", 0)
    top_channel = "Organic Search" if organic >= direct else "Direct"

    client = genai.Client(api_key=os.environ["GEMINI_API_KEY"])
    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=(
            "Paraphrase the following paragraph for a professional PowerPoint report. "
            f"Update the statistics: Organic Search has {organic} users and Direct has {direct} users. "
            f"The dominant traffic channel is {top_channel}. "
            "Keep all numbers exactly as they are. Use clear, formal business English. "
            "No em dashes, bullets, or markdown. Output one plain paragraph only.\n\n" + original
        ),
    )
    return response.text.strip()


def _shorten_page_names(page_views: dict) -> dict:
    """Shorten page names: split on hyphen and take the first part if present,
    then use Gemini for any name still longer than 20 characters."""
    load_runtime_environment()
    shortened = {}
    for name, views in page_views.items():
        short = name.split("-")[0].strip() if "-" in name else name
        shortened[short] = views

    long_names = [name for name in shortened if len(name) > 20]
    if not long_names:
        return shortened

    client = genai.Client(api_key=os.environ["GEMINI_API_KEY"])
    for name in long_names:
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=(
                "Shorten the following page name to at most 20 characters for a bar chart label. "
                "Keep it meaningful and recognizable. Output only the shortened name, nothing else.\n\n"
                + name
            ),
        )
        short = response.text.strip()
        shortened[short] = shortened.pop(name)
    return shortened


def _generate_slide5_text(countries: dict) -> str:
    """Use Gemini to paraphrase the slide 5 country traffic statement using live data."""
    load_runtime_environment()
    top5 = sorted(countries.items(), key=lambda x: x[1], reverse=True)[:5]

    parts = []
    for i, (country, users) in enumerate(top5):
        if i == 0:
            parts.append(f"{country} ({users} users)")
        else:
            parts.append(f"{country} with ({users} users)")

    raw = "Most of the traffic is coming from " + ", followed by ".join(parts) + "."

    client = genai.Client(api_key=os.environ["GEMINI_API_KEY"])
    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=(
            "Paraphrase the following sentence for a professional PowerPoint report. "
            "Keep all country names and numbers exactly as they are. "
            "Use clear, formal business English. "
            "No em dashes, bullets, or markdown. Output one plain sentence only.\n\n" + raw
        ),
    )
    return response.text.strip()



def _weekly_ranges(start_date: str, end_date: str) -> list[tuple[str, str, str]]:
    """Split a date range into 4 weekly chunks.
    Returns list of (label, week_start, week_end) in GA4 picker format e.g. 'Feb 1, 2026'.
    """
    from datetime import datetime, timedelta
    fmt = "%b %d, %Y"
    start = datetime.strptime(start_date, fmt)
    end = datetime.strptime(end_date, fmt)
    weeks = []
    week_start = start
    for i in range(1, 5):
        week_end = min(week_start + timedelta(days=6), end)
        weeks.append((
            f"Week {i}",
            week_start.strftime(fmt),
            week_end.strftime(fmt),
        ))
        week_start = week_end + timedelta(days=1)
        if week_start > end:
            break
    return weeks


def capture_screenshots_and_metrics(
    report_name: str,
    start_date: str,
    end_date: str,
) -> tuple[dict[str, Path], dict, dict, dict, dict, dict]:
    """Navigate GA4, set date range, scrape metrics, and take screenshots.
    Returns (screenshots, home_metrics, snapshot_metrics, acquisition, page_views, weekly_active_users).
    start_date/end_date format: 'Feb 1, 2026'
    """

    screenshots: dict[str, Path] = {}
    out_dir = SCREENSHOTS_DIR / report_name
    out_dir.mkdir(parents=True, exist_ok=True)

    with sync_playwright() as p:
        context = _launch_persistent_context(p, headless=False)
        try:
            page = context.new_page()
            page.bring_to_front()

            # --- Home page: set date range + scrape metrics ---
            page = _goto_ga4_section(page, report_name, "/home")
            # Dismiss any lingering search/insights panel before interacting with the page.
            try:
                page.keyboard.press("Escape")
                page.wait_for_timeout(1000)
            except Exception:
                pass
            _set_date_range(page, start_date, end_date)
            home_metrics = _scrape_home_metrics(page)

            # --- Screenshot of the home line chart only ---
            try:
                chart_el = page.locator("ga-card.card_0 ga-tab-chart")
                chart_el.wait_for(state="visible", timeout=10000)
                path = out_dir / "home_chart.png"
                chart_el.screenshot(path=str(path))
                screenshots["home_chart"] = path
            except Exception:
                path = out_dir / "home_chart.png"
                page.screenshot(path=str(path), full_page=False)
                screenshots["home_chart"] = path

            # --- Weekly active users for line chart (slide 9) ---
            weekly_active_users: dict[str, int] = {}
            for label, w_start, w_end in _weekly_ranges(start_date, end_date):
                try:
                    _set_date_range(page, w_start, w_end)
                    w_metrics = _scrape_home_metrics(page)
                    val = w_metrics.get("Active users", "0")
                    weekly_active_users[label] = int(str(val).replace(",", ""))
                except Exception:
                    weekly_active_users[label] = 0

            # Restore full date range before continuing
            _set_date_range(page, start_date, end_date)

            # --- Reports Snapshot: set date range + scrape metrics ---
            _ensure_expected_ga4_property(page, report_name)
            page.get_by_text("View reports snapshot").click()
            page.wait_for_timeout(3000)
            _ensure_expected_ga4_property(page, report_name)
            snapshot_date_btn = page.get_by_role("combobox", name="Open date range picker")
            snapshot_date_btn.wait_for(state="visible", timeout=15000)
            snapshot_date_btn.click()
            page.wait_for_timeout(1000)
            page.get_by_role("menuitem").filter(has_text="Custom").click()
            page.wait_for_timeout(1000)
            start_input = page.get_by_label("Start date")
            start_input.wait_for(state="visible", timeout=10000)
            start_input.click()
            start_input.select_text()
            start_input.fill(start_date)
            page.keyboard.press("Tab")
            page.wait_for_timeout(500)
            end_input = page.get_by_label("End date")
            end_input.click()
            end_input.select_text()
            end_input.fill(end_date)
            page.keyboard.press("Tab")
            page.wait_for_timeout(500)
            page.get_by_role("button", name="Apply").click()
            page.wait_for_timeout(3000)
            snapshot_metrics = _scrape_snapshot_metrics(page)

            # --- Snapshot KPI card screenshot (Active users / New users / Avg engagement + line chart) ---
            try:
                card_el = page.locator("ga-card[data-guidedhelpid='summary']").first
                card_el.wait_for(state="visible", timeout=10000)
                page.mouse.move(0, 0)
                page.mouse.click(0, 0)
                page.wait_for_timeout(800)
                path = out_dir / "snapshot_card.png"
                card_el.screenshot(path=str(path))
                screenshots["snapshot_card"] = path
            except Exception:
                pass

            # --- Countries table screenshot for slide 6 ---
            try:
                page.get_by_text("View countries").click()
                page.wait_for_timeout(4000)
                _ensure_expected_ga4_property(page, report_name)
                row_num_col = page.locator("th.cdk-column-__row_index__").first
                end_col = page.locator("th.cdk-column-DEFAULT-engagedSessionsPerUser").first
                row_num_col.wait_for(state="visible", timeout=10000)
                page.keyboard.press("End")
                page.wait_for_timeout(1000)
                page.mouse.wheel(0, 3000)
                page.wait_for_timeout(1000)
                start_box = row_num_col.bounding_box()
                end_box = end_col.bounding_box()
                table_box = page.locator("table.adv-table").bounding_box()
                clip = {
                    "x": start_box["x"],
                    "y": table_box["y"],
                    "width": (end_box["x"] + end_box["width"]) - start_box["x"],
                    "height": table_box["height"],
                }
                path = out_dir / "countries_table.png"
                page.screenshot(path=str(path), clip=clip, full_page=True)
                screenshots["countries_table"] = path
            except Exception:
                pass

            # --- User acquisition page: scrape channels + generate pie chart ---
            acquisition: dict[str, int] = {}
            try:
                page.go_back()
                page.wait_for_timeout(3000)
                _ensure_expected_ga4_property(page, report_name)
                page.get_by_role("button", name="View user acquisition", exact=True).click()
                page.wait_for_timeout(4000)
                _ensure_expected_ga4_property(page, report_name)
                body = page.locator("body").inner_text()
                for line in [l.strip() for l in body.splitlines() if l.strip()]:
                    match = re.match(r"^\d+\t(.+?)\t(\d+)\s*\(", line)
                    if match:
                        acquisition[match.group(1).strip()] = int(match.group(2))
                pie_data = {k: v for k, v in acquisition.items() if k in ("Organic Search", "Direct")}
                if pie_data:
                    pie_path = out_dir / "traffic_pie.png"
                    generate_traffic_source_pie_chart(pie_data, str(pie_path))
                    screenshots["traffic_pie"] = pie_path
            except Exception:
                pass

            # --- Pages and screens table screenshot for slide 8 + scrape page views for slide 9 ---
            page_views: dict[str, int] = {}
            try:
                page.go_back()
                page.wait_for_timeout(3000)
                _ensure_expected_ga4_property(page, report_name)
                page.get_by_role("button", name="View pages and screens", exact=True).click()
                page.wait_for_timeout(4000)
                _ensure_expected_ga4_property(page, report_name)
                row_num_col = page.locator("th.cdk-column-__row_index__").first
                end_col = page.locator("th.cdk-column-DEFAULT-userEngagementDurationPerUser").first
                row_num_col.wait_for(state="visible", timeout=10000)
                page.keyboard.press("End")
                page.wait_for_timeout(1000)
                page.mouse.wheel(0, 3000)
                page.wait_for_timeout(1000)
                start_box = row_num_col.bounding_box()
                end_box = end_col.bounding_box()
                table_box = page.locator("table.adv-table").bounding_box()
                clip = {
                    "x": start_box["x"],
                    "y": table_box["y"],
                    "width": (end_box["x"] + end_box["width"]) - start_box["x"],
                    "height": table_box["height"],
                }
                path = out_dir / "pages_table.png"
                page.screenshot(path=str(path), clip=clip, full_page=True)
                screenshots["pages_table"] = path

                # Scrape top 4 page names + view counts
                body = page.locator("body").inner_text()
                for line in [l.strip() for l in body.splitlines() if l.strip()]:
                    match = re.match(r"^\d+\t(.+?)\t(\d+)\s*", line)
                    if match and len(page_views) < 4:
                        page_views[match.group(1).strip()] = int(match.group(2).strip())
            except Exception:
                pass

            # --- Screenshots for remaining GA4 sections ---
            for label, fragment in GA4_SECTIONS:
                page = _goto_ga4_section(page, report_name, fragment)
                path = out_dir / f"{label}.png"
                page.screenshot(path=str(path), full_page=False)
                screenshots[label] = path

        finally:
            context.close()

    return screenshots, home_metrics, snapshot_metrics, acquisition, page_views, weekly_active_users


# Matches full dates like "02 March 2026" or "1 July 2025"
DATE_PATTERN = re.compile(r"\b\d{1,2}\s+\w+\s+\d{4}\b")
# Matches month+year only like "February 2026"
MONTH_YEAR_PATTERN = re.compile(r"\b(January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{4}\b")


def _performance_month(date_range: str) -> str:
    """Extract 'Month YYYY' from a date range string e.g. '1 February 2026 - 28 February 2026' -> 'February 2026'"""
    match = MONTH_YEAR_PATTERN.search(date_range)
    return match.group(0) if match else ""


def _replace_text_in_slide(slide, replacements: dict[str, str], report_date: str, performance_month: str) -> None:
    """Replace text in all shapes of a slide while preserving formatting.
    - Any DD Month YYYY pattern -> report_date
    - Any Month YYYY pattern -> performance_month (applied after full date swap to avoid double replacing)
    - Any explicit replacements in the replacements dict
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                text = run.text
                # Replace full dates first (e.g. "02 March 2026")
                text = DATE_PATTERN.sub(report_date, text)
                # Then replace month+year (e.g. "February 2026") with the performance month
                if performance_month:
                    text = MONTH_YEAR_PATTERN.sub(performance_month, text)
                # Apply explicit replacements
                for old, new in replacements.items():
                    text = text.replace(old, new)
                run.text = text


def _replace_image_in_slide(slide, image_path: Path, shape_index: int = None, shape_name: str = None) -> None:
    """Replace an existing picture shape with a new image at the same position/size.
    Looks up by shape_name if provided, otherwise falls back to shape_index."""
    pictures = [s for s in slide.shapes if s.shape_type == 13]  # 13 = picture

    pic = None
    if shape_name:
        for p in pictures:
            if p.name == shape_name:
                pic = p
                break
    if pic is None and shape_index is not None and shape_index < len(pictures):
        pic = pictures[shape_index]
    if pic is None:
        return

    left, top, width, height = pic.left, pic.top, pic.width, pic.height
    pic._element.getparent().remove(pic._element)
    slide.shapes.add_picture(str(image_path), left, top, width, height)


def _get_previous_month_data(report_name: str) -> str:
    """Extract text from the previous month's template to give Gemini context."""
    template_path = TEMPLATES_DIR / TEMPLATES[report_name]
    prs = Presentation(str(template_path))
    lines = []
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if line:
                    lines.append(line)
    return "\n".join(lines)


def _generate_recommendations(
    report_name: str,
    current_stats: str,
    previous_data: str,
    date_range: str,
) -> list[str]:
    """Use Gemini to compare current vs previous month and generate 3 recommendations."""
    client = genai.Client()

    prompt = f"""
You are a digital analytics expert writing a monthly website performance report for {report_name.replace("_", " ").title()}.

Previous month data:
{previous_data}

Current month ({date_range}) stats:
{current_stats}

Compare the two months and write exactly 3 concise, actionable recommendations.
Format each as a single sentence starting with a number e.g. "1. ...", "2. ...", "3. ..."
Focus on traffic trends, user engagement, and acquisition channels.
"""

    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=prompt,
    )
    text = response.text.strip()

    recommendations = []
    for line in text.splitlines():
        line = line.strip()
        if re.match(r"^\d+\.", line):
            recommendations.append(line)

    return recommendations[:3]


UNSUPPORTED_REPORTS = {"union_hardware"}


def _fill_text_run(para, new_text: str) -> None:
    """Replace a paragraph's text with new_text, preserving the first run's rPr formatting."""
    from pptx.oxml.ns import qn
    from copy import deepcopy
    import lxml.etree as etree

    if not para.runs:
        return
    first_rpr = para.runs[0]._r.find(qn("a:rPr"))
    saved_rpr = deepcopy(first_rpr) if first_rpr is not None else None
    for child in list(para._p):
        if child.tag != qn("a:pPr"):
            para._p.remove(child)
    new_r = etree.SubElement(para._p, qn("a:r"))
    if saved_rpr is not None:
        new_r.insert(0, saved_rpr)
    new_t = etree.SubElement(new_r, qn("a:t"))
    new_t.text = new_text




def generate_report(
    report_name: str,
    date_range: str,
    report_date: str,
    start_date: str,
    end_date: str,
    _stage_callback=None,
) -> Path:
    """Full pipeline: screenshots + metrics → slide 3 paragraph → pptx edit → recommendations → save.
    date_range: human label e.g. '1 February 2026 - 28 February 2026'
    start_date / end_date: GA4 picker format e.g. 'Feb 1, 2026'
    """
    if report_name in UNSUPPORTED_REPORTS:
        raise NotImplementedError(
            f"'{report_name}' has a unique template structure and is not yet supported. "
            "A dedicated generator will be built for it."
        )

    # Dispatch to 2026 pipeline for new-format templates
    from .generator_2026 import TEMPLATES_2026, generate_report_2026
    if report_name in TEMPLATES_2026:
        return generate_report_2026(
            report_name=report_name,
            date_range=date_range,
            report_date=report_date,
            start_date=start_date,
            end_date=end_date,
        )

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    # Step 1: Capture screenshots + scrape GA4 metrics
    screenshots, home_metrics, snapshot_metrics, acquisition, page_views, weekly_active_users = capture_screenshots_and_metrics(
        report_name, start_date, end_date
    )

    # Step 2: Generate Slide 3 paragraph via Gemini
    slide3_paragraph = _generate_slide3_paragraph(report_name, home_metrics, snapshot_metrics)

    # Step 3: Load template
    template_path = TEMPLATES_DIR / TEMPLATES[report_name]
    prs = Presentation(str(template_path))

    # Step 4: Derive the performance month from the date range e.g. "February 2026"
    performance_month = _performance_month(date_range)

    # Step 5: Apply text replacements to all slides
    for slide in prs.slides:
        _replace_text_in_slide(slide, {}, report_date, performance_month)

    # Step 6: Replace Slide 3 paragraph (slide index 2, 0-based)
    # Preserves the formatting (font, size, bold, color) of the first run in the paragraph
    if len(prs.slides) > 2:
        slide3 = prs.slides[2]
        for shape in slide3.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                full_text = "".join(run.text for run in para.runs)
                if "past month" in full_text.lower() or "attracted" in full_text.lower():
                    if para.runs:
                        # Keep the first run's XML formatting, put all new text there
                        first_run = para.runs[0]
                        first_run.text = slide3_paragraph
                        # Remove all subsequent runs from the paragraph XML
                        from pptx.oxml.ns import qn
                        p_elem = para._p
                        for r_elem in p_elem.findall(qn("a:r"))[1:]:
                            p_elem.remove(r_elem)
                    break

    # Step 6b: Replace Slide 4 text using Gemini-paraphrased metrics
    if len(prs.slides) > 3:
        from pptx.oxml.ns import qn as _qn
        slide4_texts = _generate_slide4_texts(home_metrics)
        slide4 = prs.slides[3]
        for shape in slide4.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(r.text for r in para.runs)
                if "line chart" in text.lower() or "displays the key metrics" in text.lower():
                    replacement = slide4_texts["intro"]
                elif text.strip().lower().startswith("users :") or text.strip().lower().startswith("users:"):
                    replacement = slide4_texts["users"]
                elif "new users" in text.lower():
                    replacement = slide4_texts["new_users"]
                elif "engagement time" in text.lower() or ("average" in text.lower() and "spent" in text.lower()):
                    replacement = slide4_texts["engagement"]
                else:
                    continue
                if para.runs:
                    para.runs[0].text = replacement
                    for r_elem in para._p.findall(_qn("a:r"))[1:]:
                        para._p.remove(r_elem)

    # Step 6c: Generate country bar chart for slide 5 (top 5 countries only)
    countries = snapshot_metrics.get("countries", {})
    if countries:
        top5 = dict(sorted(countries.items(), key=lambda x: x[1], reverse=True)[:5])
        country_chart_path = SCREENSHOTS_DIR / report_name / "country_chart.png"
        generate_country_bar_chart(top5, str(country_chart_path))
        screenshots["country_chart"] = country_chart_path

    # Step 6d: Replace Slide 5 country text (shape 1, para 0)
    if len(prs.slides) > 4 and countries:
        from pptx.oxml.ns import qn as _qn2
        slide5_text = _generate_slide5_text(countries)
        slide5 = prs.slides[4]
        for shape in slide5.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(r.text for r in para.runs)
                if "traffic is coming from" in text.lower() or "most of the traffic" in text.lower():
                    if para.runs:
                        para.runs[0].text = slide5_text
                        for r_elem in para._p.findall(_qn2("a:r"))[1:]:
                            para._p.remove(r_elem)
                    break

    # Step 6e: Update slide 7 acquisition table (shape 3) with live data
    if len(prs.slides) > 6 and acquisition:
        slide7 = prs.slides[6]
        for shape in slide7.shapes:
            if shape.shape_type == 19:  # TABLE
                table = shape.table
                organic = acquisition.get("Organic Search", "")
                direct = acquisition.get("Direct", "")
                if len(table.rows) > 1:
                    table.rows[1].cells[1].text = str(organic)
                if len(table.rows) > 2:
                    table.rows[2].cells[1].text = str(direct)
                break

    # Step 6f: Replace slide 7 description paragraph (shape 4, para 0) with Gemini paraphrase
    if len(prs.slides) > 6 and acquisition:
        from pptx.oxml.ns import qn as _qn3
        slide7 = prs.slides[6]
        slide7_text = _generate_slide7_text(slide7, acquisition)
        for shape in slide7.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(r.text for r in para.runs)
                if "traffic channels" in text.lower() or "pie chart reveals" in text.lower():
                    if para.runs:
                        para.runs[0].text = slide7_text
                        for r_elem in para._p.findall(_qn3("a:r"))[1:]:
                            para._p.remove(r_elem)
                    break

    # Step 6g: Generate user type pie chart for slide 9 (new vs returning users)
    new_users = home_metrics.get("New users")
    returning_users = home_metrics.get("Returning users")
    if new_users and returning_users:
        try:
            user_type_data = {"New Users": int(new_users), "Returning Users": int(returning_users)}
            user_type_chart_path = SCREENSHOTS_DIR / report_name / "user_type_pie.png"
            generate_user_type_pie_chart(user_type_data, str(user_type_chart_path))
            screenshots["user_type_pie"] = user_type_chart_path
        except Exception:
            pass

    # Step 6h: Generate page views bar chart for slide 9
    if page_views:
        shortened_page_views = _shorten_page_names(page_views)
        sorted_page_views = dict(sorted(shortened_page_views.items(), key=lambda x: x[1], reverse=True))
        page_views_chart_path = SCREENSHOTS_DIR / report_name / "page_views_chart.png"
        generate_page_views_bar_chart(sorted_page_views, str(page_views_chart_path))
        screenshots["page_views_chart"] = page_views_chart_path

    # Step 7: Replace screenshots
    # Format: slide_idx (0-based) -> (screenshot_label, picture_index_on_slide)
    screenshot_slide_map = {
        0: ("home_chart",      "Picture 14"),  # Slide 1 - cover GA4 line chart
        3: ("home_chart",      "Picture 10"),  # Slide 4 - line chart
        4: ("country_chart",   "Picture 8"),   # Slide 5 - country bar chart
        5: ("countries_table", "Picture 6"),   # Slide 6 - countries table
        6: ("traffic_pie",     "Picture 9"),   # Slide 7 - traffic pie chart
        7: ("pages_table",     "Picture 10"),  # Slide 8 - pages and screens table
    }
    for slide_idx, (label, name) in screenshot_slide_map.items():
        if label in screenshots and slide_idx < len(prs.slides):
            _replace_image_in_slide(prs.slides[slide_idx], screenshots[label], shape_name=name)

    # Step 6i: Generate weekly active users line chart for slide 9
    if weekly_active_users:
        line_chart_path = SCREENSHOTS_DIR / report_name / "weekly_line_chart.png"
        generate_line_chart(weekly_active_users, str(line_chart_path))
        screenshots["weekly_line_chart"] = line_chart_path

    # Slide 9 has three image replacements — use shape names to avoid index shifting
    if len(prs.slides) > 8:
        slide9 = prs.slides[8]
        if "user_type_pie" in screenshots:
            _replace_image_in_slide(slide9, screenshots["user_type_pie"], shape_name="Picture 20")
        if "weekly_line_chart" in screenshots:
            _replace_image_in_slide(slide9, screenshots["weekly_line_chart"], shape_name="Picture 32")
        if "page_views_chart" in screenshots:
            _replace_image_in_slide(slide9, screenshots["page_views_chart"], shape_name="Picture 21")

    # Step 8: Get previous month data and generate Gemini recommendations
    previous_data = _get_previous_month_data(report_name)
    current_stats = (
        f"Date range: {date_range}\nReport date: {report_date}\n"
        f"Active users: {home_metrics.get('Active users', 'N/A')}\n"
        f"New users: {home_metrics.get('New users', 'N/A')}\n"
        f"Channels: {snapshot_metrics.get('channels', {})}\n"
        f"Top country: {snapshot_metrics.get('top_country', 'N/A')}"
    )
    recommendations = _generate_recommendations(
        report_name, current_stats, previous_data, date_range
    )

    # Step 9: Replace recommendations on the last slide
    # Each recommendation is a separate paragraph in Rectangle 4, starting with "1.", "2.", "3."
    # We preserve font family and size by copying the first run's rPr, then replacing all runs
    # with a single run carrying the full recommendation text.
    if recommendations:
        from pptx.oxml.ns import qn as _qn_rec
        from copy import deepcopy
        import lxml.etree as _etree
        last_slide = prs.slides[-1]
        for shape in last_slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                full_text = "".join(r.text for r in para.runs).strip()
                for i, rec in enumerate(recommendations, start=1):
                    if full_text.startswith(f"{i}."):
                        if para.runs:
                            # Copy rPr from first run to preserve font family + size
                            first_rpr = para.runs[0]._r.find(_qn_rec("a:rPr"))
                            saved_rpr = deepcopy(first_rpr) if first_rpr is not None else None
                            # Remove everything except <a:pPr> (paragraph properties)
                            for child in list(para._p):
                                if child.tag != _qn_rec("a:pPr"):
                                    para._p.remove(child)
                            # Build a new single run with the saved formatting
                            new_r = _etree.SubElement(para._p, _qn_rec("a:r"))
                            if saved_rpr is not None:
                                new_r.insert(0, saved_rpr)
                            new_t = _etree.SubElement(new_r, _qn_rec("a:t"))
                            new_t.text = rec
                        break

    # Step 10: Save output
    safe_name = report_name.replace("_", "-")
    output_path = OUTPUT_DIR / f"{safe_name}-{report_date.replace(' ', '-')}.pptx"
    prs.save(str(output_path))

    return output_path
